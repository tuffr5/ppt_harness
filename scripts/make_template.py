"""Build a branded template deck — a theme with no slides worth copying.

    python scripts/make_template.py tests/fixtures/brand-template.pptx

This exists because of what a demo of `serve --from` looked like without it. Both decks in
the repository carry the stock Office 2007 theme: white, Calibri, `#4F81BD`. A recording
that claims "the palette, the faces and the grid came across from your template" and then
shows black text on white has demonstrated nothing — the viewer cannot tell a borrowed theme
from a default one, because on that template they are the same thing.

So this writes a deck whose theme is *unmistakable*: a deep indigo brand, an ordered accent
ramp that reads as a house palette rather than Office's rainbow, a near-black ink that is not
`#000000`, and Avenir Next for display. Borrow it and the first generated slide is visibly
that company's slide.

`python-pptx` has no theme API — `theme1.xml` is not part of its object model — so the file
is authored normally and the theme part is replaced in the package afterwards. That is also
exactly the shape of the thing being tested: the harness reads its theme from `theme1.xml`
and nowhere else, so a deck whose *only* distinguishing feature is that part is the honest
test of whether extraction works.

One cover slide, kept deliberately: a template with zero slides is a legal `.pptx` that
several tools refuse to open, and the cover doubles as proof of the claim `--from` makes —
you can see the slide that did *not* come across.
"""

from __future__ import annotations

import argparse
import shutil
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

#: The palette, as roles rather than as decoration. `dk1`/`lt1` are ink and paper; `accent1`
#: is what the harness reads as `brand`. The ramp is ordered — a component asking for
#: "accent 2" gets the teal every time, which is what makes a deck look designed rather than
#: coloured in.
THEME = {
    "dk1": "10161F",       # ink — near-black, never #000000: pure black reads as harsh
    "lt1": "FFFFFF",       # paper
    "dk2": "1F2A37",       # secondary ink, for rules and captions
    "lt2": "F1F5F9",       # surface — cards, table zebra, callout fills
    "accent1": "1F4FD8",   # brand
    "accent2": "0CA678",
    "accent3": "F59F00",
    "accent4": "E8590C",
    "accent5": "7048E8",
    "accent6": "1098AD",
    "hlink": "1F4FD8",
    "folHlink": "7048E8",
}

#: Present on macOS, and different enough from Calibri that a viewer can see the swap. The
#: harness resolves the stack per script and reports what it could not find, so a machine
#: without Avenir Next degrades honestly rather than silently.
MAJOR_FONT = "Avenir Next"
MINOR_FONT = "Avenir Next"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def theme_xml(name: str = "Brand") -> bytes:
    """A complete `theme1.xml`.

    Written out rather than patched into the existing one: the stock theme carries an effect
    style list and a font scheme with two dozen script-specific faces, and editing around
    them leaves a file that is half Office's opinion and half ours.
    """
    scheme = "".join(
        f'<a:{slot}><a:srgbClr val="{value}"/></a:{slot}>'
        for slot, value in THEME.items()
    )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="{A}" name="{name}">
  <a:themeElements>
    <a:clrScheme name="{name}">{scheme}</a:clrScheme>
    <a:fontScheme name="{name}">
      <a:majorFont><a:latin typeface="{MAJOR_FONT}"/>
        <a:ea typeface=""/><a:cs typeface=""/></a:majorFont>
      <a:minorFont><a:latin typeface="{MINOR_FONT}"/>
        <a:ea typeface=""/><a:cs typeface=""/></a:minorFont>
    </a:fontScheme>
    <a:fmtScheme name="{name}">
      <a:fillStyleLst>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
      </a:fillStyleLst>
      <a:lnStyleLst>
        <a:ln w="9525"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
        <a:ln w="19050"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
        <a:ln w="28575"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
      </a:lnStyleLst>
      <a:effectStyleLst>
        <a:effectStyle><a:effectLst/></a:effectStyle>
        <a:effectStyle><a:effectLst/></a:effectStyle>
        <a:effectStyle><a:effectLst/></a:effectStyle>
      </a:effectStyleLst>
      <a:bgFillStyleLst>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
      </a:bgFillStyleLst>
    </a:fmtScheme>
  </a:themeElements>
</a:theme>'''.encode()


def _cover(prs) -> None:
    """The one slide, and the one that does not travel.

    Styled from the same palette so the file reads as a real template rather than as a
    fixture: whoever opens it should see the house style, and whoever borrows it should get
    the house style without this slide's words.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Emu(int(Inches(0.34))))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(THEME["accent1"])
    bar.line.fill.background()
    bar.text_frame.text = ""

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.0), Inches(2.4))
    frame = box.text_frame
    frame.word_wrap = True
    head = frame.paragraphs[0]
    run = head.add_run()
    run.text = "Brand template"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.name = MAJOR_FONT
    run.font.color.rgb = RGBColor.from_string(THEME["dk1"])

    sub = frame.add_paragraph()
    run = sub.add_run()
    run.text = ("Palette, type and grid live here. Slides do not — "
                "start a deck with `serve --from` and none of this slide comes with it.")
    run.font.size = Pt(20)
    run.font.name = MINOR_FONT
    run.font.color.rgb = RGBColor.from_string(THEME["dk2"])


def build(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    _cover(prs)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)

    # Swap the theme part. Rewriting the zip wholesale rather than editing in place, because
    # a `.pptx` is an ordered archive and appending a second `theme1.xml` produces a file
    # PowerPoint opens and the harness reads differently.
    staged = path.with_suffix(".staging.pptx")
    with zipfile.ZipFile(path) as src, zipfile.ZipFile(
            staged, "w", zipfile.ZIP_DEFLATED) as out:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                data = theme_xml()
            out.writestr(item, data)
    shutil.move(str(staged), str(path))
    return path


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    parser.add_argument("path", type=Path, nargs="?",
                        default=Path("tests/fixtures/brand-template.pptx"))
    args = parser.parse_args()

    written = build(args.path)

    # Read it back through the harness rather than trusting the write. A template whose
    # theme does not extract is worse than no template: every deck built from it would look
    # like the default one, and nothing would say so.
    from ppt_harness.io.theme_extract import extract_theme

    theme = extract_theme(written)
    print(f"{written}  ·  {written.stat().st_size / 1000:.0f} kB")
    print(f"  brand   {theme.palette['brand']}")
    print(f"  ink/bg  {theme.palette['ink']} on {theme.palette['bg']}")
    print(f"  accents {', '.join(a for a in theme.palette['accents'])}")
    print(f"  display {theme.type.families['display']}")
    if theme.inferred:
        print(f"  inferred {', '.join(theme.inferred)}")

    with zipfile.ZipFile(written) as z:
        parsed = etree.fromstring(z.read("ppt/theme/theme1.xml"))
    assert parsed.find(f".//{{{A}}}clrScheme") is not None, "theme lost its colour scheme"
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
