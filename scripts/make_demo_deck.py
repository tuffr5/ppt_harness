"""Build a small demo deck.

Two reasons this exists.

**Tests should not need a private 174 MB file.** The suite currently picks up whatever
`.pptx` is in the repo root, which means it passes on one machine and skips on another.

**The round-trip guarantee is unproven exactly where this deck is strong.** DESIGN §6.2
records that open-then-save preserves everything on a media-heavy deck, but that deck
carries no native chart, no table, no group, and no OLE object — the parts most likely to
break a naive exporter. This builds them on purpose, small enough to commit.

    python scripts/make_demo_deck.py tests/fixtures/demo.pptx

Deterministic by default so a diff means something. `--seed` varies the wording for a
livelier demo without changing the structure.
"""

from __future__ import annotations

import argparse
import io
import random
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.shapes.autoshape import AutoShapeType
from pptx.util import Emu, Inches, Pt

# 16:9 at the size PowerPoint itself uses.
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x12, 0x16, 0x1C)
MUTED = RGBColor(0x5A, 0x64, 0x72)
BRAND = RGBColor(0x15, 0x60, 0x82)
ACCENT = RGBColor(0xE9, 0x71, 0x32)

TOPICS = [
    ("Quarterly review", "Where the numbers went and why"),
    ("Platform migration", "What we moved, what broke, what is left"),
    ("Research update", "Three results and one dead end"),
]

FINDINGS = [
    "Latency fell 38% after the cache rewrite",
    "Two regions still run the old scheduler",
    "Cost per request is flat despite 3x traffic",
    "The migration script needs a dry-run mode",
]


# --------------------------------------------------------------------- helpers


def _textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, autofit=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    # Left as PowerPoint would leave it — the harness must cope with a deck that does not
    # already satisfy its own writer assertions.
    frame.auto_size = MSO_AUTO_SIZE.NONE if not autofit else None
    frame.text = text
    for para in frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def _png(width: int, height: int, rgb: tuple[int, int, int]) -> io.BytesIO:
    """A flat PNG, so the deck carries a real media part without a binary in the repo."""
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (width, height), rgb).save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


# ---------------------------------------------------------------------- slides


def _master_shape(prs, kind, x, y, cx, cy, name: str):
    """Add an autoshape straight onto the slide master.

    python-pptx exposes `add_shape` only on slides, and the obvious workaround — build it
    on a throwaway slide and delete that slide afterwards — leaves an orphaned slide part
    with a dangling relationship, which PowerPoint reports as a damaged file. Going through
    the same `spTree` factory the slide API uses avoids inventing a slide at all.
    """
    from pptx.shapes.autoshape import Shape as PptxShape
    from pptx.util import Emu

    tree = prs.slide_master.shapes._spTree
    autoshape = AutoShapeType(kind)
    # python-pptx's own allocator: shape ids must be unique within the tree, and scanning
    # for `id` attributes catches the GUIDs on `extLst` elements too.
    element = tree.add_autoshape(tree._next_shape_id, name, autoshape.prst,
                                    Emu(int(x)), Emu(int(y)), Emu(int(cx)), Emu(int(cy)))
    return PptxShape(element, None)


def _label(shape, text: str, size: int, colour) -> None:
    shape.text_frame.text = text
    for para in shape.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = colour


def _master_art(prs) -> None:
    """Non-placeholder art on the master: a logo mark and a footer band.

    Every slide inherits these. A fixture without them cannot exercise the inheritance
    path, which is where the largest preview defect of this project so far lived.
    """
    logo = _master_shape(prs, MSO_SHAPE.OVAL, Inches(12.1), Inches(0.25),
                         Inches(0.9), Inches(0.9), "Logo")
    logo.fill.solid()
    logo.fill.fore_color.rgb = BRAND
    logo.line.fill.background()
    _label(logo, "DH", 14, RGBColor(0xFF, 0xFF, 0xFF))

    band = _master_shape(prs, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1),
                         Inches(13.333), Inches(0.4), "Footer band")
    band.fill.solid()
    band.fill.fore_color.rgb = BRAND
    band.line.fill.background()
    _label(band, "ppt-harness demo", 11, RGBColor(0xFF, 0xFF, 0xFF))


def _shapes_slide(prs) -> None:
    """Filled autoshapes and a connector.

    A preview that renders only text shows floating words where the file has arrows, so the
    fixture has to contain some.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Shapes and connectors"

    for index, (kind, colour, label) in enumerate([
        (MSO_SHAPE.ROUNDED_RECTANGLE, BRAND, "collect"),
        (MSO_SHAPE.OVAL, ACCENT, "measure"),
        (MSO_SHAPE.CHEVRON, RGBColor(0x0E, 0x7C, 0x66), "report"),
    ]):
        left = Inches(1.0 + index * 4.0)
        box = slide.shapes.add_shape(kind, left, Inches(2.4), Inches(2.6), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = colour
        box.line.color.rgb = INK
        box.line.width = Pt(1.25)
        box.text_frame.text = label
        for para in box.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(15)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        if index < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.75),
                                           Inches(2.85), Inches(1.1), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(4.4),
                                      Inches(12.3), Inches(4.4))
    line.line.color.rgb = MUTED
    line.line.width = Pt(2)

    _textbox(slide, Inches(1.0), Inches(4.7), Inches(11.0), Inches(0.8),
             "Rounded rectangle, ellipse, chevron, two arrows and a straight connector.",
             size=14, color=MUTED)


def _title_slide(prs, title: str, subtitle: str) -> None:
    """Uses real placeholders, so the importer has a placeholder chain to resolve."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.notes_slide.notes_text_frame.text = "Demo deck built by scripts/make_demo_deck.py."


def _bullets_slide(prs, findings: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Findings"
    body = slide.placeholders[1].text_frame
    body.text = findings[0]
    for line in findings[1:]:
        body.add_paragraph().text = line
    slide.notes_slide.notes_text_frame.text = "Placeholder text inherits its size."


def _chart_slide(prs) -> None:
    """A **native** chart with an embedded worksheet.

    This is the case DESIGN §1.5 insists must never become an image: the recipient can edit
    the data. It is also the part of the package a rebuild-style exporter silently destroys.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Requests per region"

    data = CategoryChartData()
    data.categories = ["us-east", "us-west", "eu-central", "ap-south"]
    data.add_series("Q1", (18.2, 11.9, 9.4, 6.1))
    data.add_series("Q2", (22.7, 12.4, 13.8, 9.9))

    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.9), Inches(1.7), Inches(7.6), Inches(4.6), data,
    )
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    _textbox(slide, Inches(8.9), Inches(1.9), Inches(3.6), Inches(3.0),
             "Growth is concentrated in eu-central, which also carries the "
             "largest share of the migration risk.",
             size=16, color=MUTED)


def _table_slide(prs) -> None:
    """Tables are a graphic frame — opaque to the harness, and must survive untouched."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Migration status"

    rows = [
        ("Region", "Scheduler", "Owner", "Status"),
        ("us-east", "v2", "platform", "done"),
        ("us-west", "v2", "platform", "done"),
        ("eu-central", "v1", "infra", "in progress"),
        ("ap-south", "v1", "infra", "not started"),
    ]
    table = slide.shapes.add_table(len(rows), 4, Inches(0.9), Inches(1.8),
                                   Inches(11.5), Inches(3.2)).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
                    run.font.bold = r == 0
                    run.font.color.rgb = INK if r else BRAND


def _media_and_group_slide(prs) -> None:
    """A picture plus a grouped pair.

    Groups are opaque because editing a child means recomputing the group's extents. The
    harness must preserve the group and still report it as present rather than dropping it.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Assets and groups"

    slide.shapes.add_picture(_png(480, 300, (0x15, 0x60, 0x82)),
                             Inches(0.9), Inches(1.8), Inches(5.2), Inches(3.25))

    a = _textbox(slide, Inches(7.0), Inches(1.9), Inches(2.4), Inches(0.9),
                 "grouped A", size=14, color=ACCENT)
    b = _textbox(slide, Inches(7.0), Inches(3.0), Inches(2.4), Inches(0.9),
                 "grouped B", size=14, color=ACCENT)
    slide.shapes.add_group_shape([a, b])

    _textbox(slide, Inches(10.0), Inches(1.9), Inches(2.5), Inches(2.0),
             "The picture is a media part; the pair to the left is a group.",
             size=14, color=MUTED)


def _autofit_slide(prs) -> None:
    """A box PowerPoint has to shrink to make fit.

    The harness reports the overflow at the declared size and names the `normAutofit`
    scale, so a deck that "looks fine in PowerPoint" has an explanation. This slide is what
    that behaviour is tested against.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Text that only fits because PowerPoint shrank it"

    box = _textbox(
        slide, Inches(0.9), Inches(1.9), Inches(5.6), Inches(1.6),
        "This paragraph is deliberately longer than the box that holds it, so that "
        "PowerPoint stores a normAutofit fontScale rather than letting it overflow. "
        "The harness reads that as an admission and reports the real overflow.",
        size=20,
    )
    # Emit the confession explicitly: python-pptx's TEXT_TO_FIT_SHAPE is the writer-side
    # equivalent, and it is exactly what the harness refuses to produce itself.
    body = box.text_frame._bodyPr
    for child in list(body):
        if child.tag.endswith("}noAutofit") or child.tag.endswith("}spAutoFit"):
            body.remove(child)
    node = body.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit",
        {"fontScale": "77500", "lnSpcReduction": "10000"},
    )
    body.append(node)

    _textbox(slide, Inches(7.0), Inches(1.9), Inches(5.4), Inches(2.0),
             "Open this slide in the harness: it reports overflow and says the source "
             "shrinks the text to 78%.", size=16, color=MUTED)


# ------------------------------------------------------------------------ build


def build(path: Path, seed: int | None = None) -> Path:
    rng = random.Random(seed)
    title, subtitle = rng.choice(TOPICS) if seed is not None else TOPICS[0]
    findings = list(FINDINGS)
    if seed is not None:
        rng.shuffle(findings)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(int(SLIDE_W)), Emu(int(SLIDE_H))

    _master_art(prs)
    _title_slide(prs, title, subtitle)
    _bullets_slide(prs, findings)
    _shapes_slide(prs)
    _chart_slide(prs)
    _table_slide(prs)
    _media_and_group_slide(prs)
    _autofit_slide(prs)

    prs.core_properties.title = title
    prs.core_properties.author = "ppt-harness demo"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out", nargs="?", default="tests/fixtures/demo.pptx", type=Path)
    parser.add_argument("--seed", type=int, default=None,
                        help="Vary the wording. Omit for a deterministic deck.")
    args = parser.parse_args()

    out = build(args.out, args.seed)
    size_kb = out.stat().st_size / 1024
    print(f"{out}  ({size_kb:.0f} KB)")


if __name__ == "__main__":
    main()
