"""Icons on `icon_row` — DESIGN §1.4, §1.5, §3.

`icon_row` was a component named for something it did not draw. It held a `title` and a list
of bare labels, its two variants were called `icon_top` and `icon_left`, and nothing anywhere
in the pipeline produced a mark — so the variants differed in `per_row` and alignment and in
nothing a reader would recognise as an icon being above a word rather than beside it. Same
class of defect as the variants that rendered identically and the `media` slot the writer
skipped, and it is checked the same way: against the **exported file**, because the file is
what the recipient opens and the preview is that file rendered.

Five claims, one per section:

- a named icon reaches the .pptx as real vector geometry, coloured through a theme role;
- a name the harness does not have is refused at the tool gate, before anything renders;
- `icon_top` and `icon_left` put the mark in genuinely different places;
- the preview draws what the file contains, from the same coordinates;
- and reopening an exported deck recognises the marks again, without a marker in the file.
"""

from __future__ import annotations

import re
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

from ppt_harness.components import icons, registry
from ppt_harness.core.session import Session
from ppt_harness.io import import_pptx
from ppt_harness.io.export_mutate import ICON_PATH_UNITS, export
from ppt_harness.render import expand, html
from ppt_harness.state.document import Block, Mode, Slide
from ppt_harness.state.theme_default import default_theme
from ppt_harness.tools import router

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
THEME = default_theme()

ITEMS = [
    {"icon": "growth", "label": "Revenue up 38%"},
    {"icon": "team", "label": "Twelve new hires"},
    {"icon": "deadline", "label": "Shipped a week early"},
]


def _add(session: Session, variant: str, items: list | None = None) -> dict:
    return router.dispatch(session, "add_slide", {
        "layout": "stack",
        "blocks": [{"region": "body", "component": "icon_row", "variant": variant,
                    "slots": {"items": items if items is not None else ITEMS}}],
    })


def _slide(variant: str, items: list | None = None) -> Slide:
    return Slide(id="s", index=0, mode=Mode.MANAGED, layout="stack", blocks=[
        Block(id="bk", region="body", component="icon_row", variant=variant,
              slots={"items": items if items is not None else ITEMS})])


def _items_slot(variant: str, items: list | None = None) -> expand.LaidOutSlot:
    laid = expand.expand_slide(THEME, _slide(variant, items))
    return next(s for s in laid if s.slot == "items")


# --------------------------------------------------------------- the vendored set


def test_every_vendored_path_parses_to_drawable_commands() -> None:
    """The table is data, so it is the one thing a code change cannot keep honest.

    A path the exporter cannot parse writes an empty `a:pathLst` — a shape with no geometry,
    which is invisible and reports no error anywhere. Every icon is walked once here so a
    bad re-vendor fails in the suite rather than on somebody's slide.
    """
    assert len(icons.names()) >= 100, "the curated set is meant to cover a deck's concepts"
    for name in icons.names():
        segments = icons.segments(name)
        assert segments, f"{name} has no drawable geometry"
        assert segments[0][0] == "M", f"{name} does not begin with a moveto"
        for command, values in segments:
            assert command in ("M", "L", "C", "Z"), f"{name} kept a {command} command"
            assert len(values) == {"M": 2, "L": 2, "C": 6, "Z": 0}[command]
            # Inside the view box, with a hair of slack for a stroke drawn on the edge.
            assert all(-1 <= v <= icons.view_box() + 1 for v in values), \
                f"{name} draws outside its own view box"


def test_the_licence_travels_with_the_artwork() -> None:
    """MIT requires the notice to ship with the copy. It is also how anyone finds out where
    these came from — a JSON of path data with no provenance is unattributable by
    construction."""
    beside = Path(icons.__file__).parent / "iconset"
    assert (beside / "LICENSE").read_text().strip().startswith("MIT License")
    assert icons.attribution()["source"].endswith("tabler-icons")
    assert icons.attribution()["licence"] == "MIT"


# ------------------------------------------------------------------- reaching the file


def test_a_named_icon_reaches_the_file_as_vector_geometry(blank: Session,
                                                          tmp_path: Path) -> None:
    """The mark is `a:custGeom`, not a picture and not a placeholder.

    Geometry rather than a raster is the whole design: python-pptx cannot read SVG at all
    (`add_asset` refuses it with `svg_unsupported`), and a PNG could not take the theme's
    colour. So the assertion is on the emitted XML — a custom path, with real curves in it,
    stroked in the block's accent.
    """
    assert _add(blank, "icon_top")["ok"]
    out = tmp_path / "icons.pptx"
    report = export(blank.deck, out)
    assert not report.violations

    with zipfile.ZipFile(out) as z:
        xml = z.read("ppt/slides/slide1.xml").decode()
    assert xml.count("<a:custGeom>") == len(ITEMS), "one shape per named icon"
    assert "<a:cubicBezTo>" in xml, "the curves flattened to straight lines"
    accent = THEME.accent(0).lstrip("#").upper()
    assert re.search(rf'<a:ln w="\d+" cap="rnd"><a:solidFill><a:srgbClr val="{accent}"',
                     xml), "an icon must resolve its colour through a theme role"


def test_an_icon_is_a_shape_the_recipient_can_still_edit(blank: Session,
                                                         tmp_path: Path) -> None:
    """Selectable, movable, recolourable — the same claim `chart` makes about its plot.

    A mark baked into the slide background, or drawn as a picture, would be none of those.
    It also has to sit inside the slide, like every other box the expander produces.
    """
    assert _add(blank, "icon_top")["ok"]
    out = tmp_path / "icons.pptx"
    export(blank.deck, out)
    prs = Presentation(str(out))
    marks = [s for s in prs.slides[0].shapes if s.name.endswith(".icon")]
    assert len(marks) == len(ITEMS)
    for mark in marks:
        assert mark.width == mark.height, "an icon's frame must stay square"
        assert mark.left >= 0 and mark.top >= 0
        assert mark.left + mark.width <= int(prs.slide_width)
        assert mark.top + mark.height <= int(prs.slide_height)


def test_the_accent_override_recolours_every_mark(blank: Session, tmp_path: Path) -> None:
    """No hex in the catalog: the icon takes whichever accent the block is set to.

    This is the property a rasterised icon set cannot have, and the reason `custGeom` was
    worth the path table.
    """
    result = _add(blank, "icon_top")
    slide_id = result["target"]
    block = blank.deck.slides[0].blocks[0]
    assert router.dispatch(blank, "set_override", {
        "slide_id": slide_id, "block_id": block.id, "key": "accent", "value": 2})["ok"]

    out = tmp_path / "accent.pptx"
    export(blank.deck, out)
    with zipfile.ZipFile(out) as z:
        xml = z.read("ppt/slides/slide1.xml").decode()
    assert THEME.accent(2).lstrip("#").upper() in xml
    assert f'cap="rnd"><a:solidFill><a:srgbClr val="{THEME.accent(0).lstrip("#").upper()}"' \
        not in xml


# ----------------------------------------------------------------------- the gate


def test_an_unknown_icon_is_refused_before_anything_renders(blank: Session) -> None:
    """A write that lands and cannot be drawn is worse than one that is refused.

    Worse here than anywhere else in the catalog, because it fails *silently*: a bad icon
    name would export as a label with no mark, the budget would pass, the export would be
    clean, and the deck would quietly contain the component's original bug. The refusal has
    to name the available set, the way `_check_media` names what a media slot takes.
    """
    result = _add(blank, "icon_top", [{"icon": "unicorn", "label": "Nope"}])
    assert result["ok"] is False
    assert result["error"] == "unknown_icon"
    assert "growth" in result["message"], "the refusal must list what is available"


def test_a_near_miss_is_offered_the_name_it_meant(blank: Session) -> None:
    """A hundred and forty-five names is a complete answer and an unusable one."""
    detail = _add(blank, "icon_top", [{"icon": "grwoth", "label": "Typo"}])["message"]
    assert "Did you mean" in detail and "growth" in detail


def test_an_item_with_no_icon_is_refused(blank: Session) -> None:
    """The icon is the component, not a decoration on it.

    Allowing a bare label would put `icon_row` straight back where it started: a row named
    for a mark with no mark on it, and two variants that are then indistinguishable, because
    where the mark goes is the only thing that distinguishes them.
    """
    result = _add(blank, "icon_top", ["Just a label"])
    assert result["ok"] is False
    assert result["error"] == "icon_required"


def test_the_gate_guards_set_slots_too(blank: Session) -> None:
    """Both doors into a slot, or the check is a suggestion."""
    slide_id = _add(blank, "icon_top")["target"]
    block = blank.deck.slides[0].blocks[0]
    result = router.dispatch(blank, "set_slots", {
        "slide_id": slide_id, "block_id": block.id,
        "patch": {"items": [{"icon": "nope", "label": "x"}]}})
    assert result["ok"] is False
    assert result["error"] == "unknown_icon"


def test_the_catalog_tells_a_model_what_the_names_are() -> None:
    """A gate that refuses names a model has no way to discover is a trap, not a gate."""
    described = registry.describe("icon_row")
    assert "growth" in described["icons"]
    assert described["slots"]["items"]["icons"] is True
    assert {v["icon"] for v in described["variants"].values()} == {"top", "left"}
    # Not in the always-on catalog: a hundred and forty-five words on every turn buys
    # nothing until a model is actually filling this component.
    assert all("icons" not in entry for entry in registry.catalog())


# ------------------------------------------------------------------ the two variants


def test_the_variants_place_the_mark_differently(blank: Session, tmp_path: Path) -> None:
    """`icon_top` and `icon_left` were two names for one rendering.

    Measured on the file, per cell: under `icon_top` the mark is above its label and shares
    the cell's horizontal centre with it; under `icon_left` it is beside the label and to
    the left of it. Comparing the two exports rather than the two `Variant` records is the
    point — the catalog claiming a difference is exactly what was wrong before.
    """
    places = {}
    for variant in ("icon_top", "icon_left"):
        session = Session.blank("v")
        assert _add(session, variant)["ok"]
        out = tmp_path / f"{variant}.pptx"
        export(session.deck, out)
        shapes = {s.name: s for s in Presentation(str(out)).slides[0].shapes}
        mark = next(s for n, s in shapes.items() if n.endswith("#0.icon"))
        label = next(s for n, s in shapes.items() if n.endswith(":items#0"))
        places[variant] = (mark, label)

    top_mark, top_label = places["icon_top"]
    assert top_mark.top + top_mark.height <= top_label.top, "icon_top must be above"
    assert abs((top_mark.left + top_mark.width / 2)
               - (top_label.left + top_label.width / 2)) < top_mark.width * 0.1

    left_mark, left_label = places["icon_left"]
    assert left_mark.left + left_mark.width <= left_label.left, "icon_left must be beside"
    assert left_mark.top < left_label.top + left_label.height, "and not above"


def test_the_mark_never_overlaps_the_words_it_labels() -> None:
    """Overlap would mean the budget was measuring a box the text does not get.

    `LaidOutSlot.cells` is what `budget.for_slot` reads, so the mark's square being taken out
    of it *there* is what keeps the gate charging a labelled mark for the room it actually
    leaves. Checked against every region the component declares, because a footer band and a
    body region are very different cells.
    """
    for layout, region in (("stack", "body"), ("hero_plus_row", "footer_row"),
                           ("full_bleed", "canvas")):
        for variant in ("icon_top", "icon_left"):
            slide = Slide(id="s", index=0, mode=Mode.MANAGED, layout=layout, blocks=[
                Block(id="bk", region=region, component="icon_row", variant=variant,
                      slots={"items": ITEMS})])
            laid = next(s for s in expand.expand_slide(THEME, slide) if s.slot == "items")
            for cell, mark in zip(laid.cells(len(ITEMS)), laid.icons(len(ITEMS)),
                                  strict=True):
                assert mark.w == mark.h > 0, f"{layout}/{variant} drew no mark"
                apart = (mark.x + mark.w <= cell.x + 0.5 or cell.x + cell.w <= mark.x + 0.5
                         or mark.y + mark.h <= cell.y + 0.5
                         or cell.y + cell.h <= mark.y + 0.5)
                assert apart, f"{layout}/{variant}: the mark sits on its own label"


def test_adding_icons_did_not_take_capacity_off_the_component() -> None:
    """The catalog's own worst case still has to fit — DESIGN §3, and `fixtures.payload`.

    A mark sized off the cell alone would have eaten the label's second line, and the budget
    reads the carved box, so the component would have started refusing content it advertises
    it can hold. `_icon_metrics` holds back the lines the label would have had anyway.
    """
    from ppt_harness.components import fixtures
    from ppt_harness.render import budget as budget_mod

    payload = fixtures.payload("icon_row")
    assert all("icon" in item for item in payload["items"])
    for variant in registry.get("icon_row").variants:
        layout, region = fixtures.region_for("icon_row")
        slide = Slide(id="s", index=0, mode=Mode.MANAGED, layout=layout, blocks=[
            Block(id="bk", region=region, component="icon_row", variant=variant,
                  slots=payload)])
        for laid in expand.expand_slide(THEME, slide):
            result = budget_mod.check_value(payload[laid.slot],
                                            budget_mod.for_slot(THEME, laid), THEME)
            assert result.ok, f"icon_row/{variant}.{laid.slot}: {result.error(THEME)}"


# ------------------------------------------------------------ preview equals export


@pytest.mark.parametrize("variant", ["icon_top", "icon_left"])
def test_the_preview_draws_the_same_marks_in_the_same_boxes(blank: Session, tmp_path: Path,
                                                            variant: str) -> None:
    """Preview-equals-export is the invariant everything else here is measured against.

    Both sides are checked against the same numbers: the preview's `<svg>` rectangles are
    compared to the shape frames in the exported package, converted through the one scale
    the writer uses. A mark the preview drew somewhere else would be a preview of a
    different slide — which is exactly how `image_full/full` and `image_full/inset` came to
    look identical in both.
    """
    assert _add(blank, variant)["ok"]
    out = tmp_path / "both.pptx"
    export(blank.deck, out)
    prs = Presentation(str(out))
    canvas_w, canvas_h = blank.theme.grid.canvas
    exported = sorted(
        (s.left / int(prs.slide_width) * canvas_w, s.top / int(prs.slide_height) * canvas_h,
          s.width / int(prs.slide_width) * canvas_w)
         for s in prs.slides[0].shapes if s.name.endswith(".icon"))

    slide = blank.deck.slides[0]
    markup = html.render_slide(blank.theme, slide, 0, 0).html
    drawn = re.findall(
        r'<svg class="icon"[^>]*style="left:([\d.]+)px; top:([\d.]+)px; '
        r'width:([\d.]+)px', markup)
    assert len(drawn) == len(ITEMS), "the preview drew a different number of marks"
    preview = sorted((float(x), float(y), float(w)) for x, y, w in drawn)

    for (ex, ey, ew), (px, py, pw) in zip(exported, preview, strict=True):
        assert abs(ex - px) < 1.5 and abs(ey - py) < 1.5 and abs(ew - pw) < 1.5


def test_the_preview_and_the_writer_read_one_path_table() -> None:
    """The `d` in the markup is the string the exporter turns into `a:custGeom`.

    Two copies of an icon's outline is the shape of the bug that makes a curve bend one way
    on screen and the other in the file — which is why the artwork is normalised to
    `M`/`L`/`C`/`Z` once, at vendor time, instead of being converted by each consumer.
    """
    slide = _slide("icon_top")
    markup = html.render_slide(THEME, slide, 0, 0).html
    for item in ITEMS:
        assert f'd="{icons.path(item["icon"])}"' in markup


def test_the_stroke_weight_is_one_ratio_in_both() -> None:
    """A mark drawn at two weights reads as two different drawings.

    The preview states the width in view-box units so the browser scales it with the artwork
    exactly as PowerPoint scales `a:ln` with its shape; both come off `icon_stroke_px`.
    """
    laid = _items_slot("icon_top")
    box = laid.icons(len(ITEMS))[0]
    assert expand.icon_stroke_px(box) == pytest.approx(
        box.w * icons.stroke_units() / icons.view_box())
    markup = html.render_slide(THEME, _slide("icon_top"), 0, 0).html
    assert f'stroke-width="{icons.stroke_units():.3f}"' in markup
    markup = html.render_slide(THEME, _slide("icon_top"), 0, 0).html
    assert f'stroke-width="{icons.stroke_units():.3f}"' in markup


# --------------------------------------------------------------------- the one-way door


def test_ejecting_a_slide_keeps_its_marks(blank: Session, tmp_path: Path) -> None:
    """`eject_slide` is one-way, so anything it drops is gone.

    The panels and the picture were already fixed for exactly this — a door that quietly
    costs you the cards is a door people stop trusting — and an `icon_row` ejected without
    its icons is the same loss with a different name: the blocks are gone by the time anyone
    notices, and there is nothing left to re-derive the marks from.
    """
    slide_id = _add(blank, "icon_top")["target"]
    before = len(ITEMS)
    assert router.dispatch(blank, "eject_slide", {"slide_id": slide_id})["ok"]
    assert blank.deck.slides[0].mode is Mode.FREEFORM

    out = tmp_path / "ejected.pptx"
    export(blank.deck, out)
    with zipfile.ZipFile(out) as z:
        xml = z.read("ppt/slides/slide1.xml").decode()
    assert xml.count("<a:custGeom>") == before, "the freeze dropped the marks"
    assert THEME.accent(0).lstrip("#").upper() in xml, "and their colour with them"

    # And the preview of the ejected slide still draws them, from the same path table.
    markup = html.render_slide(blank.theme, blank.deck.slides[0], 1280, 720).html
    assert markup.count('class="icon"') == before


# ------------------------------------------------------------------- and back again


def _icon_shapes(session: Session) -> list[str]:
    """The icon behind every recognised shape on slide 1, in document order."""
    return [s.geometry.icon for s in session.deck.slides[0].shapes
            if s.geometry is not None and s.geometry.icon]


def _slide_xml(path: Path) -> bytes:
    """Slide 1, canonicalised — the comparison DESIGN §6.2 defines for an untouched part."""
    with zipfile.ZipFile(path) as z:
        return etree.tostring(etree.fromstring(z.read("ppt/slides/slide1.xml")),
                              method="c14n2")


def _marks(path: Path) -> list[bytes]:
    """Every `a:custGeom` on slide 1, serialised, in document order."""
    with zipfile.ZipFile(path) as z:
        root = etree.fromstring(z.read("ppt/slides/slide1.xml"))
    return [etree.tostring(node) for node in root.iter(f"{{{A}}}custGeom")]


def test_an_exported_icon_is_recognised_when_the_deck_is_reopened(blank: Session,
                                                                  tmp_path: Path) -> None:
    """Export, `Session.open`, and the marks are icons again — DESIGN §7 step 3.

    Import lands every slide as `freeform`, so the marks come back as shapes carrying
    `a:custGeom`. Unrecognised, `_geometry` found no `a:prstGeom`, fell back to `preset="rect"`
    and the preview drew a **stroked rectangle** for every icon on the slide — the harness
    losing something it wrote itself and then showing a rendering the file does not contain.
    """
    assert _add(blank, "icon_top")["ok"]
    out = tmp_path / "icons.pptx"
    export(blank.deck, out)

    reopened = Session.open(out)
    assert _icon_shapes(reopened) == [item["icon"] for item in ITEMS]

    # The point of knowing: the preview of the reopened deck draws the marks, not the boxes.
    markup = html.render_slide(reopened.deck.theme, reopened.deck.slides[0],
                               *reopened.deck.theme.grid.canvas).html
    assert markup.count('class="icon"') == len(ITEMS)
    for item in ITEMS:
        assert f'd="{icons.path(item["icon"])}"' in markup


def test_geometry_that_is_not_in_the_table_is_never_given_a_name(blank: Session,
                                                                 tmp_path: Path) -> None:
    """A false positive is worse than a miss, so this is the test that matters.

    A hand-drawn freeform — same `a:custGeom` element, same path space, a shape nobody
    vendored — must come back nameless. Recognition is exact: one unit of drift anywhere in
    the path and the answer is "no icon", which leaves the shape exactly the anonymous
    freeform it was before any of this existed.
    """
    assert _add(blank, "icon_top")["ok"]
    out = tmp_path / "icons.pptx"
    export(blank.deck, out)

    prs = Presentation(str(out))
    marks = [s for s in prs.slides[0].shapes if s.name.endswith(".icon")]
    # Move one point by a single unit out of 21600 — far below anything a reader could see,
    # and the whole distance between "this is `growth`" and "this is a drawing".
    point = marks[0]._element.spPr.find(f"{{{A}}}custGeom").find(f".//{{{A}}}pt")
    point.set("x", str(int(point.get("x")) + 1))
    # The second keeps its path and is renamed to what a marker-in-the-name scheme would have
    # written, which is the case that scheme cannot survive: PowerPoint's Selection Pane lets
    # anyone type this, so it must neither be believed nor get in the way of the geometry.
    marks[1].name = "definitely_a_growth.icon"
    prs.save(str(out))

    reopened = Session.open(out)
    assert _icon_shapes(reopened) == [ITEMS[1]["icon"], ITEMS[2]["icon"]], (
        "a nudged path was still named, or a renamed shape stopped being read")


def test_the_importer_refuses_every_path_it_did_not_write() -> None:
    """The reader's ways of saying no, one element each.

    Reached here rather than through a file because each is a different malformed neighbour of
    a real icon, and building three .pptx files to exercise three `return ""` lines would test
    python-pptx rather than the refusals. The shared premise: a wrong name is worse than none,
    so anything that is not exactly the geometry this harness emits is not an icon.
    """
    real = icons.geometry_key(icons.placed("growth", ICON_PATH_UNITS))
    tag = {"M": "moveTo", "L": "lnTo", "C": "cubicBezTo"}
    body = ""
    for command, values in icons.placed("growth", ICON_PATH_UNITS):
        if command == "Z":
            body += "<a:close/>"
            continue
        points = "".join(f'<a:pt x="{x}" y="{y}"/>'
                         for x, y in zip(values[::2], values[1::2], strict=True))
        body += f"<a:{tag[command]}>{points}</a:{tag[command]}>"

    def read(paths: str) -> str:
        sp_pr = etree.fromstring(
            f'<p:spPr xmlns:p="{P}" xmlns:a="{A}"><a:custGeom><a:pathLst>{paths}'
            "</a:pathLst></a:custGeom></p:spPr>")
        return import_pptx._icon_name(sp_pr)

    square = f'<a:path w="{ICON_PATH_UNITS}" h="{ICON_PATH_UNITS}">'
    assert read(f"{square}{body}</a:path>") == "growth", "the premise: this one is real"
    assert read(f"{square}{body}</a:path>{square}{body}</a:path>") == "", "two paths"
    assert read(f'<a:path w="{ICON_PATH_UNITS}" h="800">{body}</a:path>') == "", "not square"
    assert read(f'<a:path w="1in" h="1in">{body}</a:path>') == "", "a universal measure"
    assert read(f"<a:path>{body}</a:path>") == "", "no path space at all"
    # An `a:arcTo` is the command the vendored set is normalised *away* from, so a path
    # carrying one was drawn by something else no matter how much of it matches.
    assert read(f'{square}{body}<a:arcTo wR="1" hR="1" stAng="0" swAng="1"/></a:path>') == ""
    # A `a:lnTo` carrying three points is not a `a:lnTo` — arity is checked, not assumed.
    assert read(f'{square}<a:lnTo><a:pt x="1" y="2"/><a:pt x="3" y="4"/></a:lnTo>'
                f"</a:path>") == ""
    assert icons.identify(real, ICON_PATH_UNITS) == "growth", "the table stopped agreeing"


def test_recognition_costs_the_exported_file_nothing(blank: Session, tmp_path: Path) -> None:
    """Re-export after a round trip is the same file — DESIGN §6.2.

    Recognition reads the geometry that was already there rather than writing a marker beside
    it, so the writer is untouched and this is close to free. It is asserted anyway, because
    "export mutates, never rebuilds" is the assumption every imported deck's SmartArt and
    sensitivity label rest on, and an icon is the harness's own geometry sitting in a file it
    is now looking at on the way in.

    Canonical equality on the slide part and byte equality on the marks, which is §6.2's own
    contract rather than a weaker one chosen to pass: whole-file identity is explicitly *not*
    the assertion there, because python-pptx reserialises every XML declaration and rebuilds
    `.rels` from its own model. The geometry itself has no such excuse — a path that drifted
    a unit per round trip would still satisfy "canonically equal to a *different* path", so
    the marks are compared byte for byte across all three exports.
    """
    assert _add(blank, "icon_top")["ok"]
    first = tmp_path / "first.pptx"
    export(blank.deck, first)

    second = tmp_path / "second.pptx"
    reopened = Session.open(first)
    export(reopened.deck, second)
    third = tmp_path / "third.pptx"
    again = Session.open(second)
    export(again.deck, third)

    assert _icon_shapes(reopened) == [item["icon"] for item in ITEMS]
    assert _icon_shapes(again) == [item["icon"] for item in ITEMS], "recognition decayed"
    assert _slide_xml(first) == _slide_xml(second), "reopening rewrote the slide"
    assert _slide_xml(second) == _slide_xml(third)
    assert len(_marks(first)) == len(ITEMS)
    assert _marks(first) == _marks(second) == _marks(third), "the geometry drifted"


def test_the_table_has_no_two_icons_the_importer_could_confuse() -> None:
    """145 paths, 145 keys. The property recognition is only sound while it holds.

    Checked against the writer's own path space and against the raw view box, because the
    second is the coarsest space any of this can be asked about — if the set is unambiguous
    at 24 units it is unambiguous at every resolution above it, and a re-vendor that
    introduced a duplicate would fail here rather than mislabelling somebody's slide.
    """
    for units in (ICON_PATH_UNITS, int(icons.view_box())):
        keys = {icons.geometry_key(icons.placed(name, units)): name
                for name in icons.names()}
        assert len(keys) == len(icons.names()), f"two icons share a path at {units} units"
        for key, name in keys.items():
            assert icons.identify(key, units) == name
    assert icons.identify("M 0 0 L 1 1", ICON_PATH_UNITS) == "", "an invented path was named"
    assert icons.identify("", 0) == "", "a path space of nothing named something"
