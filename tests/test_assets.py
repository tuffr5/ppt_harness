"""Getting a picture into a deck — `add_asset`, DESIGN §1.5.

The gap this closes is worth stating, because the tests are shaped by it: `Session.assets`
had exactly one filler, import, so a deck the harness *generated* could not hold a picture at
all. The only thing that ever made one work was a fallback that read a slot's `asset_id` as a
path on this machine — which is to say the deck's content lived outside the deck, and the
picture vanished the moment the file was opened anywhere else.

So what is under test is not really "the tool returns a key". It is:

- the key is enough, on its own, to put the image in the exported file *and* in the preview
  — the one equivalence the whole harness is measured against;
- the same picture added twice is one picture, because a deck that grows every time somebody
  re-adds their logo is a deck nobody can mail;
- a rejected write is the cheapest failure, so what is not an image is refused by its bytes
  and said plainly; and
- it undoes, like every other write.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from ppt_harness.core.session import Session
from ppt_harness.io import media
from ppt_harness.io.export_mutate import export
from ppt_harness.state.document import Mode
from ppt_harness.tools import router

PICTURE = (320, 200)


def _png(path: Path, size: tuple[int, int] = PICTURE,
         colour: tuple[int, int, int] = (0x15, 0x60, 0x82)) -> Path:
    """A real image, generated rather than committed: a binary in the repository is a thing
    to explain, and every test here cares only about its bytes and its proportions."""
    from PIL import Image

    Image.new("RGB", size, colour).save(path)
    return path


@pytest.fixture
def picture(tmp_path: Path) -> Path:
    return _png(tmp_path / "q3 revenue.png")


def _add(session: Session, path: Path, **extra) -> dict:
    return router.dispatch(session, "add_asset", {"path": str(path), **extra})


def _place_media(session: Session, key: str, alt: str = "A blue rectangle") -> None:
    assert router.dispatch(session, "add_slide", {"layout": "stack", "blocks": [
        {"region": "body", "component": "image_full", "variant": "full",
         "slots": {"media": {"asset_id": key, "alt": alt}}}]})["ok"]


def _pictures(path: Path, index: int = 0):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return [s for s in Presentation(str(path)).slides[index].shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


# ------------------------------------------------------------------- the ingest


def test_an_image_is_added_and_reachable_by_its_key(blank: Session, picture: Path) -> None:
    """The whole point: bytes in, a name back, and the name resolves like any other asset."""
    result = _add(blank, picture)
    assert result["ok"], result.get("message")

    key = result["after"]["asset_id"]
    assert key in blank.assets
    content_type, blob = blank.assets[key]
    assert content_type == "image/png"
    assert blob == picture.read_bytes()
    # Reachable by the two paths that matter: the exporter's resolver and the preview's.
    assert media.resolve(key, blank.assets) is not None
    assert (blank.asset_data_uri(key) or "").startswith("data:image/png;base64,")


def test_what_it_learned_about_the_image_comes_back_with_the_key(blank: Session,
                                                                 picture: Path) -> None:
    """The caller is about to choose a slot, and these are the facts that decide.

    A 3:1 panorama in a square slot letterboxes to two thick bands of background, and the
    moment to learn that is before the slide is built rather than from the render after.
    """
    after = _add(blank, picture)["after"]
    assert after["media_type"] == "image/png"
    assert (after["width_px"], after["height_px"]) == PICTURE
    assert after["aspect"] == pytest.approx(PICTURE[0] / PICTURE[1], rel=1e-3)
    assert after["bytes"] == picture.stat().st_size


def test_the_key_is_legible_and_carries_the_filename(blank: Session, picture: Path) -> None:
    """A pure hash dedupes perfectly and reads like line noise in an op log somebody is
    trying to follow. The stem carries the meaning; the digest settles collisions."""
    key = _add(blank, picture)["after"]["asset_id"]
    assert key.startswith("q3-revenue-"), key
    assert media.MAX_ASSET_BYTES  # the limit is shared with import, not restated here


def test_the_same_image_twice_does_not_duplicate_it(blank: Session, picture: Path) -> None:
    """Identity is the bytes. Re-adding a picture — the logo, on every third slide — must
    not grow the deck, and a model that retries a call it is unsure landed must not either."""
    first = _add(blank, picture)["after"]
    second = _add(blank, picture)["after"]

    assert second["asset_id"] == first["asset_id"]
    assert second["reused"] is True
    assert len(blank.assets) == 1


def test_the_same_bytes_under_another_filename_are_still_one_asset(blank: Session,
                                                                   tmp_path: Path) -> None:
    """Content-addressed means the directory it was copied from is not part of its identity."""
    one = _png(tmp_path / "chart.png")
    two = tmp_path / "copy" / "chart.png"
    two.parent.mkdir()
    two.write_bytes(one.read_bytes())

    assert _add(blank, one)["after"]["asset_id"] == _add(blank, two)["after"]["asset_id"]
    assert len(blank.assets) == 1


def test_a_caller_can_name_the_key(blank: Session, picture: Path) -> None:
    """A script that wants `logo` in its slot payloads should be able to say `logo`."""
    result = _add(blank, picture, key="logo")
    assert result["after"]["asset_id"] == "logo"
    assert "logo" in blank.assets


def test_a_name_already_holding_a_different_picture_is_refused(blank: Session,
                                                               tmp_path: Path) -> None:
    """Silently rebinding a name would change every slide already using it."""
    _add(blank, _png(tmp_path / "a.png"), key="logo")
    result = _add(blank, _png(tmp_path / "b.png", colour=(0xC0, 0x20, 0x20)), key="logo")
    assert result["ok"] is False
    assert result["error"] == "key_taken"
    assert "another name" in result["message"]


def test_re_adding_the_same_picture_under_its_own_name_is_not_a_conflict(
    blank: Session, picture: Path
) -> None:
    assert _add(blank, picture, key="logo")["ok"]
    again = _add(blank, picture, key="logo")
    assert again["ok"] and again["after"]["reused"] is True


def test_a_key_that_could_not_be_a_filename_is_refused(blank: Session,
                                                       picture: Path) -> None:
    """A key is also a file under the workspace's `assets/`; a slash in one escapes it."""
    result = _add(blank, picture, key="../../etc/passwd")
    assert result["ok"] is False
    assert result["error"] == "bad_key"


# ------------------------------------------------------------------- the refusals


def test_a_non_image_is_refused_with_what_would_work(blank: Session,
                                                     tmp_path: Path) -> None:
    bad = tmp_path / "notes.txt"
    bad.write_text("not an image")
    result = _add(blank, bad)
    assert result["ok"] is False
    assert result["error"] == "not_an_image"
    assert ".png" in result["message"]
    assert not blank.assets, "a refused write must not have landed"


def test_the_media_type_comes_from_the_bytes_not_the_extension(blank: Session,
                                                               tmp_path: Path) -> None:
    """A `.png` that is really a JPEG is ordinary on a laptop full of downloads, and a
    package declaring the wrong content type is one PowerPoint offers to repair."""
    from PIL import Image

    lying = tmp_path / "actually-a-jpeg.png"
    Image.new("RGB", PICTURE, (0x15, 0x60, 0x82)).save(lying, format="JPEG")
    assert _add(blank, lying)["after"]["media_type"] == "image/jpeg"

    honest = tmp_path / "no-suffix"
    honest.write_bytes(_png(tmp_path / "src.png").read_bytes())
    assert _add(blank, honest)["after"]["media_type"] == "image/png"


def test_svg_is_refused_by_name_with_the_way_out(blank: Session, tmp_path: Path) -> None:
    """PowerPoint shows an SVG only as a vector part plus a rasterised fallback Office
    generates; python-pptx cannot even measure one. Its own code, because "rasterise it"
    is something the caller can act on and "not an image" is not."""
    svg = tmp_path / "diagram.svg"
    svg.write_text('<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10"/>')
    result = _add(blank, svg)
    assert result["ok"] is False
    assert result["error"] == "svg_unsupported"
    assert "PNG" in result["message"]


def test_a_missing_file_is_refused(blank: Session) -> None:
    result = router.dispatch(blank, "add_asset", {"path": "/nope/missing.png"})
    assert result["ok"] is False
    assert result["error"] == "no_such_image"


def test_an_image_past_the_limit_is_refused_by_the_same_number_import_uses(
    blank: Session, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """A deck whose assets are readable on import but unwritable by the tool would be a rule
    nobody could act on, so the limit is one constant with two users."""
    monkeypatch.setattr(media, "MAX_ASSET_BYTES", 512)
    result = _add(blank, _png(tmp_path / "big.png", size=(600, 400)))
    assert result["ok"] is False
    assert result["error"] == "image_too_large"
    assert "MB" in result["message"]


# --------------------------------------------------------------------- the file


def test_a_media_slot_naming_the_key_exports_the_picture(blank: Session,
                                                         tmp_path: Path,
                                                         picture: Path) -> None:
    """End to end, and on the file, because the file is what the recipient opens.

    Before `add_asset` this was impossible on a generated deck: nothing put bytes in the
    store, so the slot could only name a path and the deck stopped being self-contained.
    """
    key = _add(blank, picture)["after"]["asset_id"]
    _place_media(blank, key, alt="A blue rectangle")

    out = tmp_path / "media.pptx"
    report = export(blank.deck, out, assets=blank.assets)
    assert not report.violations

    found = _pictures(out)
    assert found, "the media slot exported no picture"
    assert found[0]._element._nvXxPr.cNvPr.get("descr") == "A blue rectangle"
    assert found[0].image.blob == picture.read_bytes(), "the file got different bytes"
    assert found[0].image.content_type == "image/png"


def test_the_picture_part_and_its_relationship_both_land(blank: Session, tmp_path: Path,
                                                         picture: Path) -> None:
    """A picture is two things in the package. One without the other is a broken file."""
    _place_media(blank, _add(blank, picture)["after"]["asset_id"])
    out = tmp_path / "parts.pptx"
    export(blank.deck, out, assets=blank.assets)

    with zipfile.ZipFile(out) as z:
        assert [n for n in z.namelist() if n.startswith("ppt/media/")]
        rels = z.read("ppt/slides/_rels/slide1.xml.rels").decode()
        assert "image" in rels


def test_the_export_and_the_preview_show_the_same_picture(blank: Session, tmp_path: Path,
                                                          picture: Path) -> None:
    """The invariant everything else here is measured against.

    Both sides now answer the same question from the same place — the deck's assets — which
    is exactly what removing the path fallback bought.
    """
    import base64

    key = _add(blank, picture)["after"]["asset_id"]
    _place_media(blank, key)
    out = tmp_path / "both.pptx"
    export(blank.deck, out, assets=blank.assets)

    markup = blank.render_html(blank.deck.slides[0].id, inline_assets=True)
    inlined = base64.b64encode(picture.read_bytes()).decode()
    assert inlined in markup, "the preview is not showing the picture the file has"
    assert _pictures(out)[0].image.blob == picture.read_bytes()


def test_a_slot_naming_a_path_is_refused_at_the_write(blank: Session,
                                                      picture: Path) -> None:
    """The fallback `add_asset` replaces, asserted gone — and caught at the cheapest moment.

    A path in a slot used to work on the machine that wrote it and nowhere else, silently.
    It is never a valid key, so the write is refused with the way out rather than left to
    surface as a violation at export.
    """
    result = router.dispatch(blank, "add_slide", {"layout": "stack", "blocks": [
        {"region": "body", "component": "image_full",
         "slots": {"media": {"asset_id": str(picture), "alt": "A blue rectangle"}}}]})
    assert result["ok"] is False
    assert result["error"] == "asset_is_a_path"
    assert "add_asset" in result["message"]
    assert blank.asset_data_uri(str(picture)) is None, "the preview still reads the disk"


def test_a_key_behind_nothing_is_reported_rather_than_drawn(blank: Session,
                                                            tmp_path: Path) -> None:
    """The honest failure, on both sides.

    A slot may legitimately name a key before `add_asset` has been called — authoring the
    slides before the pictures is an ordinary way to work — so this is not refused at the
    write. What must not happen is a successful-looking export with nothing on the slide.
    """
    from ppt_harness.io.writer_assertions import FidelityError

    _place_media(blank, "not-added-yet-1234abcd")
    out = tmp_path / "missing.pptx"
    with pytest.raises(FidelityError, match="slot_not_written"):
        export(blank.deck, out, assets=blank.assets)

    result = router.dispatch(blank, "export", {"path": str(out)})
    assert any("slot_not_written" in v for v in result["violations"]), result
    assert not _pictures(out)
    assert "missing" in blank.render_html(blank.deck.slides[0].id)


# --------------------------------------------------------------------- undo, state


def test_undo_removes_the_asset(blank: Session, picture: Path) -> None:
    key = _add(blank, picture)["after"]["asset_id"]
    assert router.dispatch(blank, "undo")["ok"]
    assert key not in blank.assets

    assert router.dispatch(blank, "redo")["ok"]
    assert blank.assets[key][1] == picture.read_bytes(), "redo lost the bytes"


def test_the_op_log_does_not_carry_the_image(blank: Session, picture: Path) -> None:
    """An op that carried base64 would put megabytes into every `journal.jsonl` line and hold
    them for the life of the session, twice over, for something the workspace already writes
    to `assets/`. The op names a digest; the store's pool holds the picture."""
    import json

    _add(blank, picture)
    op = blank.store.log.ops[-1]
    assert op.op == "add_asset"
    line = json.dumps(op.model_dump(mode="json"))
    assert len(line) < 400, "the bytes leaked into the op"
    assert op.patch["sha1"] and "content_type" in op.patch


def test_assets_stay_out_of_the_document_model(blank: Session, picture: Path) -> None:
    """The deck is dumped whole for every invertible op; image bytes cannot ride along."""
    _add(blank, picture)
    assert "assets" not in blank.deck.model_dump(mode="json")


def test_an_added_asset_survives_a_restore(blank: Session, tmp_path: Path,
                                           picture: Path) -> None:
    """Assets used to be written once, at attach, on the grounds that nothing mutates an
    imported image. A picture added mid-session then lived only in memory, and a restart
    resumed a deck whose slide named an asset that was no longer there."""
    from ppt_harness.state.workspace import Workspace

    workspace = Workspace(tmp_path / "ws")
    workspace.attach(blank.store)
    key = _add(blank, picture)["after"]["asset_id"]
    _place_media(blank, key)

    restored, report = workspace.restore()
    assert report["resumed"]
    assert restored.assets[key][1] == picture.read_bytes()


# ------------------------------------------------------- the tools that use a key


def test_add_image_ingests_rather_than_pointing(imported: Session, picture: Path) -> None:
    """`add_image` recorded a path and nothing else, so its picture was in the exported file
    and absent from every render the model was verifying against — `render/html.py` draws
    `shape.asset`, which was empty. It now files the bytes and names the key."""
    slide_id = next(s.id for s in imported.deck.slides if s.mode is Mode.FREEFORM)
    result = router.dispatch(imported, "add_image",
                             {"slide_id": slide_id, "region": "left",
                              "path": str(picture), "alt": "A blue rectangle"})
    assert result["ok"], result.get("message")

    key = result["after"]["asset_id"]
    assert imported.assets[key][1] == picture.read_bytes()

    shape = next(s for s in imported.slide(slide_id).shapes if s.asset == key)
    assert shape.source, "provenance is still recorded"
    markup = imported.render_html(slide_id, inline_assets=True)
    assert imported.asset_data_uri(key).split(",")[1] in markup


def test_add_image_and_add_asset_agree_on_a_picture(imported: Session,
                                                    picture: Path) -> None:
    """One picture, one set of bytes, one key — whichever door it came through."""
    slide_id = next(s.id for s in imported.deck.slides if s.mode is Mode.FREEFORM)
    key = router.dispatch(imported, "add_asset", {"path": str(picture)})["after"]["asset_id"]
    added = router.dispatch(imported, "add_image",
                            {"slide_id": slide_id, "region": "left",
                             "path": str(picture), "alt": "x"})
    assert added["after"]["asset_id"] == key


def test_undo_after_add_image_takes_the_shape_and_the_asset_together(
    imported: Session, picture: Path
) -> None:
    """One turn. Half of it landing leaves either a picture nothing points at or a shape
    pointing at nothing."""
    slide_id = next(s.id for s in imported.deck.slides if s.mode is Mode.FREEFORM)
    before = len(imported.slide(slide_id).shapes)
    key = router.dispatch(imported, "add_image",
                          {"slide_id": slide_id, "region": "left", "path": str(picture),
                           "alt": "x"})["after"]["asset_id"]

    assert router.dispatch(imported, "undo")["ok"]
    assert key not in imported.assets
    assert len(imported.slide(slide_id).shapes) == before


def test_a_stream_is_fresh_for_every_slot_naming_the_same_asset(blank: Session,
                                                                tmp_path: Path,
                                                                picture: Path) -> None:
    """python-pptx reads the file object to the end, so a second slot handed the same buffer
    would be written as a zero-byte picture part."""
    key = _add(blank, picture)["after"]["asset_id"]
    _place_media(blank, key)
    _place_media(blank, key)

    out = tmp_path / "twice.pptx"
    export(blank.deck, out, assets=blank.assets)
    for index in (0, 1):
        assert _pictures(out, index)[0].image.blob == picture.read_bytes()
    # One image part, not two: the package keys media by content, and so does the store.
    with zipfile.ZipFile(out) as z:
        assert len([n for n in z.namelist() if n.startswith("ppt/media/")]) == 1


def test_the_probe_reads_proportions_through_the_writers_own_decoder(tmp_path: Path) -> None:
    """Two answers to "how wide is this image" is the shape of the bug that makes an ejected
    slide stop matching the managed one it came from."""
    blob = _png(tmp_path / "wide.png", size=(400, 100)).read_bytes()
    read = media.probe(blob)
    assert read is not None
    assert read.aspect == pytest.approx(4.0, rel=1e-3)
    assert media.probe(b"not an image at all") is None
    assert media.probe(io.BytesIO().getvalue()) is None


def test_list_assets_recovers_a_key_the_model_lost(blank: Session, tmp_path: Path) -> None:
    """A content-addressed key is unguessable by design, so losing one has to be survivable.

    The alternatives without this are re-ingesting a file whose path the model may not have
    kept, or naming a key that is not there and finding out at export.
    """
    one = _png(tmp_path / "wide.png", size=(600, 200))
    key = router.dispatch(blank, "add_asset", {"path": str(one)})["after"]["asset_id"]

    listed = router.dispatch(blank, "list_assets")
    assert [a["asset_id"] for a in listed["assets"]] == [key]
    assert listed["count"] == 1
    entry = listed["assets"][0]
    assert (entry["width_px"], entry["height_px"], entry["aspect"]) == (600, 200, 3.0), \
        "the aspect ratio is the fact that decides whether a picture will letterbox"


def test_list_assets_never_returns_the_bytes(blank: Session, tmp_path: Path) -> None:
    """This lands in a model's context. A base64 image there would cost more than every
    other read tool in the harness together."""
    router.dispatch(blank, "add_asset", {"path": str(_png(tmp_path / "p.png"))})

    listed = router.dispatch(blank, "list_assets")
    assert len(str(listed)) < 400, "the listing is carrying image data"
    assert all(not isinstance(v, bytes) for a in listed["assets"] for v in a.values())


def test_an_empty_deck_lists_nothing_rather_than_failing(blank: Session) -> None:
    """Silence and emptiness read the same to a model, so the count is stated."""
    listed = router.dispatch(blank, "list_assets")
    assert listed["assets"] == [] and listed["count"] == 0


def test_an_added_picture_survives_its_source_file_being_deleted(
        imported: Session, tmp_path: Path) -> None:
    """The last route that depended on a path outside the deck.

    `_add_shape` wrote a freeform picture from `shape.source` — where the file was read from
    on the machine that read it — so a deck exported after the original moved or was mailed
    onward carried a picture nobody else could see. The bytes are in the store now; the path
    is provenance. Deleting the file before exporting is the whole test.
    """
    source = _png(tmp_path / "vanishing.png")
    slide_id = imported.deck.slides[0].id
    added = router.dispatch(imported, "add_image", {
        "slide_id": slide_id, "path": str(source), "alt": "A square", "region": "body"})
    assert added["ok"], added.get("message")

    source.unlink()
    out = router.dispatch(imported, "export", {"path": str(tmp_path / "gone.pptx")})
    assert out["ok"], out.get("violations")

    with zipfile.ZipFile(tmp_path / "gone.pptx") as bundle:
        assert [n for n in bundle.namelist() if n.startswith("ppt/media/")], \
            "the picture did not reach the file"


def test_removing_an_asset_a_slide_uses_is_refused(blank: Session, picture: Path) -> None:
    """An asset dropped from under a slide does not fail here — it fails at export, as a
    slot the writer could not build, on a slide that looked finished. Refusing at the gate is
    the same trade the budget makes: the cheapest failure is the one that never rendered.

    The refusal names the places, because "it is in use" without saying where leaves the
    caller to go and search for them.
    """
    key = router.dispatch(blank, "add_asset", {"path": str(picture)})["after"]["asset_id"]
    router.dispatch(blank, "add_slide", {"layout": "stack", "blocks": [
        {"region": "body", "component": "image_full", "variant": "full",
         "slots": {"media": {"asset_id": key, "alt": "A picture"}}}]})

    refused = router.dispatch(blank, "remove_asset", {"key": key})
    assert refused["ok"] is False
    assert refused["error"] == "asset_in_use"
    assert blank.deck.slides[0].id in refused["message"], "the refusal has to say where"
    assert key in blank.store.assets, "a refused removal removed something"


def test_an_unused_asset_can_be_taken_back(blank: Session, tmp_path: Path) -> None:
    """What the tool is actually for: an ingest the caller wants to undo — a wrong file, or a
    duplicate under a name they would rather reuse. Not housekeeping; `export` writes only
    referenced media, so an orphan costs the exported deck nothing."""
    key = router.dispatch(
        blank, "add_asset", {"path": str(_png(tmp_path / "spare.png"))})["after"]["asset_id"]
    removed = router.dispatch(blank, "remove_asset", {"key": key})
    assert removed["ok"] and key not in blank.store.assets
    assert router.dispatch(blank, "list_assets")["count"] == 0


def test_removing_an_asset_is_undoable(blank: Session, tmp_path: Path) -> None:
    """The bytes are kept, not just the key — otherwise undo restores a name behind nothing
    and the next export reports a slot it could not build."""
    path = _png(tmp_path / "spare.png")
    key = router.dispatch(blank, "add_asset", {"path": str(path)})["after"]["asset_id"]
    original = blank.store.assets[key][1]

    router.dispatch(blank, "remove_asset", {"key": key})
    router.dispatch(blank, "undo")

    assert blank.store.assets[key][1] == original, "undo restored the key but not the bytes"


def test_removing_an_asset_that_is_not_there_says_so(blank: Session) -> None:
    """Silence would read as success, and the caller would carry on believing a key is gone
    that was never present under that name."""
    result = router.dispatch(blank, "remove_asset", {"key": "not-a-key"})
    assert result["ok"] is False and result["error"] == "no_such_asset"
