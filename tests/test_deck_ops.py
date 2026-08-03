"""Deck-level operations — PLAN A5.

Changes to the deck rather than to anything on a slide. Three claims carry weight:

**A duplicate is a copy of the part, not a rebuild from the model.** The harness does not
understand every shape it can duplicate, so a copy assembled from what it *does* understand
would quietly drop the SmartArt.

**Resizing the canvas invalidates absolute geometry, and says so.** Every imported frame is
EMU measured against the old canvas. Rescaling silently would produce a deck that looks
right here and wrong to whoever opens it.

**A theme role is re-validated as a whole.** The theme is checked once at load so managed
slides *cannot* fail contrast; changing a role without rechecking retires that guarantee.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from ppt_harness.core.session import Session
from ppt_harness.state.document import Mode
from ppt_harness.tools import router


def _export(session: Session, tmp_path: Path, name: str = "deck.pptx") -> Path:
    out = tmp_path / name
    result = router.dispatch(session, "export", {"path": str(out)})
    assert result["ok"], result.get("message")
    return out


def _kinds(path: Path, index: int) -> list[str]:
    return [str(s.shape_type).split()[0]
            for s in Presentation(str(path)).slides[index].shapes]


# ------------------------------------------------------------------ duplicate


def test_a_duplicate_lands_after_the_original(imported: Session) -> None:
    source = imported.deck.slides[1]
    before = len(imported.deck.slides)
    result = router.dispatch(imported, "duplicate_slide", {"slide_id": source.id})
    assert result["ok"]
    assert len(imported.deck.slides) == before + 1
    assert imported.deck.slides[2].id == result["target"]


def test_every_id_is_reissued(imported: Session) -> None:
    """A copy sharing shape ids with its original makes every later edit ambiguous — the
    harness would find two shapes for one name and pick whichever came first."""
    source = imported.deck.slides[1]
    original_ids = {s.id for s in source.shapes}
    result = router.dispatch(imported, "duplicate_slide", {"slide_id": source.id})
    copy = imported.deck.slide(result["target"])

    assert copy.id != source.id
    assert not ({s.id for s in copy.shapes} & original_ids)
    assert len({s.id for s in copy.shapes}) == len(copy.shapes)


def test_a_duplicate_copies_the_part_not_the_model(imported: Session,
                                                   tmp_path: Path) -> None:
    """The load-bearing claim: shapes the harness does not model survive the copy."""
    source = next((s for s in imported.deck.slides
                   if any(x.opaque for x in s.shapes) and any(x.asset for x in s.shapes)),
                  None)
    if source is None:
        pytest.skip("fixture has no slide with both an opaque shape and an asset")

    index = source.index
    router.dispatch(imported, "duplicate_slide", {"slide_id": source.id})
    out = _export(imported, tmp_path)

    assert _kinds(out, index) == _kinds(out, index + 1)
    copy = Session.open(out).deck.slides[index + 1]
    assert any(s.opaque for s in copy.shapes), "the opaque shape did not survive"
    assert any(s.asset for s in copy.shapes), "the picture did not survive"


def test_slide_order_survives_a_duplicate(imported: Session, tmp_path: Path) -> None:
    """The regression this replaces: reordering assumed the package held imported slides
    first, in original order. A duplicate is appended at the end, so every slide after the
    copy silently took the wrong content."""
    titles = [(s.shapes[0].text if s.shapes else "") for s in imported.deck.slides]
    source = imported.deck.slides[1]
    router.dispatch(imported, "duplicate_slide", {"slide_id": source.id})

    out = _export(imported, tmp_path)
    reopened = Session.open(out)
    got = [(s.shapes[0].text if s.shapes else "") for s in reopened.deck.slides]

    assert got == [*titles[:2], titles[1], *titles[2:]]


def test_undo_removes_a_duplicate(imported: Session) -> None:
    before = len(imported.deck.slides)
    router.dispatch(imported, "duplicate_slide",
                    {"slide_id": imported.deck.slides[0].id})
    assert router.dispatch(imported, "undo")["ok"]
    assert len(imported.deck.slides) == before


# ----------------------------------------------------------------------- hide


def test_hiding_keeps_the_slide_in_the_deck(imported: Session, tmp_path: Path) -> None:
    """Hidden, not removed — "drop this for now" means skip it, not lose it."""
    slide = imported.deck.slides[1]
    before = len(imported.deck.slides)
    assert router.dispatch(imported, "hide_slide", {"slide_id": slide.id})["ok"]
    assert slide.hidden is True
    assert len(imported.deck.slides) == before

    out = _export(imported, tmp_path)
    assert Session.open(out).deck.slides[1].hidden is True


def test_a_shown_slide_writes_no_attribute(imported: Session, tmp_path: Path) -> None:
    """Absent means shown. Writing the default is churn the round-trip guarantee then has
    to account for."""
    out = _export(imported, tmp_path, "untouched.pptx")
    for slide in Presentation(str(out)).slides:
        assert slide._element.get("show") is None


def test_hiding_twice_is_refused(imported: Session) -> None:
    slide = imported.deck.slides[1]
    router.dispatch(imported, "hide_slide", {"slide_id": slide.id})
    again = router.dispatch(imported, "hide_slide", {"slide_id": slide.id})
    assert again["ok"] is False
    assert again["error"] == "already_there"


def test_a_slide_can_be_shown_again(imported: Session) -> None:
    slide = imported.deck.slides[1]
    router.dispatch(imported, "hide_slide", {"slide_id": slide.id})
    assert router.dispatch(imported, "hide_slide",
                           {"slide_id": slide.id, "hidden": False})["ok"]
    assert slide.hidden is False


# --------------------------------------------------------------------- layout


def test_layout_is_managed_only(imported: Session) -> None:
    """A freeform slide has no layout frame; its shapes are where its author put them."""
    result = router.dispatch(imported, "set_layout",
                             {"slide_id": imported.deck.slides[0].id, "layout": "stack"})
    assert result["ok"] is False
    assert result["error"] == "wrong_mode"


def test_a_layout_without_a_home_for_a_block_is_refused(populated: Session) -> None:
    """Dropping a block silently would lose content, so the refusal names the blocks."""
    managed = next(s for s in populated.deck.slides if s.mode is Mode.MANAGED)
    result = router.dispatch(populated, "set_layout",
                             {"slide_id": managed.id, "layout": "title"})
    assert result["ok"] is False
    assert result["error"] in ("region_missing", "unknown_layout")
    if result["error"] == "region_missing":
        assert managed.blocks[0].id in result["message"]


# ----------------------------------------------------------------- slide size


def test_resizing_reports_what_it_invalidated(imported: Session) -> None:
    """Imported frames are absolute EMU against the old canvas. Rescaling silently would
    look right here and wrong to whoever opens it."""
    result = router.dispatch(imported, "set_slide_size", {"preset": "4:3"})
    assert result["ok"]
    assert result["invalidated"] > 0
    assert "not rescaled" in result["note"]
    assert tuple(imported.theme.grid.canvas) == (960, 720)


def test_resizing_to_the_current_size_is_refused(imported: Session) -> None:
    current = tuple(imported.theme.grid.canvas)
    preset = "16:9" if current == (1280, 720) else "4:3"
    router.dispatch(imported, "set_slide_size", {"preset": preset})
    again = router.dispatch(imported, "set_slide_size", {"preset": preset})
    assert again["ok"] is False
    assert again["error"] == "already_there"


def test_an_unknown_preset_names_the_ones_that_exist(imported: Session) -> None:
    result = router.dispatch(imported, "set_slide_size", {"preset": "widescreen"})
    assert result["ok"] is False
    assert "16:9" in result["message"]


def test_undo_restores_the_canvas(imported: Session) -> None:
    before = tuple(imported.theme.grid.canvas)
    router.dispatch(imported, "set_slide_size", {"preset": "4:3"})
    assert router.dispatch(imported, "undo")["ok"]
    assert tuple(imported.theme.grid.canvas) == before


# ---------------------------------------------------------------- theme roles


def test_a_role_can_be_corrected(imported: Session) -> None:
    """Theme extraction guesses several roles; this is how a person fixes one."""
    result = router.dispatch(imported, "set_theme_role",
                             {"role": "brand", "value": "#0B5FA5"})
    assert result["ok"], result.get("message")
    assert imported.theme.palette["brand"] == "#0B5FA5"


def test_correcting_a_role_clears_its_inferred_flag(imported: Session) -> None:
    """`inferred` exists so a user can see what was guessed. A corrected role is no longer
    a guess."""
    router.dispatch(imported, "set_theme_role",
                    {"role": "ink_muted", "value": "#5A6472"})
    assert "palette.ink_muted" not in imported.theme.inferred


def test_a_dependent_role_is_re_derived(imported: Session) -> None:
    """`brand_ink` exists only in relation to `brand`. Leaving it behind fails contrast
    against a colour nobody chose."""
    router.dispatch(imported, "set_theme_role", {"role": "brand", "value": "#111111"})
    from ppt_harness.state.theme_default import contrast

    palette = imported.theme.palette
    assert contrast(str(palette["brand_ink"]), str(palette["brand"])) >= 4.5


def test_a_change_that_breaks_contrast_is_refused_with_the_ratio(imported: Session) -> None:
    """The theme is validated once at load so managed slides cannot fail contrast. A role
    changed without rechecking would quietly retire that guarantee."""
    result = router.dispatch(imported, "set_theme_role",
                             {"role": "ink", "value": "#FEFEFE"})
    assert result["ok"] is False
    assert result["error"] == "contrast_failed"
    assert result["problems"]


@pytest.mark.parametrize("value", ["red", "#12345", "0B5FA5", "#GGGGGG"])
def test_a_colour_that_is_not_a_colour_is_refused(imported: Session, value: str) -> None:
    result = router.dispatch(imported, "set_theme_role",
                             {"role": "brand", "value": value})
    assert result["ok"] is False
    assert result["error"] == "bad_colour"


def test_an_ordered_list_role_is_refused(imported: Session) -> None:
    """`accents` is a sequence — item N of a list slot takes accent N — so setting it to a
    single colour would silently collapse the palette."""
    result = router.dispatch(imported, "set_theme_role",
                             {"role": "accents", "value": "#0B5FA5"})
    assert result["ok"] is False
    assert result["error"] == "not_a_single_colour"


def test_undo_restores_the_theme(imported: Session) -> None:
    before = dict(imported.theme.palette)
    router.dispatch(imported, "set_theme_role", {"role": "brand", "value": "#0B5FA5"})
    assert router.dispatch(imported, "undo")["ok"]
    assert imported.theme.palette == before


# ------------------------------------------------------------------- package


@pytest.mark.parametrize("operation", ["duplicate", "hide", "resize", "theme"])
def test_deck_operations_leave_the_package_valid(imported: Session, tmp_path: Path,
                                                 operation: str) -> None:
    slide_id = imported.deck.slides[1].id
    if operation == "duplicate":
        router.dispatch(imported, "duplicate_slide", {"slide_id": slide_id})
    elif operation == "hide":
        router.dispatch(imported, "hide_slide", {"slide_id": slide_id})
    elif operation == "resize":
        router.dispatch(imported, "set_slide_size", {"preset": "4:3"})
    else:
        router.dispatch(imported, "set_theme_role",
                        {"role": "brand", "value": "#0B5FA5"})

    out = _export(imported, tmp_path, f"{operation}.pptx")
    reopened = Session.open(out)
    assert len(reopened.deck.slides) == len(imported.deck.slides)
