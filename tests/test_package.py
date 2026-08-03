"""Package-level audit and repair — `io/package.py`.

Below the document model: parts, content types, and the relationship graph. python-pptx
cannot see these, which is why an edited deck rots here and nowhere else — a deleted picture
leaves its relationship and its image behind, PowerPoint does not complain, and the debris
compounds across sessions.

Every test here exists because a draft got it wrong:

- The first audit reported three problems on a *pristine* file, having modelled dangling
  relationships as a blocklist of types that are unreferenced by design. A validator that
  cries wolf on a clean file is worse than none.
- The first sweep deleted every relationship that audit called dangling, and promptly
  severed `notesSlide -> slide` and dropped `printerSettings`.

So the two things pinned hardest are: clean files stay clean, and the sweep touches only
what it was built for.
"""

from __future__ import annotations

import re
import zipfile

import pytest

from ppt_harness.core.session import Session
from ppt_harness.io import package
from ppt_harness.tools import router


@pytest.fixture
def edited(tmp_path, fixture_deck):
    """The fixture with one picture deleted — the defect this module exists for."""
    session = Session.open(fixture_deck)
    found = next(((sl, sh) for sl in session.deck.slides for sh in sl.shapes
                  if sh.type == "picture"), None)
    if found is None:
        pytest.skip("fixture has no picture to delete")
    _, shape = found
    assert router.dispatch(session, "delete_shape", {"shape_id": shape.id})["ok"]
    out = tmp_path / "edited.pptx"
    assert router.dispatch(session, "export", {"path": str(out)})["ok"]
    return out


# ------------------------------------------------------------------ no false alarms


def test_a_pristine_deck_audits_clean(fixture_deck) -> None:
    """The one that matters most. Every check added here must pass this, or the report
    becomes noise and the noise trains people to skip it."""
    report = package.audit(fixture_deck)
    assert report.clean, [str(f) for f in report.findings]


def test_relationships_bound_by_design_are_not_dangling(fixture_deck) -> None:
    """`notesSlide -> slide`, `presentation -> printerSettings`, layouts, masters and themes
    are reached by relationship alone and cited nowhere in XML. An earlier draft flagged all
    of them."""
    kinds = {f.kind for f in package.audit(fixture_deck).findings}
    assert "dangling_relationship" not in kinds


# ---------------------------------------------------------------- catching the real thing


def test_a_deleted_picture_leaves_a_dangling_relationship(edited) -> None:
    findings = package.audit(edited).findings
    assert [f.kind for f in findings] == ["dangling_relationship"]
    assert "image" in findings[0].detail


def test_the_source_is_baselined_away(edited, fixture_deck) -> None:
    """A deck that arrived broken must not have its inherited faults billed to this edit —
    the failure mode of every linter nobody runs."""
    report = package.audit(edited, original=fixture_deck)
    assert [f.kind for f in report.findings] == ["dangling_relationship"]


def test_baselining_reports_what_it_suppressed(tmp_path, fixture_deck) -> None:
    """Silently dropping inherited problems would be its own dishonesty."""
    broken = tmp_path / "broken.pptx"
    broken.write_bytes(fixture_deck.read_bytes())
    report = package.audit(broken, original=fixture_deck)
    assert report.clean
    assert report.as_dict()["problems"] == []


# ------------------------------------------------------------------------- the sweep


def test_sweep_removes_the_orphan_and_its_relationship(edited) -> None:
    removed = package.sweep(edited)
    assert any("image" in r for r in removed)
    assert package.audit(edited).clean
    with zipfile.ZipFile(edited) as zf:
        assert not [n for n in zf.namelist() if n.startswith("ppt/media/")]


def test_sweep_never_severs_a_notes_backlink(edited) -> None:
    """The bug that made the first sweep dangerous: `notesSlide -> slide` is unreferenced in
    XML by design, and cutting it orphans the speaker notes."""
    package.sweep(edited)
    with zipfile.ZipFile(edited) as zf:
        rels = [n for n in zf.namelist() if n.startswith("ppt/notesSlides/_rels/")]
        for name in rels:
            body = zf.read(name).decode()
            assert "../slides/slide" in body, f"{name} lost its backlink to its slide"


def test_sweep_leaves_the_presentation_relationships_alone(edited, fixture_deck) -> None:
    def rel_count(path):
        with zipfile.ZipFile(path) as zf:
            return zf.read("ppt/_rels/presentation.xml.rels").decode().count("<Relationship")

    before = rel_count(fixture_deck)
    package.sweep(edited)
    assert rel_count(edited) == before


def test_a_swept_deck_still_opens(edited) -> None:
    """The only check that matters if the others are wrong."""
    package.sweep(edited)
    from pptx import Presentation

    assert len(Presentation(str(edited)).slides) > 0


def test_sweeping_a_clean_deck_changes_nothing(tmp_path, fixture_deck) -> None:
    copy = tmp_path / "copy.pptx"
    copy.write_bytes(fixture_deck.read_bytes())
    before = copy.read_bytes()
    assert package.sweep(copy) == []
    assert copy.read_bytes() == before, "a no-op sweep must not even rewrite the zip"


# ------------------------------------------------------------------------- structure


def test_a_missing_target_is_reported_not_swept(tmp_path, fixture_deck) -> None:
    """A part that should exist and does not is a problem to report. Deleting the
    relationship would hide it."""
    broken = tmp_path / "missing.pptx"
    with zipfile.ZipFile(fixture_deck) as src, zipfile.ZipFile(broken, "w") as out:
        for item in src.infolist():
            if item.filename.startswith("ppt/media/"):
                continue  # drop the image, keep the relationship pointing at it
            out.writestr(item, src.read(item.filename))
    kinds = [f.kind for f in package.audit(broken).findings]
    assert "missing_target" in kinds


def _mutated(src, dest, fn):
    """A copy of `src` with one thing deliberately broken.

    Breaking a real package beats hand-writing a minimal one: a synthetic file passes checks
    for reasons that have nothing to do with the check under test, and then the check looks
    fine while being unable to fire on anything real.
    """
    with zipfile.ZipFile(src) as z, zipfile.ZipFile(dest, "w") as out:
        for item in z.infolist():
            body = fn(item.filename, z.read(item.filename))
            if body is not None:
                out.writestr(item, body)
    return dest


def test_a_slide_missing_from_the_presentation_list_is_reported(tmp_path, fixture_deck):
    """The part exists and is related, so every per-relationship check passes — and the
    slide is invisible to a reader while still counting toward the file."""
    broken = _mutated(fixture_deck, tmp_path / "a.pptx", lambda n, b:
                      re.sub(rb"<p:sldId[^>]*/>", b"", b, count=1)
                      if n == "ppt/presentation.xml" else b)
    assert "slide_not_in_presentation" in {f.kind for f in package.audit(broken).findings}


def test_a_duplicate_relationship_id_is_reported(tmp_path, fixture_deck) -> None:
    """PowerPoint resolves an ambiguous rId one way; python-pptx may resolve it the other."""
    def dup(name, body):
        if name != "ppt/_rels/presentation.xml.rels":
            return body
        first = re.search(rb"<Relationship[^>]*/>", body)
        return body.replace(first.group(0), first.group(0) * 2, 1)

    broken = _mutated(fixture_deck, tmp_path / "b.pptx", dup)
    assert "duplicate_relationship_id" in {f.kind for f in package.audit(broken).findings}


def test_a_mislabelled_picture_is_reported(tmp_path, fixture_deck) -> None:
    """No relationship breaks and no schema is violated — PowerPoint just renders nothing
    where the picture was, which is why nothing else catches it."""
    broken = _mutated(fixture_deck, tmp_path / "c.pptx", lambda n, b:
                      b.replace(b'ContentType="image/png"', b'ContentType="image/jpeg"')
                      if n == "[Content_Types].xml" else b)
    findings = {f.kind for f in package.audit(broken).findings}
    assert "content_type_mismatch" in findings


def test_a_target_escaping_the_package_is_reported(tmp_path, fixture_deck) -> None:
    broken = _mutated(fixture_deck, tmp_path / "d.pptx", lambda n, b:
                      b.replace(b'Target="../media/', b'Target="../../../../etc/')
                      if n.endswith("slide6.xml.rels") else b)
    assert "target_escapes_package" in {f.kind for f in package.audit(broken).findings}


def test_a_truncated_part_is_reported(tmp_path, fixture_deck) -> None:
    """The realistic writer failure — a half-finished write, a bad encoding — and the reason
    full schema validation is not worth vendoring megabytes of XSDs for: our writes go
    through python-pptx's object model, which cannot emit misordered elements, so what is
    left to catch is a part that does not parse at all.
    """
    broken = _mutated(fixture_deck, tmp_path / "trunc.pptx", lambda n, b:
                      b[:len(b) // 2] if n == "ppt/slides/slide2.xml" else b)
    findings = package.audit(broken).findings
    assert [f.kind for f in findings] == ["malformed_xml"]
    assert findings[0].part == "ppt/slides/slide2.xml"


def test_a_package_with_no_content_types_says_so(tmp_path) -> None:
    empty = tmp_path / "empty.pptx"
    with zipfile.ZipFile(empty, "w") as zf:
        zf.writestr("ppt/presentation.xml", "<p:presentation/>")
    findings = package.audit(empty).findings
    assert findings[0].kind == "no_content_types"
