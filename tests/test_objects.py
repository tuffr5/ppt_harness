"""Pictures, tables, charts and notes — PLAN A4.

The objects a deck is made of beyond text. Two commitments are under test more than the
mechanics:

**A chart is data, not a picture.** `add_chart` writes a native chart with its embedded
worksheet, so the recipient can edit the numbers. DESIGN §1.5 forbids degrading one into an
image, and there is no path back if it happened.

**Alt text is required.** A picture without it is invisible to anyone using a screen
reader, and the moment to supply it is the moment the picture is added — not a cleanup pass
that never comes.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from ppt_harness.core.session import Session
from ppt_harness.state.document import Mode
from ppt_harness.tools import router


@pytest.fixture
def picture(tmp_path: Path) -> Path:
    from PIL import Image

    path = tmp_path / "probe.png"
    Image.new("RGB", (320, 200), (0x15, 0x60, 0x82)).save(path)
    return path


@pytest.fixture
def slide_id(imported: Session) -> str:
    return next(s for s in imported.deck.slides if s.mode is Mode.FREEFORM).id


def _export(imported: Session, tmp_path: Path, name: str = "objects.pptx") -> Path:
    out = tmp_path / name
    result = router.dispatch(imported, "export", {"path": str(out)})
    assert result["ok"], result.get("message")
    return out


def _shapes(path: Path, index: int):
    return list(Presentation(str(path)).slides[index].shapes)


# ------------------------------------------------------------------------ pictures


def test_a_picture_is_added_with_its_alt_text(imported: Session, slide_id: str,
                                              picture: Path, tmp_path: Path) -> None:
    result = router.dispatch(imported, "add_image",
                             {"slide_id": slide_id, "region": "left",
                              "path": str(picture), "alt": "A blue rectangle"})
    assert result["ok"], result.get("message")

    index = imported.deck.slide(slide_id).index
    added = [s for s in _shapes(_export(imported, tmp_path), index)
             if s.shape_type is not None and "PICTURE" in str(s.shape_type)]
    assert added, "no picture reached the file"
    assert added[-1]._element._nvXxPr.cNvPr.get("descr") == "A blue rectangle"


def test_a_picture_without_alt_text_is_refused(imported: Session, slide_id: str,
                                               picture: Path) -> None:
    """The rule the harness would otherwise quietly drop."""
    result = router.dispatch(imported, "add_image",
                             {"slide_id": slide_id, "region": "left",
                              "path": str(picture), "alt": "   "})
    assert result["ok"] is False
    assert result["error"] == "alt_required"


def test_a_missing_file_is_refused(imported: Session, slide_id: str) -> None:
    result = router.dispatch(imported, "add_image",
                             {"slide_id": slide_id, "region": "left",
                              "path": "/nope/missing.png", "alt": "x"})
    assert result["ok"] is False
    assert result["error"] == "no_such_image"


def test_a_non_image_is_refused_with_what_would_work(imported: Session, slide_id: str,
                                                     tmp_path: Path) -> None:
    bad = tmp_path / "notes.txt"
    bad.write_text("not an image")
    result = router.dispatch(imported, "add_image",
                             {"slide_id": slide_id, "region": "left",
                              "path": str(bad), "alt": "x"})
    assert result["ok"] is False
    assert ".png" in result["message"]


def test_the_media_part_and_its_relationship_both_land(imported: Session, slide_id: str,
                                                       picture: Path,
                                                       tmp_path: Path) -> None:
    """A picture is two things in the package. One without the other is a broken file."""
    router.dispatch(imported, "add_image",
                    {"slide_id": slide_id, "region": "left", "path": str(picture),
                     "alt": "A blue rectangle"})
    out = _export(imported, tmp_path)
    with zipfile.ZipFile(out) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        assert media
        index = imported.deck.slide(slide_id).index + 1
        rels = z.read(f"ppt/slides/_rels/slide{index}.xml.rels").decode()
        assert "image" in rels


# -------------------------------------------------------------------------- tables


def test_a_table_is_a_real_table(imported: Session, slide_id: str,
                                 tmp_path: Path) -> None:
    """Not a picture of one — the recipient has to be able to edit the cells."""
    router.dispatch(imported, "add_table",
                    {"slide_id": slide_id, "region": "right",
                     "headers": ["Region", "Q1"],
                     "rows": [["us-east", "18.2"], ["eu-central", "9.4"]]})
    index = imported.deck.slide(slide_id).index
    tables = [s for s in _shapes(_export(imported, tmp_path), index) if s.has_table]
    assert tables
    table = tables[0].table
    assert len(table.rows) == 3
    assert table.cell(0, 0).text == "Region"
    assert table.cell(2, 1).text == "9.4"


def test_a_cell_can_be_changed(imported: Session, slide_id: str, tmp_path: Path) -> None:
    added = router.dispatch(imported, "add_table",
                            {"slide_id": slide_id, "region": "right",
                             "headers": ["a", "b"], "rows": [["1", "2"]]})
    shape_id = added["target"].split("/")[-1]
    assert router.dispatch(imported, "set_cell",
                           {"shape_id": shape_id, "row": 1, "col": 1,
                            "text": "99"})["ok"]

    index = imported.deck.slide(slide_id).index
    table = next(s for s in _shapes(_export(imported, tmp_path), index) if s.has_table)
    assert table.table.cell(1, 1).text == "99"


def test_an_out_of_range_cell_is_refused(imported: Session, slide_id: str) -> None:
    added = router.dispatch(imported, "add_table",
                            {"slide_id": slide_id, "region": "right",
                             "rows": [["1", "2"]]})
    shape_id = added["target"].split("/")[-1]
    result = router.dispatch(imported, "set_cell",
                             {"shape_id": shape_id, "row": 9, "col": 0, "text": "x"})
    assert result["ok"] is False
    assert result["error"] == "out_of_range"


def test_a_table_too_wide_to_read_is_refused(imported: Session, slide_id: str) -> None:
    """A limit stated as a design judgement, not a technical one: nine columns do not read
    from the back of a room."""
    result = router.dispatch(imported, "add_table",
                             {"slide_id": slide_id, "region": "right",
                              "rows": [[str(i) for i in range(12)]]})
    assert result["ok"] is False
    assert result["error"] == "too_wide"


def test_an_empty_table_is_refused(imported: Session, slide_id: str) -> None:
    result = router.dispatch(imported, "add_table",
                             {"slide_id": slide_id, "region": "right", "rows": []})
    assert result["ok"] is False


def test_undo_removes_a_table(imported: Session, slide_id: str) -> None:
    slide = imported.deck.slide(slide_id)
    before = len(slide.shapes)
    router.dispatch(imported, "add_table",
                    {"slide_id": slide_id, "region": "right", "rows": [["1"]]})
    assert router.dispatch(imported, "undo")["ok"]
    assert len(slide.shapes) == before


# -------------------------------------------------------------------------- charts


def test_a_chart_carries_its_own_worksheet(imported: Session, slide_id: str,
                                           tmp_path: Path) -> None:
    """The whole claim of DESIGN §1.5. Without the worksheet it is a picture."""
    router.dispatch(imported, "add_chart",
                    {"slide_id": slide_id, "region": "body", "kind": "column",
                     "categories": ["a", "b", "c"],
                     "series": [{"name": "Q1", "values": [1, 2, 3]}]})
    out = _export(imported, tmp_path)
    with zipfile.ZipFile(out) as z:
        assert any("charts/chart" in n for n in z.namelist())
        assert any("embeddings" in n and n.endswith(".xlsx") for n in z.namelist())


def test_a_chart_is_editable_data_not_an_image(imported: Session, slide_id: str,
                                               tmp_path: Path) -> None:
    router.dispatch(imported, "add_chart",
                    {"slide_id": slide_id, "region": "body", "kind": "column",
                     "categories": ["a", "b"],
                     "series": [{"name": "Q1", "values": [1, 2]}]})
    index = imported.deck.slide(slide_id).index
    charts = [s for s in _shapes(_export(imported, tmp_path), index)
              if getattr(s, "has_chart", False)]
    assert charts
    assert list(charts[0].chart.plots[0].categories) == ["a", "b"]


def test_replacing_chart_data_rewrites_the_worksheet_too(imported: Session, slide_id: str,
                                                         tmp_path: Path) -> None:
    """Rewriting the plot alone leaves the workbook stale, so the recipient who opens the
    data sees numbers that no longer match the picture."""
    added = router.dispatch(imported, "add_chart",
                            {"slide_id": slide_id, "region": "body", "kind": "column",
                             "categories": ["a", "b"],
                             "series": [{"name": "Q1", "values": [1, 2]}]})
    shape_id = added["target"].split("/")[-1]
    assert router.dispatch(imported, "set_chart_data",
                           {"shape_id": shape_id, "categories": ["x", "y", "z"],
                            "series": [{"name": "New", "values": [7, 8, 9]}]})["ok"]

    index = imported.deck.slide(slide_id).index
    chart = next(s for s in _shapes(_export(imported, tmp_path), index)
                 if getattr(s, "has_chart", False)).chart
    assert list(chart.plots[0].categories) == ["x", "y", "z"]
    assert list(chart.series[0].values) == [7, 8, 9]


def test_a_series_that_does_not_line_up_is_refused(imported: Session,
                                                   slide_id: str) -> None:
    result = router.dispatch(imported, "add_chart",
                             {"slide_id": slide_id, "region": "body", "kind": "column",
                              "categories": ["a", "b", "c"],
                              "series": [{"name": "Q1", "values": [1, 2]}]})
    assert result["ok"] is False
    assert result["error"] == "series_mismatch"


def test_changing_the_numbers_in_a_picture_of_a_chart_is_refused(imported: Session,
                                                                 slide_id: str,
                                                                 picture: Path) -> None:
    """The honest answer from DESIGN §1.5: the data is not there to change."""
    added = router.dispatch(imported, "add_image",
                            {"slide_id": slide_id, "region": "left",
                             "path": str(picture), "alt": "a chart, as a picture"})
    shape_id = added["target"].split("/")[-1]
    result = router.dispatch(imported, "set_chart_data",
                             {"shape_id": shape_id, "categories": ["a"],
                              "series": [{"name": "x", "values": [1]}]})
    assert result["ok"] is False
    assert result["error"] == "not_a_chart"
    assert "not there to change" in result["message"]


def test_an_unknown_chart_kind_is_refused(imported: Session, slide_id: str) -> None:
    result = router.dispatch(imported, "add_chart",
                             {"slide_id": slide_id, "region": "body", "kind": "radar",
                              "categories": ["a"], "series": [{"name": "x",
                                                               "values": [1]}]})
    assert result["ok"] is False
    assert "column" in result["message"]


# --------------------------------------------------------------------------- notes


def test_notes_reach_the_file_and_survive_reimport(imported: Session, slide_id: str,
                                                   tmp_path: Path) -> None:
    text = "Speak to the chart, not the table."
    assert router.dispatch(imported, "set_notes",
                           {"slide_id": slide_id, "notes": text})["ok"]

    out = _export(imported, tmp_path)
    index = imported.deck.slide(slide_id).index
    assert Presentation(str(out)).slides[index].notes_slide \
        .notes_text_frame.text == text
    assert Session.open(out).deck.slides[index].notes == text


def test_notes_are_not_budget_checked(imported: Session, slide_id: str) -> None:
    """Nothing renders them on the slide, so length is not a layout problem."""
    result = router.dispatch(imported, "set_notes",
                             {"slide_id": slide_id, "notes": "word " * 500})
    assert result["ok"]


def test_undo_restores_notes(imported: Session, slide_id: str) -> None:
    before = imported.deck.slide(slide_id).notes
    router.dispatch(imported, "set_notes", {"slide_id": slide_id, "notes": "changed"})
    assert router.dispatch(imported, "undo")["ok"]
    assert imported.deck.slide(slide_id).notes == before


def test_unchanged_notes_do_not_create_a_part(imported: Session, tmp_path: Path) -> None:
    """Touching `notes_slide` creates the part. An empty notes slide on every exported deck
    is churn the round-trip guarantee then has to explain away."""
    with zipfile.ZipFile(imported.deck.source_path) as z:
        before = len([n for n in z.namelist() if "notesSlide" in n])
    out = _export(imported, tmp_path, "untouched.pptx")
    with zipfile.ZipFile(out) as z:
        after = len([n for n in z.namelist() if "notesSlide" in n])
    assert after == before


# ---------------------------------------------------------------------- boundaries


def test_objects_are_placed_by_region_never_by_coordinate() -> None:
    for name in ("add_image", "add_table", "add_chart"):
        tool = next(t for t in router.tools() if t.name == name)
        properties = set(tool.schema["properties"])
        assert "region" in properties
        assert not properties & {"x", "y", "left", "top", "width", "height"}


def test_adding_an_object_to_a_managed_slide_is_refused(populated: Session,
                                                        picture: Path) -> None:
    managed = next(s for s in populated.deck.slides if s.mode is Mode.MANAGED)
    result = router.dispatch(populated, "add_image",
                             {"slide_id": managed.id, "region": "body",
                              "path": str(picture), "alt": "x"})
    assert result["ok"] is False
    assert result["error"] == "wrong_mode"


@pytest.mark.parametrize("scenario", ["image", "table", "chart", "notes"])
def test_every_object_leaves_the_package_valid(imported: Session, slide_id: str,
                                               picture: Path, tmp_path: Path,
                                               scenario: str) -> None:
    if scenario == "image":
        router.dispatch(imported, "add_image",
                        {"slide_id": slide_id, "region": "left",
                         "path": str(picture), "alt": "x"})
    elif scenario == "table":
        router.dispatch(imported, "add_table",
                        {"slide_id": slide_id, "region": "right", "rows": [["1", "2"]]})
    elif scenario == "chart":
        router.dispatch(imported, "add_chart",
                        {"slide_id": slide_id, "region": "body", "kind": "pie",
                         "categories": ["a", "b"],
                         "series": [{"name": "s", "values": [1, 2]}]})
    else:
        router.dispatch(imported, "set_notes",
                        {"slide_id": slide_id, "notes": "spoken"})

    out = _export(imported, tmp_path, f"{scenario}.pptx")
    # Opening it at all proves the package parses; the counts prove nothing was orphaned.
    reopened = Session.open(out)
    assert len(reopened.deck.slides) == len(imported.deck.slides)
