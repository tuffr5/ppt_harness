"""The component catalog — DESIGN §3, PLAN B1.

Sixteen components, each tested at the bounds it declares for itself. The corpus is
generated from the registry rather than written out, so a component added tomorrow is
covered tomorrow and a declaration that drifts from reality fails here rather than on
someone's slide.

Two claims carry the weight:

**A declared maximum is a promise.** `max_items: 8` means eight items fit. If they do not,
the number is wrong — and the number is what the model reads before deciding what to build.

**Degradation preserves content.** A chain may only run between components whose slot
shapes overlap, so the repair ladder can walk it without dropping a column on the floor.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest
from pptx import Presentation

from ppt_harness.components import fixtures, overrides, registry
from ppt_harness.core.session import Session
from ppt_harness.tools import router

CASES = fixtures.every_case()
KEYS = sorted(registry.COMPONENTS)


def _case_id(case: tuple[str, str, str, bool]) -> str:
    key, variant, script, optional = case
    return f"{key}-{variant}-{script}" + ("" if optional else "-required-only")


def _add(session: Session, key: str, variant: str, slots: dict) -> dict:
    layout, region = fixtures.region_for(key)
    return router.dispatch(session, "add_slide", {"layout": layout, "blocks": [
        {"region": region, "component": key, "variant": variant, "slots": slots}]})


# ------------------------------------------------------------------- the corpus


@pytest.mark.parametrize("case", CASES, ids=[_case_id(c) for c in CASES])
def test_a_component_honours_the_maximum_it_declares(blank: Session, case) -> None:
    """`max_items: 8` is a promise the model reads before deciding what to build. A
    component that cannot keep it has the wrong number, not a hard slide."""
    key, variant, script, optional = case
    result = _add(blank, key, variant,
                  fixtures.payload(key, script, optional=optional))
    assert result["ok"], result.get("message")


@pytest.mark.parametrize("case", CASES, ids=[_case_id(c) for c in CASES])
def test_content_past_the_maximum_is_refused(blank: Session, case) -> None:
    """The cheapest failure in the system, and it has to actually fire."""
    key, variant, script, _ = case
    if key in ("image_full",):
        pytest.skip("a picture scales; it has no text to overflow")
    result = _add(blank, key, variant, fixtures.oversized(key, script))
    assert result["ok"] is False
    assert result["error"] == "budget_exceeded"
    assert "options:" in result["message"], "a rejection must say what to do instead"


@pytest.mark.parametrize("key", KEYS)
def test_every_component_has_a_home(key: str) -> None:
    """A component declaring a region no layout provides can never be placed."""
    layout, region = fixtures.region_for(key)
    assert region in registry.layout(layout).regions


@pytest.mark.parametrize("key", KEYS)
def test_every_component_describes_itself(key: str) -> None:
    comp = registry.get(key)
    assert comp.purpose.strip(), f"{key} has no purpose line"
    assert comp.variants, f"{key} has no variants"
    assert comp.slots, f"{key} has no slots"
    assert any(s.required for s in comp.slots.values()), f"{key} requires nothing"


@pytest.mark.parametrize("key", KEYS)
def test_no_component_carries_a_font_size(key: str) -> None:
    """Slots name a theme role. A size here is the moment the type scale stops meaning
    anything."""
    for spec in registry.get(key).slots.values():
        assert isinstance(spec.role, str) and spec.role


# --------------------------------------------------------------------- chains


@pytest.mark.parametrize("key", KEYS)
def test_degradation_preserves_content(key: str) -> None:
    """A chain that dropped a slot shape would be a data loss dressed as a layout change."""
    comp = registry.get(key)
    required = {spec.shape for spec in comp.slots.values() if spec.required}
    for nxt in registry.degradation_chain(key):
        available = set(registry.get(nxt).slot_shapes().values())
        assert required <= available, f"{key} -> {nxt} drops {sorted(required - available)}"


@pytest.mark.parametrize("key", KEYS)
def test_a_chain_terminates(key: str) -> None:
    chain = registry.degradation_chain(key)
    assert len(chain) == len(set(chain)), f"{key} cycles: {chain}"
    assert len(chain) < len(registry.COMPONENTS)


def test_the_ladder_has_somewhere_to_go() -> None:
    """If nothing degraded, the repair ladder would have exactly one rung — shorten the
    words — which is the rung it is meant to avoid."""
    with_chains = [k for k in KEYS if registry.degradation_chain(k)]
    assert len(with_chains) >= 5, with_chains


@pytest.mark.parametrize("key", KEYS)
def test_a_swap_cannot_lose_a_required_slot(key: str) -> None:
    """`set_component` is only allowed within a slot shape. Anything else drops a slot on
    the floor, which is a data loss dressed as a layout change."""
    required = {spec.shape for spec in registry.get(key).slots.values() if spec.required}
    for other in registry.swappable(key):
        assert required <= set(registry.get(other).slot_shapes().values())


# ----------------------------------------------------------------- the objects


def test_a_tabular_slot_is_written_as_a_real_table(blank: Session,
                                                   tmp_path: Path) -> None:
    """The bug this replaces: a managed `data_table` exported a Python dict repr onto the
    slide, because the writer only knew how to write text."""
    assert _add(blank, "data_table", "plain",
                {"tabular": {"headers": ["Region", "Q1"],
                             "rows": [["us-east", "18.2"]]}})["ok"]
    out = tmp_path / "table.pptx"
    router.dispatch(blank, "export", {"path": str(out)})

    shapes = list(Presentation(str(out)).slides[0].shapes)
    assert any(s.has_table for s in shapes), "no real table was written"
    for shape in shapes:
        if shape.has_text_frame:
            assert "headers" not in shape.text_frame.text, "a dict repr reached the slide"


def test_a_chart_slot_is_written_with_its_worksheet(blank: Session,
                                                    tmp_path: Path) -> None:
    """A chart without its worksheet is a picture, and DESIGN §1.5 forbids the harness
    producing one."""
    assert _add(blank, "chart", "wide",
                {"chart": {"kind": "bar", "labels": ["a", "b"],
                           "series": [{"name": "Q1", "values": [1, 2]}]}})["ok"]
    out = tmp_path / "chart.pptx"
    router.dispatch(blank, "export", {"path": str(out)})

    shapes = list(Presentation(str(out)).slides[0].shapes)
    assert any(getattr(s, "has_chart", False) for s in shapes)
    with zipfile.ZipFile(out) as z:
        assert any("embeddings" in n and n.endswith(".xlsx") for n in z.namelist())


def test_objects_are_not_audited_as_text_boxes(blank: Session, tmp_path: Path) -> None:
    """A table's cells and a chart's plot are PowerPoint's to lay out — the one place it is
    allowed to, because the alternative costs the recipient an editable object."""
    _add(blank, "data_table", "plain",
         {"tabular": {"headers": ["a"], "rows": [["1"]]}})
    result = router.dispatch(blank, "export", {"path": str(tmp_path / "audit.pptx")})
    assert result["ok"], result.get("message")
    assert "fidelity clean" in result["summary"]


# -------------------------------------------------------------- bounded overrides


@pytest.mark.parametrize("key", sorted(overrides.OVERRIDES))
def test_every_override_is_clamped_at_both_ends(key: str) -> None:
    """The property that makes these safe to hand a model: nothing here can be set to a
    value that collides slots or leaves the palette."""
    knob = overrides.OVERRIDES[key]
    for wild in (-999, 999, "nonsense", None, 1e9):
        clamped = knob.clamp(wild)
        assert clamped in knob.values or (
            isinstance(clamped, (int, float))
            and min(knob.values) <= clamped <= max(knob.values)
        ), f"{key}.clamp({wild!r}) = {clamped!r}"


def test_an_index_override_stays_an_index() -> None:
    """`accent: 2.0` reads as a scale factor in the op log, and indexes a palette by a
    float."""
    assert overrides.OVERRIDES["accent"].clamp(2) == 2
    assert isinstance(overrides.OVERRIDES["accent"].clamp("3"), int)


def test_unknown_knobs_are_dropped_not_rejected() -> None:
    """A model that invents `spacing` should still get a slide."""
    assert overrides.clamp_all({"spacing": 4, "density": "compact"}) == {
        "density": "compact"}


def test_no_override_carries_a_colour_or_a_size() -> None:
    """An override holding a hex value or a point size would be `set_font` renamed."""
    for knob in overrides.OVERRIDES.values():
        for value in knob.values:
            assert not (isinstance(value, str) and value.startswith("#"))
        assert "size" not in knob.key or knob.key == "media_scale"


def test_emphasis_resolves_through_the_theme(blank: Session) -> None:
    """It names an ink *role*; the palette answers."""
    theme = blank.theme
    quiet = overrides.ink_for(theme, {"emphasis": "quiet"})
    strong = overrides.ink_for(theme, {"emphasis": "strong"})
    normal = overrides.ink_for(theme, {})
    assert {quiet, strong, normal} <= set(
        v for v in theme.palette.values() if isinstance(v, str))
    assert quiet != normal != strong


def test_an_accent_override_wraps_within_the_palette(blank: Session) -> None:
    accents = blank.theme.palette["accents"]
    assert overrides.accent_for(blank.theme, {"accent": 99}) in accents


def test_density_reaches_the_geometry(blank: Session) -> None:
    from ppt_harness.state.document import Block, Mode, Slide

    slots = {"title": "A heading", "items": ["one", "two", "three"]}
    slide = Slide(id="s", index=0, mode=Mode.MANAGED, layout="stack", blocks=[
        Block(id="b", region="body", component="card_grid", variant="2x2", slots=slots)])
    from ppt_harness.render import expand

    normal = [s.box.h for s in expand.expand_slide(blank.theme, slide)]
    slide.blocks[0].overrides = {"density": "compact"}
    compact = [s.box.h for s in expand.expand_slide(blank.theme, slide)]
    assert sum(compact) > sum(normal), "compact bought no room"


def test_the_override_tool_reports_a_clamp(blank: Session) -> None:
    """Asking for 9.0 and being told "already 1.0" is confusing until you know why."""
    added = router.dispatch(blank, "add_slide", {"layout": "stack", "blocks": [
        {"region": "body", "component": "bullets", "slots": {"items": ["a", "b"]}}]})
    slide_id = added["target"]
    block_id = blank.deck.slide(slide_id).blocks[0].id

    result = router.dispatch(blank, "set_override",
                             {"slide_id": slide_id, "block_id": block_id,
                              "key": "accent", "value": 99})
    assert result["ok"]
    assert result["clamped_from"] == 99
    assert blank.deck.slide(slide_id).blocks[0].overrides["accent"] == 5


def test_the_design_catalog_table_matches_the_registry() -> None:
    """DESIGN §3's table is a transcription, and transcriptions rot.

    This one had listed `image-bg`, `with-kicker` and `lead-in` — variants the code has never
    had — while calling `image_full`'s pair `bleed` and `inset` at a time when neither drew
    anything. A model reads these names as a contract and will ask for what the document
    offers, so a catalog that overstates fails at the tool gate for a reason the reader was
    told was legal.

    Compared as sets per component, because the table also carries emphasis and prose that
    are a document's business rather than the registry's.
    """
    import re
    from pathlib import Path

    from ppt_harness.components import registry

    design = (Path(__file__).resolve().parent.parent / "DESIGN.md").read_text()
    # Sliced to §3 first: several other tables in this document have the same shape, and a
    # loose pattern picks up modes, export modes and skill kinds as though they were
    # components.
    section = design[design.index("## 3. Component catalog"):]
    section = section[:section.index("\n\n", section.index("|---"))]
    rows = re.findall(r"^\| `(\w+)` \| [^|]* \| ([^|]*) \|", section, re.M)
    documented = {key: {v.strip().strip("*").strip() for v in variants.split(",")}
                  for key, variants in rows}

    assert set(documented) == set(registry.COMPONENTS), (
        f"only in DESIGN: {sorted(set(documented) - set(registry.COMPONENTS))}; "
        f"only in the registry: {sorted(set(registry.COMPONENTS) - set(documented))}"
    )
    for key, component in registry.COMPONENTS.items():
        assert documented[key] == set(component.variants), (
            f"{key}: DESIGN says {sorted(documented[key])}, "
            f"the registry has {sorted(component.variants)}"
        )
