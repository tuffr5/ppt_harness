"""Paragraph and run properties — PLAN A1.

Everything here is something OOXML already stores and the harness had no way to say.
Alignment and list style live on `<a:pPr>`; hyperlinks hang off `<a:rPr>`, the same run
mechanism emphasis uses.

The boundary these tests defend: **which words** and **what shape** are the model's to say;
**how big** and **what face** are the theme's. `set_align` exists and `set_font` does not,
and that is not an accident of what got built first.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from ppt_harness.core.session import Session
from ppt_harness.state import richtext as rt
from ppt_harness.tools import router

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


@pytest.fixture
def target(imported: Session) -> str:
    """A fresh box with room, so a budget refusal never masks the property under test."""
    slide = imported.deck.slides[1]
    result = router.dispatch(imported, "add_textbox",
                             {"slide_id": slide.id, "region": "body",
                              "text": "alpha\nbeta\ngamma"})
    assert result["ok"], result.get("message")
    return result["target"]


def _shape(imported: Session, target: str):
    _, shape, _ = imported.store.resolve_text_target(target)
    return shape


def _export(imported: Session, tmp_path: Path) -> Presentation:
    out = tmp_path / "props.pptx"
    assert router.dispatch(imported, "export", {"path": str(out)})["ok"]
    return Presentation(str(out))


def _paragraphs(prs: Presentation, needle: str):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and needle in shape.text_frame.text:
                return list(shape.text_frame.paragraphs)
    return []


# ----------------------------------------------------------------------- set_align


@pytest.mark.parametrize("align", ["center", "right", "justify", "left"])
def test_alignment_round_trips_to_ooxml(imported: Session, target: str, tmp_path: Path,
                                        align: str) -> None:
    assert router.dispatch(imported, "set_align",
                           {"target": target, "align": align})["ok"]
    assert _shape(imported, target).align == align

    expected = {"left": None, "center": "ctr", "right": "r", "justify": "just"}[align]
    for para in _paragraphs(_export(imported, tmp_path), "alpha"):
        p_pr = para._p.find(f"{{{A}}}pPr")
        assert (p_pr.get("algn") if p_pr is not None else None) == expected


def test_an_unknown_alignment_is_refused(imported: Session, target: str) -> None:
    result = router.dispatch(imported, "set_align",
                             {"target": target, "align": "middle"})
    assert result["ok"] is False
    assert "center" in result["message"], "a refusal should name the alternatives"


def test_undo_restores_the_alignment(imported: Session, target: str) -> None:
    before = _shape(imported, target).align
    router.dispatch(imported, "set_align", {"target": target, "align": "right"})
    assert router.dispatch(imported, "undo")["ok"]
    assert _shape(imported, target).align == before


# ------------------------------------------------------------------------ set_list


def test_a_bulleted_list_writes_buchar(imported: Session, target: str,
                                       tmp_path: Path) -> None:
    assert router.dispatch(imported, "set_list",
                           {"target": target, "kind": "bullet", "level": 1})["ok"]
    for para in _paragraphs(_export(imported, tmp_path), "alpha"):
        p_pr = para._p.find(f"{{{A}}}pPr")
        assert p_pr.get("lvl") == "1"
        assert p_pr.find(f"{{{A}}}buChar") is not None


def test_a_numbered_list_writes_buautonum(imported: Session, target: str,
                                          tmp_path: Path) -> None:
    router.dispatch(imported, "set_list", {"target": target, "kind": "number"})
    for para in _paragraphs(_export(imported, tmp_path), "alpha"):
        assert para._p.find(f"{{{A}}}pPr").find(f"{{{A}}}buAutoNum") is not None


def test_clearing_a_list_writes_bunone_explicitly(imported: Session, target: str,
                                                  tmp_path: Path) -> None:
    """Omitting the declaration inherits the layout's bullet, which is how a paragraph the
    user cleared comes back bulleted."""
    router.dispatch(imported, "set_list", {"target": target, "kind": "bullet"})
    router.dispatch(imported, "set_list", {"target": target, "kind": "none"})
    for para in _paragraphs(_export(imported, tmp_path), "alpha"):
        p_pr = para._p.find(f"{{{A}}}pPr")
        assert p_pr.find(f"{{{A}}}buNone") is not None
        assert p_pr.find(f"{{{A}}}buChar") is None


def test_only_one_bullet_declaration_survives_a_change(imported: Session, target: str,
                                                       tmp_path: Path) -> None:
    """A paragraph carries at most one. Appending without clearing leaves both, and
    PowerPoint's choice between them is not ours to predict."""
    for kind in ("bullet", "number", "bullet"):
        router.dispatch(imported, "set_list", {"target": target, "kind": kind})
    for para in _paragraphs(_export(imported, tmp_path), "alpha"):
        p_pr = para._p.find(f"{{{A}}}pPr")
        declarations = sum(len(p_pr.findall(f"{{{A}}}{tag}"))
                           for tag in ("buNone", "buChar", "buAutoNum"))
        assert declarations == 1


def test_an_out_of_range_level_is_refused(imported: Session, target: str) -> None:
    assert router.dispatch(imported, "set_list",
                           {"target": target, "kind": "bullet", "level": 9})["ok"] is False


def test_the_preview_draws_a_real_list(imported: Session, target: str) -> None:
    """Markers drawn by the renderer, not typed into the text — otherwise the words the
    harness measures are not the words on the slide."""
    router.dispatch(imported, "set_list", {"target": target, "kind": "bullet"})
    markup = imported.render_html(target.split("/")[0])
    assert "<ul" in markup and "<li>" in markup
    assert "•" not in markup, "the bullet must not be a character in the text"


# ------------------------------------------------------------------------ set_link


def test_a_link_is_attached_to_the_span(imported: Session, target: str,
                                        tmp_path: Path) -> None:
    assert router.dispatch(imported, "set_link",
                           {"target": target, "span": "gamma",
                            "url": "https://example.com"})["ok"]

    linked = [r for r in _shape(imported, target).runs if r.link]
    assert [r.text for r in linked] == ["gamma"]

    for para in _paragraphs(_export(imported, tmp_path), "gamma"):
        for run in para.runs:
            r_pr = run._r.find(f"{{{A}}}rPr")
            has_link = r_pr is not None and r_pr.find(f"{{{A}}}hlinkClick") is not None
            assert has_link == (run.text == "gamma")


def test_a_link_can_be_removed(imported: Session, target: str) -> None:
    router.dispatch(imported, "set_link",
                    {"target": target, "span": "gamma", "url": "https://example.com"})
    router.dispatch(imported, "set_link", {"target": target, "span": "gamma", "url": ""})
    assert not any(r.link for r in _shape(imported, target).runs)


def test_a_url_powerpoint_cannot_follow_is_refused(imported: Session, target: str) -> None:
    result = router.dispatch(imported, "set_link",
                             {"target": target, "span": "gamma", "url": "javascript:x"})
    assert result["ok"] is False
    assert result["error"] == "bad_url"


def test_a_span_that_is_not_there_is_refused_with_the_text(imported: Session,
                                                           target: str) -> None:
    """So the model can correct itself rather than guess again."""
    result = router.dispatch(imported, "set_link",
                             {"target": target, "span": "delta", "url": "https://x.com"})
    assert result["ok"] is False
    assert "alpha" in result["message"]


# -------------------------------------------------------------------- set_emphasis


def test_emphasis_applies_to_a_span_without_rewriting_the_text(imported: Session,
                                                               target: str) -> None:
    """Rewriting the whole string to change one word is how unrelated edits get lost."""
    before = _shape(imported, target).text
    assert router.dispatch(imported, "set_emphasis",
                           {"target": target, "span": "beta", "mark": "bold"})["ok"]
    shape = _shape(imported, target)
    assert shape.text == before
    assert [r.text for r in shape.runs if r.bold] == ["beta"]


def test_emphasis_can_be_turned_off(imported: Session, target: str) -> None:
    router.dispatch(imported, "set_emphasis",
                    {"target": target, "span": "beta", "mark": "bold"})
    router.dispatch(imported, "set_emphasis",
                    {"target": target, "span": "beta", "mark": "bold", "on": False})
    assert not any(r.bold for r in _shape(imported, target).runs)


def test_none_clears_every_mark(imported: Session, target: str) -> None:
    for mark in ("bold", "italic", "underline"):
        router.dispatch(imported, "set_emphasis",
                        {"target": target, "span": "beta", "mark": mark})
    router.dispatch(imported, "set_emphasis",
                    {"target": target, "span": "beta", "mark": "none"})
    assert all(r.plain or r.link for r in _shape(imported, target).runs)


def test_every_occurrence_is_marked(imported: Session) -> None:
    slide = imported.deck.slides[1]
    added = router.dispatch(imported, "add_textbox",
                            {"slide_id": slide.id, "region": "body",
                             "text": "one two one"})["target"]
    result = router.dispatch(imported, "set_emphasis",
                             {"target": added, "span": "one", "mark": "italic"})
    assert result["after"]["occurrences"] == 2


def test_marks_apply_by_text_not_by_offset(imported: Session, target: str) -> None:
    """A model that has just read a slide knows the words, not the character offsets — and
    offsets go stale the moment anything else changes."""
    tool = next(t for t in router.tools() if t.name == "set_emphasis")
    properties = set(tool.schema["properties"])
    assert "span" in properties
    assert not properties & {"start", "end", "offset", "index"}


# ---------------------------------------------------------------------- boundaries


def test_managed_slots_take_their_properties_from_the_component(populated: Session) -> None:
    """A component owns its own alignment; overriding it per call is how decks drift."""
    from ppt_harness.state.document import Mode

    slide = next(s for s in populated.deck.slides if s.mode is Mode.MANAGED)
    block = slide.blocks[0]
    slot = next(iter(block.slots))
    result = router.dispatch(populated, "set_align",
                             {"target": f"{slide.id}/{block.id}/{slot}",
                              "align": "center"})
    assert result["ok"] is False
    assert result["error"] in ("managed_slot", "wrong_mode")


def test_no_text_property_tool_exposes_a_font_or_a_size() -> None:
    names = {"set_align", "set_list", "set_link", "set_emphasis"}
    for tool in router.tools():
        if tool.name not in names:
            continue
        blob = str(tool.schema).lower()
        assert "font" not in blob
        assert "size" not in blob


def test_properties_survive_a_full_round_trip(imported: Session, target: str,
                                              tmp_path: Path) -> None:
    router.dispatch(imported, "set_align", {"target": target, "align": "center"})
    router.dispatch(imported, "set_list", {"target": target, "kind": "bullet"})
    router.dispatch(imported, "set_emphasis",
                    {"target": target, "span": "beta", "mark": "bold"})

    out = tmp_path / "round.pptx"
    router.dispatch(imported, "export", {"path": str(out)})
    reopened = Session.open(out)

    shape = next(s for slide in reopened.deck.slides for s in slide.shapes
                 if s.text and "alpha" in s.text)
    assert shape.align == "center"
    assert any(r.bold and r.text == "beta" for r in shape.runs)
    assert rt.to_plain(shape.runs).replace("\n", " ").strip() == "alpha beta gamma"
