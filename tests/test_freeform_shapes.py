"""Shape-level editing on imported slides — PLAN A2.

Until these tools existed an imported deck accepted exactly one verb, `set_text`, which is
not enough to call it editable.

The load-bearing claim is **reversibility**. The harness can delete shapes it does not
understand, so undo has to restore the whole thing — including its position in the stack,
because document order *is* z-order in OOXML and a shape that comes back on top of what it
used to sit behind is not the same slide.
"""

from __future__ import annotations

import re
import zipfile
from pathlib import Path

import pytest
from lxml import etree

from ppt_harness.core.session import Session
from ppt_harness.state.document import Mode
from ppt_harness.tools import router

P = "http://schemas.openxmlformats.org/presentationml/2006/main"


@pytest.fixture
def slide(imported: Session):
    """The busiest freeform slide — the one with something worth deleting."""
    return max((s for s in imported.deck.slides if s.mode is Mode.FREEFORM),
               key=lambda s: len(s.shapes))


def _editable(slide):
    return next(s for s in slide.shapes if s.text and not s.opaque and s.ooxml_id)


# ------------------------------------------------------------------ delete_shape


def test_deleting_a_shape_removes_it_from_the_model(imported: Session, slide) -> None:
    victim = _editable(slide)
    before = len(slide.shapes)
    assert router.dispatch(imported, "delete_shape", {"shape_id": victim.id})["ok"]
    assert len(slide.shapes) == before - 1
    assert slide.shape(victim.id) is None


def test_a_deletion_is_recorded_for_the_exporter(imported: Session, slide) -> None:
    """Export patches the original package, so a deletion is an instruction rather than
    something the in-memory model can carry on its own."""
    victim = _editable(slide)
    router.dispatch(imported, "delete_shape", {"shape_id": victim.id})
    assert victim.ooxml_id in slide.removed


def test_undo_restores_a_deleted_shape_in_its_original_position(imported: Session,
                                                                slide) -> None:
    """Document order is z-order. A shape that returns on top of what it sat behind is not
    the shape that was deleted."""
    victim = slide.shapes[len(slide.shapes) // 2]
    index = slide.shapes.index(victim)
    router.dispatch(imported, "delete_shape", {"shape_id": victim.id})
    assert router.dispatch(imported, "undo")["ok"]

    restored = slide.shape(victim.id)
    assert restored is not None
    assert slide.shapes.index(restored) == index
    assert restored.model_dump(mode="json") == victim.model_dump(mode="json")


def test_an_opaque_shape_can_be_deleted_and_restored(imported: Session) -> None:
    """The harness deletes what it does not understand, so undo must not depend on
    understanding it."""
    found = next(((s, x) for s in imported.deck.slides for x in s.shapes if x.opaque),
                 None)
    if found is None:
        pytest.skip("fixture has no opaque shapes")
    slide, victim = found
    snapshot = victim.model_dump(mode="json")
    router.dispatch(imported, "delete_shape", {"shape_id": victim.id})
    router.dispatch(imported, "undo")
    assert slide.shape(victim.id).model_dump(mode="json") == snapshot


def test_deleting_from_a_managed_slide_is_refused(populated: Session) -> None:
    managed = next(s for s in populated.deck.slides if s.mode is Mode.MANAGED)
    result = router.dispatch(populated, "delete_shape", {"shape_id": f"{managed.id}_x"})
    assert result["ok"] is False


def test_the_deleted_shape_is_gone_from_the_exported_file(imported: Session, slide,
                                                          tmp_path: Path) -> None:
    victim = _editable(slide)
    words = victim.text
    router.dispatch(imported, "delete_shape", {"shape_id": victim.id})
    out = tmp_path / "deleted.pptx"
    assert router.dispatch(imported, "export", {"path": str(out)})["ok"]

    from pptx import Presentation

    texts = [s.text_frame.text for s in Presentation(str(out)).slides[slide.index].shapes
             if s.has_text_frame]
    assert not any(words and words in t for t in texts)


# ------------------------------------------------------------------- add_textbox


def test_a_text_box_is_placed_by_region_not_by_coordinate(imported: Session,
                                                          slide) -> None:
    from ppt_harness.render import expand

    result = router.dispatch(imported, "add_textbox",
                             {"slide_id": slide.id, "region": "body", "text": "Placed"})
    assert result["ok"], result.get("message")

    added = slide.shapes[-1]
    cx, cy = imported.slide_size_emu()
    canvas = imported.theme.grid.canvas
    expected = expand.region_by_name(imported.theme, "body").emu(*canvas, cx, cy)
    assert (added.frame.x, added.frame.y, added.frame.cx, added.frame.cy) == expected


def test_an_unknown_region_is_refused_with_the_ones_that_exist(imported: Session,
                                                               slide) -> None:
    result = router.dispatch(imported, "add_textbox",
                             {"slide_id": slide.id, "region": "nowhere", "text": "x"})
    assert result["ok"] is False
    assert "body" in result["message"], "a refusal should name the alternatives"


def test_the_region_enum_matches_the_registry() -> None:
    """Hand-listing regions lets the tool and the expander drift apart."""
    from ppt_harness.components import registry

    tool = next(t for t in router.tools() if t.name == "add_textbox")
    offered = set(tool.schema["properties"]["region"]["enum"])
    known = {n for frame in registry.LAYOUTS.values() for n in frame.regions}
    assert offered == known


def test_an_added_box_understands_emphasis(imported: Session, slide) -> None:
    router.dispatch(imported, "add_textbox",
                    {"slide_id": slide.id, "region": "body", "text": "a **b** c"})
    added = slide.shapes[-1]
    assert added.text == "a b c"
    assert any(r.bold for r in added.runs)


def test_an_added_box_reaches_the_exported_file(imported: Session, slide,
                                                tmp_path: Path) -> None:
    router.dispatch(imported, "add_textbox",
                    {"slide_id": slide.id, "region": "body",
                     "text": "Written by the harness"})
    out = tmp_path / "added.pptx"
    router.dispatch(imported, "export", {"path": str(out)})

    from pptx import Presentation

    texts = [s.text_frame.text for s in Presentation(str(out)).slides[slide.index].shapes
             if s.has_text_frame]
    assert any("Written by the harness" in t for t in texts)


def test_undo_removes_an_added_box(imported: Session, slide) -> None:
    before = len(slide.shapes)
    router.dispatch(imported, "add_textbox",
                    {"slide_id": slide.id, "region": "body", "text": "temporary"})
    assert router.dispatch(imported, "undo")["ok"]
    assert len(slide.shapes) == before


# --------------------------------------------------------------- duplicate_shape


def test_a_duplicate_is_offset_so_it_can_be_seen(imported: Session, slide) -> None:
    original = _editable(slide)
    assert router.dispatch(imported, "duplicate_shape", {"shape_id": original.id})["ok"]
    copy = slide.shapes[slide.shapes.index(original) + 1]
    assert copy.id != original.id
    assert (copy.frame.x, copy.frame.y) != (original.frame.x, original.frame.y)
    assert copy.text == original.text


def test_a_duplicate_is_new_to_the_file(imported: Session, slide) -> None:
    """`ooxml_id` 0 is what tells the exporter to create it rather than patch something."""
    original = _editable(slide)
    router.dispatch(imported, "duplicate_shape", {"shape_id": original.id})
    copy = slide.shapes[slide.shapes.index(original) + 1]
    assert copy.ooxml_id == 0
    assert copy.dirty, "a shape not in the file must be written on export"


def test_an_opaque_shape_is_not_duplicated(imported: Session) -> None:
    """It cannot be copied faithfully, and a silent approximation is worse than a refusal."""
    found = next(((s, x) for s in imported.deck.slides for x in s.shapes if x.opaque), None)
    if found is None:
        pytest.skip("fixture has no opaque shapes")
    _, shape = found
    result = router.dispatch(imported, "duplicate_shape", {"shape_id": shape.id})
    assert result["ok"] is False
    assert result["error"] == "shape_opaque"


# ------------------------------------------------------------------- set_z_order


@pytest.mark.parametrize(("position", "expected"), [("front", -1), ("back", 0)])
def test_z_order_moves_a_shape_through_the_stack(imported: Session, slide,
                                                 position: str, expected: int) -> None:
    shape = slide.shapes[len(slide.shapes) // 2]
    assert router.dispatch(imported, "set_z_order",
                           {"shape_id": shape.id, "position": position})["ok"]
    assert slide.shapes[expected] is shape


def test_moving_a_shape_already_in_place_is_refused(imported: Session, slide) -> None:
    """"Already there" is a fact worth saying; a silent success invites the model to try
    again and call it progress."""
    result = router.dispatch(imported, "set_z_order",
                             {"shape_id": slide.shapes[0].id, "position": "back"})
    assert result["ok"] is False
    assert result["error"] == "already_there"


def test_undo_restores_the_stacking_order(imported: Session, slide) -> None:
    order = [s.id for s in slide.shapes]
    router.dispatch(imported, "set_z_order",
                    {"shape_id": slide.shapes[1].id, "position": "front"})
    assert router.dispatch(imported, "undo")["ok"]
    assert [s.id for s in slide.shapes] == order


# ----------------------------------------------------------------- the package (L2)


def _counts(path: Path) -> tuple[int, int, int]:
    with zipfile.ZipFile(path) as z:
        parts = [n for n in z.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)]
        presentation = etree.fromstring(z.read("ppt/presentation.xml"))
        ids = presentation.findall(f".//{{{P}}}sldIdLst/{{{P}}}sldId")
        rels = etree.fromstring(z.read("ppt/_rels/presentation.xml.rels"))
        slide_rels = [e for e in rels if (e.get("Type") or "").endswith("/slide")]
    return len(parts), len(ids), len(slide_rels)


@pytest.mark.parametrize("operation", ["delete", "add", "duplicate", "reorder"])
def test_shape_operations_leave_the_package_valid(imported: Session, slide,
                                                  tmp_path: Path, operation: str) -> None:
    """An orphaned part makes PowerPoint offer to *repair* the file rather than open it."""
    if operation == "delete":
        router.dispatch(imported, "delete_shape", {"shape_id": _editable(slide).id})
    elif operation == "add":
        router.dispatch(imported, "add_textbox",
                        {"slide_id": slide.id, "region": "body", "text": "x"})
    elif operation == "duplicate":
        router.dispatch(imported, "duplicate_shape", {"shape_id": _editable(slide).id})
    else:
        router.dispatch(imported, "set_z_order",
                        {"shape_id": slide.shapes[1].id, "position": "front"})

    out = tmp_path / f"{operation}.pptx"
    assert router.dispatch(imported, "export", {"path": str(out)})["ok"]
    parts, ids, rels = _counts(out)
    assert parts == ids == rels == len(imported.deck.slides)


def test_deleting_every_editable_shape_still_exports(imported: Session, slide,
                                                     tmp_path: Path) -> None:
    """The degenerate case. An empty slide is a legitimate slide."""
    for shape in [s for s in slide.shapes if not s.opaque and s.ooxml_id]:
        router.dispatch(imported, "delete_shape", {"shape_id": shape.id})
    out = tmp_path / "emptied.pptx"
    assert router.dispatch(imported, "export", {"path": str(out)})["ok"]
    parts, ids, rels = _counts(out)
    assert parts == ids == rels
