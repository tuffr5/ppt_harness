"""Rich text — bridging to OOXML runs.

OOXML has no "bold text": a paragraph is a list of runs, each carrying its own properties.
The harness modelled text as a flat string, so it could only ever write one run — and a
request to bold something produced literal asterisks on the slide.

The rule these tests hold to: **emphasis is content, typography is the theme's.** Which
words are stressed is part of what the author is saying; what face and size they are set in
is not, which is why `bold` is expressible here and `set_font` still does not exist.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from ppt_harness.core.session import Session
from ppt_harness.state import richtext as rt
from ppt_harness.tools import router

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _flags(run: rt.Run) -> str:
    marks = "".join(k[0] for k in ("bold", "italic", "underline", "strike")
                    if getattr(run, k))
    return marks + run.script


# ------------------------------------------------------------------------ parsing


@pytest.mark.parametrize(("source", "expected"), [
    ("Make **this** bold", [("Make ", ""), ("this", "b"), (" bold", "")]),
    ("*italic*", [("italic", "i")]),
    ("__also bold__", [("also bold", "b")]),
    ("~~struck~~", [("struck", "s")]),
    ("<b>html</b>", [("html", "b")]),
    ("<u>under</u>", [("under", "u")]),
    ("H<sub>2</sub>O", [("H", ""), ("2", "sub"), ("O", "")]),
    ("x<sup>2</sup>", [("x", ""), ("2", "super")]),
])
def test_markup_becomes_runs(source: str, expected: list[tuple[str, str]]) -> None:
    assert [(r.text, _flags(r)) for r in rt.parse(source)] == expected


def test_emphasis_nests() -> None:
    runs = rt.parse("**bold with *italic* inside**")
    assert [(r.text, _flags(r)) for r in runs] == [
        ("bold with ", "b"), ("italic", "bi"), (" inside", "b")]


def test_arithmetic_is_not_emphasis() -> None:
    """`2 * 3 = 6` is a sentence, not a formatting instruction. Over-eager parsing would
    silently delete characters from a slide."""
    assert [(r.text, _flags(r)) for r in rt.parse("2 * 3 = 6")] == [("2 * 3 = 6", "")]


def test_plain_text_is_a_single_unmarked_run() -> None:
    runs = rt.parse("nothing special here")
    assert len(runs) == 1
    assert runs[0].plain


def test_adjacent_runs_with_the_same_properties_merge() -> None:
    """Fewer runs is a smaller file, and closer to what someone selecting text expects."""
    runs = rt.parse("**a****b**")
    assert [(r.text, _flags(r)) for r in runs] == [("ab", "b")]


def test_plain_projection_drops_the_markup() -> None:
    assert rt.to_plain(rt.parse("**Bold** and *italic*")) == "Bold and italic"


def test_markup_round_trips() -> None:
    """A model reading a slide must see the emphasis it is being asked to change."""
    for source in ("**b** and *i*", "~~gone~~", "H<sub>2</sub>O", "plain"):
        assert rt.to_plain(rt.parse(rt.to_markup(rt.parse(source)))) == rt.to_plain(
            rt.parse(source))


# ------------------------------------------------------------------------ the tool


def test_set_text_stores_runs_and_plain_words(imported: Session) -> None:
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    result = router.dispatch(imported, "set_text",
                             {"target": f"{slide.id}/{shape.id}", "text": "a **b** c"})
    assert result["ok"], result.get("message")
    assert shape.text == "a b c", "the plain projection is what budgets and search use"
    assert [(r.text, _flags(r)) for r in shape.runs] == [("a ", ""), ("b", "b"), (" c", "")]


def test_the_budget_measures_what_renders_not_what_was_typed(imported: Session) -> None:
    """`<u>under</u>` is five characters on the slide and thirteen in the argument.
    Budgeting the markup rejects text that fits comfortably."""
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    budget = imported.budget_for(f"{slide.id}/{shape.id}")

    bare = "word " * 6
    marked = "".join(f"<u>{w}</u> " for w in ["word"] * 6)
    from ppt_harness.render import budget as budget_mod

    assert budget_mod.check(bare, budget, imported.theme).ok
    assert budget_mod.check(rt.to_plain(rt.parse(marked)), budget, imported.theme).ok


def test_unemphasised_text_stores_no_runs(imported: Session) -> None:
    """An empty `runs` means "nothing to remember" — the plain string already says it."""
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    router.dispatch(imported, "set_text",
                    {"target": f"{slide.id}/{shape.id}", "text": "just words"})
    assert shape.runs == []


def test_undo_restores_the_emphasis_too(imported: Session) -> None:
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    router.dispatch(imported, "set_text",
                    {"target": f"{slide.id}/{shape.id}", "text": "**first**"})
    router.dispatch(imported, "set_text",
                    {"target": f"{slide.id}/{shape.id}", "text": "second"})
    router.dispatch(imported, "undo")
    assert shape.text == "first"
    assert [(r.text, _flags(r)) for r in shape.runs] == [("first", "b")]


# ------------------------------------------------------------------------- export


def test_emphasis_is_written_as_real_ooxml_runs(imported: Session, tmp_path: Path) -> None:
    """The whole point. One run per span, with `b`/`i`/`u`/`strike`/`baseline` set."""
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    router.dispatch(imported, "set_text", {
        "target": f"{slide.id}/{shape.id}",
        "text": "**B** *i* <u>u</u> ~~s~~ H<sub>2</sub>",
    })
    out = tmp_path / "rich.pptx"
    router.dispatch(imported, "export", {"path": str(out)})

    found: dict[str, dict[str, str]] = {}
    for target in Presentation(str(out)).slides[0].shapes:
        if not target.has_text_frame or "B" not in target.text_frame.text:
            continue
        for para in target.text_frame.paragraphs:
            for run in para.runs:
                node = run._r.find(f"{{{A}}}rPr")
                attrs = dict(node.attrib) if node is not None else {}
                found[run.text] = {k: v for k, v in attrs.items()
                                   if k in ("b", "i", "u", "strike", "baseline")}

    assert found.get("B") == {"b": "1"}
    assert found.get("i") == {"i": "1"}
    assert found.get("u") == {"u": "sng"}
    assert found.get("s") == {"strike": "sngStrike"}
    assert found.get("2") == {"baseline": "-25000"}


def test_no_markup_ever_reaches_the_file(imported: Session, tmp_path: Path) -> None:
    """The bug this replaces: `**this**` written through as literal characters."""
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    router.dispatch(imported, "set_text",
                    {"target": f"{slide.id}/{shape.id}", "text": "Make **this** bold"})
    out = tmp_path / "clean.pptx"
    router.dispatch(imported, "export", {"path": str(out)})

    for target in Presentation(str(out)).slides[0].shapes:
        if target.has_text_frame:
            assert "**" not in target.text_frame.text


# ------------------------------------------------------------------------ preview


def test_the_preview_shows_the_emphasis(imported: Session) -> None:
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    router.dispatch(imported, "set_text",
                    {"target": f"{slide.id}/{shape.id}", "text": "**b** *i* <u>u</u>"})
    markup = imported.render_html(slide.id)
    assert "<strong>b</strong>" in markup
    assert "<em>i</em>" in markup
    assert "<u>u</u>" in markup


def test_preview_text_is_still_escaped(imported: Session) -> None:
    """Emphasis is markup the harness understands; everything else is words."""
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    router.dispatch(imported, "set_text",
                    {"target": f"{slide.id}/{shape.id}", "text": "**a** <script>x</script>"})
    markup = imported.render_html(slide.id)
    assert "<script>" not in markup
    assert "&lt;script&gt;" in markup


# ------------------------------------------------------------------------- import


def test_emphasis_already_in_a_file_is_read(imported: Session) -> None:
    """So a model can see what is bold before deciding what to change, and so editing one
    word does not silently flatten the rest of the paragraph."""
    emphasised = [s for slide in imported.deck.slides for s in slide.shapes if s.runs]
    for shape in emphasised:
        assert rt.to_plain(shape.runs).strip() == (shape.text or "").strip()


def test_an_imported_shape_starts_clean_however_it_was_emphasised(imported: Session) -> None:
    """Emphasis counts towards `dirty`, so the file's own emphasis must be remembered as the
    baseline — otherwise every bolded shape in the deck exports on the first save."""
    assert not any(s.dirty for slide in imported.deck.slides for s in slide.shapes)


# ------------------------------------------------------- emphasis as the only change


def _plain_shape(session: Session):
    slide = session.deck.slides[0]
    return slide, next(s for s in slide.shapes if s.text and not s.opaque and not s.runs)


def test_bolding_a_word_is_a_change(imported: Session) -> None:
    """The words are identical afterwards; the slide is not."""
    slide, shape = _plain_shape(imported)
    first = shape.text.split()[0]
    router.dispatch(imported, "set_text", {
        "target": f"{slide.id}/{shape.id}",
        "text": shape.text.replace(first, f"**{first}**", 1),
    })
    assert shape.text == shape.origin_text
    assert shape.dirty


def test_emphasis_alone_reaches_the_file(imported: Session, tmp_path: Path) -> None:
    """Export patches dirty shapes only, so a bold that does not register as a change is a
    bold the preview shows and the file has never heard of."""
    slide, shape = _plain_shape(imported)
    first = shape.text.split()[0]
    router.dispatch(imported, "set_text", {
        "target": f"{slide.id}/{shape.id}",
        "text": shape.text.replace(first, f"**{first}**", 1),
    })
    out = tmp_path / "emphasis-only.pptx"
    router.dispatch(imported, "export", {"path": str(out)})

    bolded = set()
    for target in Presentation(str(out)).slides[0].shapes:
        if not target.has_text_frame:
            continue
        for para in target.text_frame.paragraphs:
            for run in para.runs:
                if run.font.bold:
                    bolded.add(run.text)
    assert first in bolded


def test_undoing_the_emphasis_makes_the_shape_clean_again(imported: Session) -> None:
    """`dirty` is a comparison, not a latch — that is what keeps the writer from hardening
    shapes the user ended up not changing."""
    slide, shape = _plain_shape(imported)
    router.dispatch(imported, "set_text",
                    {"target": f"{slide.id}/{shape.id}", "text": f"**{shape.text}**"})
    assert shape.dirty
    router.dispatch(imported, "undo")
    assert not shape.dirty


# ------------------------------------------------------------------ managed slides


@pytest.fixture
def bulleted(blank: Session) -> Session:
    router.dispatch(blank, "add_slide", {"layout": "stack", "blocks": [
        {"region": "body", "component": "bullets",
         "slots": {"items": ["Revenue up **12%**", "Costs flat"]}},
    ]})
    return blank


def _slot(session: Session) -> str:
    slide = session.deck.slides[-1]
    return f"{slide.id}/{slide.blocks[0].id}/items"


def test_a_managed_slot_takes_emphasis_too(bulleted: Session) -> None:
    """A slide the harness composed has no less right to a stressed word than one it
    imported. The slot holds the markup; the preview holds the emphasis."""
    markup = bulleted.render_html(bulleted.deck.slides[-1].id)
    assert "<strong>12%</strong>" in markup
    assert "**" not in markup


def test_managed_emphasis_is_written_as_runs(bulleted: Session, tmp_path: Path) -> None:
    """The bug this closes: `frame.text = ...` is one run, so the asterisks were the slide."""
    out = tmp_path / "managed.pptx"
    router.dispatch(bulleted, "export", {"path": str(out)})

    seen: dict[str, bool] = {}
    for slide in Presentation(str(out)).slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            assert "**" not in shape.text_frame.text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    seen[run.text] = bool(run.font.bold)
    assert seen.get("12%") is True
    assert seen.get("Revenue up ") is False


def test_managed_markup_is_neither_measured_nor_searched(bulleted: Session) -> None:
    """Asterisks are storage, not content: counted against a budget they reject text that
    fits, and matched by search they hide the slide that says the word."""
    from ppt_harness.render import budget as budget_mod

    b = bulleted.budget_for(_slot(bulleted))
    assert budget_mod.check_value(["Revenue up **12%**"], b, bulleted.theme).used_em == (
        budget_mod.check_value(["Revenue up 12%"], b, bulleted.theme).used_em)

    hits = [h for found in bulleted.search("12%") for h in found["hits"]]
    assert hits and all("**" not in h["snippet"] for h in hits)


def test_a_managed_slot_still_escapes_everything_else(bulleted: Session) -> None:
    """Emphasis is markup the harness understands; the rest is words."""
    router.dispatch(bulleted, "set_text",
                    {"target": _slot(bulleted), "text": "**a** <script>x</script>"})
    markup = bulleted.render_html(bulleted.deck.slides[-1].id)
    assert "<strong>a</strong>" in markup
    assert "<script>" not in markup
    assert "&lt;script&gt;" in markup
