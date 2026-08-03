"""Export round-trip — DESIGN §6.2.

"Export mutates, never rebuilds" is a load-bearing assumption, not a style preference. If
python-pptx ever stops passing unmodeled parts through, every imported deck silently loses
SmartArt, media, animations, or its sensitivity label. These tests fail loudly instead.

Byte-identity is deliberately *not* the assertion. python-pptx reserializes every XML
declaration and rebuilds `.rels` and `[Content_Types].xml` from its own model. Both are
lossless, so the contract is: canonical XML equality on untouched XML parts, byte equality
on binaries, and set equality on relationships and content types.

Point PPT_HARNESS_FIXTURE at a real deck to run these. A deck carrying SmartArt, a native
chart, transitions, or an OLE object is worth more here than a large one.
"""

from __future__ import annotations

import os
import re
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

RELS = "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"


def _fixture() -> Path:
    """The committed demo deck, or whatever `PPT_HARNESS_FIXTURE` names.

    Never whatever `.pptx` is lying in the repo root: that made the suite test a different
    file per machine, and depend on one that is not in the repository at all.
    """
    env = os.environ.get("PPT_HARNESS_FIXTURE")
    if env:
        return Path(env)
    demo = Path(__file__).resolve().parents[1] / "tests" / "fixtures" / "demo.pptx"
    if not demo.exists():
        pytest.skip("tests/fixtures/demo.pptx is missing; "
                    "run `python scripts/make_demo_deck.py`")
    return demo


def _parts(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as z:
        return {n: z.read(n) for n in z.namelist()}


def _c14n(raw: bytes) -> bytes:
    return etree.tostring(etree.fromstring(raw), method="c14n2")


def _is_xml(name: str) -> bool:
    return name.endswith(".xml") or name.endswith(".rels")


@pytest.fixture(scope="module")
def src() -> Path:
    return _fixture()


@pytest.fixture(scope="module")
def saved(src: Path, tmp_path_factory: pytest.TempPathFactory) -> Path:
    out = tmp_path_factory.mktemp("roundtrip") / "noop.pptx"
    Presentation(str(src)).save(str(out))
    return out


def test_no_part_is_dropped_or_added(src: Path, saved: Path) -> None:
    before, after = set(_parts(src)), set(_parts(saved))
    assert not before - after, f"dropped: {sorted(before - after)}"
    assert not after - before, f"added: {sorted(after - before)}"


def test_binary_parts_survive_byte_identical(src: Path, saved: Path) -> None:
    """Media, thumbnails, embedded objects. No excuse for these to change at all."""
    before, after = _parts(src), _parts(saved)
    changed = [n for n in before if not _is_xml(n) and before[n] != after.get(n)]
    assert not changed, f"binary parts rewritten: {changed}"


def test_unmodeled_xml_survives_canonically(src: Path, saved: Path) -> None:
    """Slides, layouts, masters, notes, theme, and any part python-pptx does not model.

    `.rels` and `[Content_Types].xml` are excluded — they are rebuilt by design and get
    their own set-equality tests below.
    """
    before, after = _parts(src), _parts(saved)
    changed = [
        n
        for n in before
        if _is_xml(n) and not n.endswith(".rels") and n != "[Content_Types].xml"
        if _c14n(before[n]) != _c14n(after[n])
    ]
    assert not changed, f"XML rewritten beyond serialization: {changed}"


def test_every_relationship_survives(src: Path, saved: Path) -> None:
    """Rebuilt parts may reorder, but a lost rId is a broken video or a missing image."""
    before, after = _parts(src), _parts(saved)

    def rels(raw: bytes) -> dict[str, tuple[str | None, ...]]:
        root = etree.fromstring(raw)
        return {
            e.get("Id"): (e.get("Type"), e.get("Target"), e.get("TargetMode"))
            for e in root.iter(RELS)
        }

    for name in (n for n in before if n.endswith(".rels")):
        assert rels(before[name]) == rels(after[name]), f"relationships changed in {name}"


def test_content_types_survive_as_a_set(src: Path, saved: Path) -> None:
    before, after = _parts(src), _parts(saved)

    def types(raw: bytes) -> set[tuple[str | None, ...]]:
        root = etree.fromstring(raw)
        return {
            (e.tag, e.get("Extension"), e.get("PartName"), e.get("ContentType")) for e in root
        }

    assert types(before["[Content_Types].xml"]) == types(after["[Content_Types].xml"])


def test_editing_one_run_touches_only_its_slide(src: Path, tmp_path: Path) -> None:
    """The blast radius of a write is the write. Anything wider is a bug in the exporter."""
    prs = Presentation(str(src))
    edited_slide: int | None = None
    for index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        run.text = "ROUNDTRIP-MARKER"
                        edited_slide = index
                        break
                if edited_slide is not None:
                    break
            if edited_slide is not None:
                break
        if edited_slide is not None:
            break
    if edited_slide is None:
        pytest.skip("fixture has no editable text run")

    out = tmp_path / "edited.pptx"
    prs.save(str(out))

    baseline = tmp_path / "noop.pptx"
    Presentation(str(src)).save(str(baseline))

    expected = f"ppt/slides/slide{edited_slide + 1}.xml"
    base, after = _parts(baseline), _parts(out)
    differing = {n for n in base if base[n] != after[n]}
    assert differing == {expected}, f"expected only {expected}, got {sorted(differing)}"


def _slide_consistency(path: Path) -> tuple[int, int, int]:
    """(slide parts, sldId entries, slide relationships) for a package."""
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    parts = _parts(path)
    slides = [n for n in parts if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)]
    presentation = etree.fromstring(parts["ppt/presentation.xml"])
    ids = presentation.findall(f".//{{{P}}}sldIdLst/{{{P}}}sldId")
    rels = etree.fromstring(parts["ppt/_rels/presentation.xml.rels"])
    slide_rels = [e for e in rels if (e.get("Type") or "").endswith("/slide")]
    return len(slides), len(ids), len(slide_rels)


def test_the_fixture_package_is_internally_consistent(src: Path) -> None:
    """A slide part with no `sldId` entry is an orphan, and PowerPoint offers to *repair*
    the file rather than open it. Cheap to check, and invisible until someone double-clicks."""
    parts, ids, rels = _slide_consistency(src)
    assert parts == ids == rels, (
        f"{src.name}: {parts} slide parts, {ids} sldId entries, {rels} relationships"
    )


def test_export_produces_an_internally_consistent_package(src: Path, tmp_path: Path) -> None:
    from ppt_harness.core.session import Session
    from ppt_harness.tools import router

    session = Session.open(src)
    router.dispatch(session, "add_slide", {
        "layout": "stack",
        "blocks": [
            {"region": "header", "component": "slide_title", "slots": {"title": "Added"}},
            {"region": "body", "component": "bullets", "slots": {"items": ["One", "Two"]}},
        ],
    })
    out = tmp_path / "consistent.pptx"
    router.dispatch(session, "export", {"path": str(out)})

    parts, ids, rels = _slide_consistency(out)
    assert parts == ids == rels, f"{parts} parts, {ids} sldId, {rels} rels"
    assert ids == len(session.deck.slides)
