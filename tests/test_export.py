"""Mutating export and the writer assertions — DESIGN §6.2.

Two separate claims are under test.

**Preservation**: patching the original package leaves everything unmodeled untouched —
media, sensitivity labels, change-tracking parts — and confines the blast radius of an edit
to the slide that changed. `test_roundtrip.py` proves that for python-pptx generally; this
proves it for the harness's own exporter.

**Fidelity**: every text box the harness writes carries `noAutofit`, `spcPts`, and explicit
insets. These are asserted against the emitted XML, not the object model, because the XML is
what the recipient opens.
"""

from __future__ import annotations

import io
import itertools
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation

from ppt_harness.components import decoration
from ppt_harness.core.session import Session
from ppt_harness.io import writer_assertions as fidelity
from ppt_harness.io.export_mutate import ExportError, export
from ppt_harness.state.document import Author, Mode, Slide
from ppt_harness.tools import router

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _parts(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as z:
        return {n: z.read(n) for n in z.namelist()}


def _c14n(raw: bytes) -> bytes:
    return etree.tostring(etree.fromstring(raw), method="c14n2")


def _add_managed(session: Session) -> str:
    result = router.dispatch(session, "add_slide", {
        "layout": "stack",
        "blocks": [
            {"region": "header", "component": "slide_title", "slots": {"title": "Findings"}},
            {"region": "body", "component": "bullets",
             "slots": {"items": ["First", "Second", "Third"]}},
        ],
    })
    assert result["ok"], result
    return result["target"]


# ------------------------------------------------------------------ preservation


@pytest.fixture
def exported(imported: Session, tmp_path: Path) -> Path:
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    with imported.transaction(Author.MODEL) as turn:
        imported.store.write(turn, "set_text", f"{slide.id}/{shape.id}",
                             {"text": "Edited"}, Author.MODEL)
    _add_managed(imported)
    out = tmp_path / "exported.pptx"
    export(imported.deck, out)
    return out


def test_binary_parts_are_untouched(exported: Path, fixture_deck: Path) -> None:
    """Media is the expensive thing an export must never rewrite."""
    before, after = _parts(fixture_deck), _parts(exported)
    for name, blob in before.items():
        if name.startswith("ppt/media") or name.endswith((".jpeg", ".png", ".mp4")):
            assert after[name] == blob, f"{name} was rewritten"


def test_unmodelled_parts_survive(exported: Path, fixture_deck: Path) -> None:
    """Change-tracking, revision info, and any sensitivity label the organisation requires."""
    before, after = _parts(fixture_deck), _parts(exported)
    for name in before:
        if name.startswith(("ppt/changesInfos", "docMetadata")) or name in (
            "ppt/revisionInfo.xml", "docProps/custom.xml"
        ):
            assert after.get(name) == before[name], f"{name} lost or rewritten"


def test_export_only_adds_parts_it_needs(exported: Path, fixture_deck: Path) -> None:
    before, after = set(_parts(fixture_deck)), set(_parts(exported))
    assert not before - after, f"dropped {sorted(before - after)}"
    added = after - before
    assert all("slide" in name for name in added), f"unexpected additions: {sorted(added)}"


def test_untouched_slides_are_not_rewritten(imported: Session, tmp_path: Path,
                                            fixture_deck: Path) -> None:
    """Editing slide 1 must not disturb slide 3.

    Canonical, not byte, equality — python-pptx reserializes every XML declaration
    (`"` to `'`, CRLF to LF). That is the contract DESIGN §6.2 states, and asserting bytes
    here would fail on serialization noise while proving nothing extra about content.
    """
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    with imported.transaction(Author.MODEL) as turn:
        imported.store.write(turn, "set_text", f"{slide.id}/{shape.id}",
                             {"text": "Only this one"}, Author.MODEL)
    out = tmp_path / "one.pptx"
    export(imported.deck, out)

    before, after = _parts(fixture_deck), _parts(out)
    others = [n for n in before if n.startswith("ppt/slides/slide") and "slide1." not in n]
    assert others, "fixture has only one slide; this proves nothing"
    for name in others:
        assert _c14n(after[name]) == _c14n(before[name]), \
            f"{name} changed but nothing on it did"


def test_the_edited_slide_is_the_only_one_that_differs(imported: Session, tmp_path: Path,
                                                       fixture_deck: Path) -> None:
    """The converse: the edit must actually reach the XML, or the test above is vacuous."""
    slide = imported.deck.slides[0]
    shape = next(s for s in slide.shapes if s.text and not s.opaque)
    with imported.transaction(Author.MODEL) as turn:
        imported.store.write(turn, "set_text", f"{slide.id}/{shape.id}",
                             {"text": "UNIQUE-MARKER-9f2a"}, Author.MODEL)
    out = tmp_path / "marked.pptx"
    export(imported.deck, out)

    after = _parts(out)
    carrying = [n for n, blob in after.items()
                if n.startswith("ppt/slides/slide") and b"UNIQUE-MARKER-9f2a" in blob]
    assert carrying == ["ppt/slides/slide1.xml"]


def test_the_edit_actually_lands(exported: Path) -> None:
    text = " ".join(
        shape.text_frame.text
        for shape in Presentation(str(exported)).slides[0].shapes
        if shape.has_text_frame
    )
    assert "Edited" in text


def test_exporting_over_the_original_is_refused(imported: Session,
                                                fixture_deck: Path) -> None:
    with pytest.raises(ExportError, match="refusing"):
        export(imported.deck, fixture_deck)


def test_an_imported_slide_that_lost_its_original_fails_loudly(imported: Session,
                                                               tmp_path: Path) -> None:
    """A freeform slide on an *imported* deck is an overlay on real OOXML.

    Narrowed from "a freeform slide always needs an original", which was not true and had
    not been true since `eject_slide` shipped: ejection produces a freeform slide on a deck
    that never had a package, and the old rule made such a deck unexportable for ever.

    The invariant worth keeping is this one. If the deck came from a file, a slide with no
    part in that file is corruption, and rebuilding it from the model would quietly drop
    every shape the harness does not understand — which is the one thing patching the
    original package exists to prevent.
    """
    imported.deck.slides.append(Slide(id="orphan", index=99, mode=Mode.FREEFORM, shapes=[]))
    with pytest.raises(ExportError, match="no original"):
        export(imported.deck, tmp_path / "x.pptx")


# --------------------------------------------------------------------- fidelity


def test_a_written_slide_satisfies_every_assertion(exported: Path) -> None:
    prs = Presentation(str(exported))
    written = [s for s in prs.slides[-1].shapes if s.has_text_frame]
    assert written, "the managed slide produced no text boxes"
    for shape in written:
        assert fidelity.check_frame(shape.text_frame._bodyPr,
                                    shape.text_frame.paragraphs) == []


def _added_slide_xml(exported: Path, source: Path) -> str:
    """The slide part the harness appended.

    Derived, never hardcoded: the index depends on how many slides the fixture already had,
    and asserting against `slide6.xml` silently inspects an imported slide the moment the
    fixture changes size.
    """
    before, after = set(_parts(source)), _parts(exported)
    added = sorted(n for n in set(after) - before
                   if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
    assert len(added) == 1, f"expected exactly one new slide part, got {added}"
    return after[added[0]].decode()


def test_no_percent_line_spacing_is_emitted(exported: Path, fixture_deck: Path) -> None:
    """`spcPct` resolves against font ascent/descent, so it cannot match CSS line-height."""
    slide = _added_slide_xml(exported, fixture_deck)
    assert "spcPct" not in slide
    assert "spcPts" in slide


def test_autofit_is_off_and_stated_positively(exported: Path, fixture_deck: Path) -> None:
    slide = _added_slide_xml(exported, fixture_deck)
    assert "noAutofit" in slide
    assert "normAutofit" not in slide


def test_insets_are_present_not_merely_zero(exported: Path, fixture_deck: Path) -> None:
    """Omitting them does not mean zero — PowerPoint defaults eat 14.4pt of width."""
    root = etree.fromstring(_added_slide_xml(exported, fixture_deck).encode())
    checked = sum(1 for body in root.iter(f"{{{A}}}bodyPr")
                  if all(body.get(k) is not None for k in ("lIns", "tIns", "rIns", "bIns")))
    assert checked > 0


def test_the_post_export_check_catches_a_planted_violation(exported: Path,
                                                           tmp_path: Path) -> None:
    """Prove the check can fail. A gate that never fires is not a gate."""
    prs = Presentation(str(exported))
    shape = next(s for s in prs.slides[-1].shapes if s.has_text_frame)
    body = shape.text_frame._bodyPr
    for key in ("lIns", "tIns", "rIns", "bIns"):
        if body.get(key) is not None:
            del body.attrib[key]
    broken = tmp_path / "broken.pptx"
    prs.save(str(broken))

    violations = fidelity.assert_fidelity(broken)
    assert any(v.rule == "implicit_insets" for v in violations)
    with pytest.raises(fidelity.FidelityError):
        fidelity.raise_for_violations(violations)


def test_the_check_is_scoped_to_shapes_we_wrote(exported: Path) -> None:
    """An imported deck's untouched boxes are the original author's and are not ours to
    hold to this contract; failing on them would make the check unusable."""
    everything = fidelity.assert_fidelity(exported)
    ours = fidelity.assert_fidelity(exported, only_shapes=set())
    assert ours == []
    assert len(everything) > 0, "the fixture's own slides should not satisfy our contract"


def test_a_no_op_export_is_clean(imported: Session, tmp_path: Path) -> None:
    """Nothing written means nothing to check. An empty written-set must not be read as
    "audit everything", or every untouched import fails on the original author's shapes.
    """
    report = export(imported.deck, tmp_path / "noop.pptx")
    assert report.shapes_patched == 0 and report.shapes_added == 0
    assert report.violations == []
    assert report.ok


def test_shape_ids_are_scoped_per_slide(exported: Path) -> None:
    """Shape ids repeat across slides, so a flat id set would pull in untouched shapes."""
    prs = Presentation(str(exported))
    seen: dict[int, int] = {}
    collisions = 0
    for index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            sid = int(shape.shape_id)
            if sid in seen and seen[sid] != index:
                collisions += 1
            seen[sid] = index
    assert collisions > 0, "fixture no longer exercises cross-slide id reuse"


# ------------------------------------------------------------------- from blank


def test_a_deck_with_no_original_exports_at_theme_size(blank: Session,
                                                       tmp_path: Path) -> None:
    _add_managed(blank)
    out = tmp_path / "new.pptx"
    report = export(blank.deck, out)
    assert report.ok

    prs = Presentation(str(out))
    expected_w, expected_h = blank.slide_size_emu()
    assert int(prs.slide_width) == expected_w
    assert int(prs.slide_height) == expected_h
    assert len(prs.slides.__iter__.__self__._sldIdLst) == 1


def test_managed_geometry_lands_inside_the_slide(blank: Session, tmp_path: Path) -> None:
    _add_managed(blank)
    out = tmp_path / "geom.pptx"
    export(blank.deck, out)
    prs = Presentation(str(out))
    for shape in prs.slides[0].shapes:
        assert shape.left >= 0 and shape.top >= 0
        assert shape.left + shape.width <= int(prs.slide_width)
        assert shape.top + shape.height <= int(prs.slide_height)


# ------------------------------------------------- managed objects and multi-column slots


def _managed_chart_slide(session: Session) -> dict:
    return router.dispatch(session, "add_slide", {
        "layout": "stack",
        "blocks": [
            {"region": "header", "component": "slide_title",
             "slots": {"title": "Churn doubled in EMEA"}},
            {"region": "body", "component": "chart", "variant": "wide", "slots": {"chart": {
                "kind": "bar", "categories": ["Q1", "Q2", "Q3"],
                "series": [{"name": "EMEA", "values": [4.1, 6.0, 8.3]}]}}},
        ],
    })


def test_a_managed_chart_reaches_the_file(blank: Session, tmp_path: Path) -> None:
    """The writer read `labels` where every other layer says `categories`.

    A chart slot filled exactly as the schema documents produced an empty `ChartSpec`, was
    dropped, and `add_slide` still answered ok — the slide measured clean and exported a
    title with nothing under it. Nothing disagreed, because the one component that knew was
    silent.
    """
    assert _managed_chart_slide(blank)["ok"]
    out = tmp_path / "chart.pptx"
    report = export(blank.deck, out)

    graphics = [s for s in Presentation(str(out)).slides[0].shapes if s.has_chart]
    assert graphics, "the chart slot exported no chart"
    assert [c for c in graphics[0].chart.plots[0].categories] == ["Q1", "Q2", "Q3"]
    assert not report.violations


def test_a_slot_that_writes_nothing_is_reported(blank: Session, tmp_path: Path) -> None:
    """Silence is the failure mode that let the chart bug live.

    A filled slot the writer cannot build is content the user will not find in the file, so
    it has to reach the export report rather than being skipped.
    """
    assert _managed_chart_slide(blank)["ok"]
    # Past the tool gate, straight into the model: a payload the writer cannot use.
    slide = blank.deck.slides[0]
    block = next(b for b in slide.blocks if b.component == "chart")
    block.slots["chart"] = {"kind": "bar", "categories": ["Q1"], "series": []}

    # `strict=False`, as the `export` tool calls it: the point is that the report
    # *says so*, not that the writer refuses to finish.
    report = export(blank.deck, tmp_path / "dropped.pptx", strict=False)
    assert [v for v in report.violations if v.rule == "slot_not_written"], \
        "a dropped slot exported silently"


def test_a_row_of_stats_is_written_across_not_down(blank: Session, tmp_path: Path) -> None:
    """`per_row` was declared by twelve variants and read by nothing.

    Every list component rendered as one vertical column, so `stat_row` and `bullets` were
    the same slide. The check is on the *file*, because the preview is the file rendered.
    """
    assert router.dispatch(blank, "add_slide", {
        "layout": "stack",
        "blocks": [{"region": "body", "component": "stat_row", "variant": "flat", "slots": {
            "items": [{"value": "8.3%", "label": "EMEA churn"},
                      {"value": "+41%", "label": "Expansion"},
                      {"value": "108%", "label": "Retention"}]}}],
    })["ok"]

    out = tmp_path / "stats.pptx"
    export(blank.deck, out)
    boxes = [s for s in Presentation(str(out)).slides[0].shapes if s.has_text_frame]
    assert len(boxes) == 3, "a row of three stats is three boxes, not one"
    assert len({b.left for b in boxes}) == 3, "the cells share an x; they were stacked"
    assert len({b.top for b in boxes}) == 1, "the cells are a row, so they share a y"

    figures = [b.text_frame.paragraphs[0].runs[0].text for b in boxes]
    assert "8.3%" in figures, "the figure was dropped; only the label survived"


def _place(session: Session, component: str, variant: str, slots: dict) -> None:
    assert router.dispatch(session, "add_slide", {"layout": "stack", "blocks": [
        {"region": "body", "component": component, "variant": variant,
         "slots": slots}]})["ok"]


def _rendering(slide) -> list[tuple]:
    """What the file says about a slide, less the ids that differ per session.

    Shape kind, slot, box, and — for a table, whose geometry the recipient's style decides —
    which rows are painted. Anything a reader could see the difference in.
    """
    out = []
    for shape in slide.shapes:
        row: list = [str(shape.shape_type), shape.name.split(":")[-1],
                     shape.left, shape.top, shape.width, shape.height]
        if getattr(shape, "has_table", False):
            row.append(tuple(str(shape.table.cell(r, 0).fill.type)
                             for r in range(len(shape.table.rows))))
            row.append(shape.table.horz_banding)
        out.append(tuple(row))
    return out


#: A real image, generated rather than committed, because the only thing any of these tests
#: cares about is its proportions — and a binary in the repository is a thing to explain.
ASSET = "a"
PICTURE = (320, 200)


def _png(size: tuple[int, int] = PICTURE) -> bytes:
    from PIL import Image as PILImage

    buffer = io.BytesIO()
    PILImage.new("RGB", size, (0x15, 0x60, 0x82)).save(buffer, format="PNG")
    return buffer.getvalue()


def _with_asset(session: Session, size: tuple[int, int] = PICTURE) -> dict:
    """Put a picture in the deck's assets and hand back the slot payload that names it.

    Assets live on the *store*, not on the deck — they are deliberately outside the document
    model — so this is also what the export tool does when it threads `session.assets`
    through to the writer.
    """
    session.store.assets[ASSET] = ("image/png", _png(size))
    return {"asset_id": ASSET, "alt": "A photograph of the thing"}


STATS = [{"value": "8.3%", "label": "EMEA churn"}, {"value": "+41%", "label": "Expansion"}]
GRID = {"headers": ["Region", "Q1"],
        "rows": [["us-east", "18.2"], ["us-west", "11.9"], ["eu-central", "9.4"]]}
MEDIA = {"asset_id": ASSET, "alt": "A photograph of the thing"}
SIBLINGS = [
    ("stat_row", "flat", "carded", {"items": STATS}),
    ("comparison", "split", "table", {"left": ["Now"], "right": ["Later"]}),
    ("data_table", "plain", "zebra", {"tabular": GRID}),
    ("image_split", "image_left", "image_right",
     {"media": MEDIA, "prose": "What the picture shows"}),
    ("image_full", "bleed", "inset", {"media": MEDIA}),
]


@pytest.mark.parametrize("component,first,second,slots", SIBLINGS,
                         ids=[f"{c}-{a}-{b}" for c, a, b, _ in SIBLINGS])
def test_variants_the_catalog_offers_reach_the_file_as_different_slides(
    blank: Session, tmp_path: Path, component: str, first: str, second: str, slots: dict
) -> None:
    """The catalog was promising renderings the writer never made.

    These pairs expanded to the same boxes and exported the same shapes, so a model that
    asked for `carded` was told it got one and shipped `flat`. On the file, because the file
    is what the recipient opens.

    `image_full` is here because it was the worst of them: `bleed` and `inset` differ only in
    a margin around a picture, and the writer placed no picture at all, so the two variants
    were the same *empty* slide.
    """
    _with_asset(blank)
    _place(blank, component, first, slots)
    _place(blank, component, second, slots)
    out = tmp_path / "siblings.pptx"
    export(blank.deck, out, assets=blank.assets)

    one, two = list(Presentation(str(out)).slides)[:2]
    assert _rendering(one) != _rendering(two)


# ------------------------------------------------------------------------- media


def _picture(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    found = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert found, "the media slot exported no picture"
    return found[0]


def test_a_media_slot_puts_a_real_picture_in_the_file(blank: Session,
                                                     tmp_path: Path) -> None:
    """The writer skipped `media` entirely, and skipping it retired two components.

    `image_full` and `image_split` are the whole of the catalog's imagery, and every slide
    built from either exported as its caption and a blank rectangle. Checked on the file,
    with the alt text, because `descr` is the only part of a picture a screen reader can
    reach and `add_image` has refused a picture without one since it shipped.
    """
    payload = _with_asset(blank)
    _place(blank, "image_full", "bleed", {"media": payload})
    out = tmp_path / "media.pptx"
    report = export(blank.deck, out, assets=blank.assets)

    picture = _picture(Presentation(str(out)).slides[0])
    assert picture._element._nvXxPr.cNvPr.get("descr") == payload["alt"]
    assert not report.violations


def test_a_picture_keeps_its_proportions_and_stays_inside_its_box(blank: Session,
                                                                  tmp_path: Path) -> None:
    """Letterboxed, never cropped and never stretched — see `io/media.py`.

    Two claims, and both matter for the same reason: the harness chose the box but has never
    seen the image. Overflowing it would put a picture over the slide's own margins, and
    filling it by distortion or by cropping would be the writer deciding which part of
    somebody's photograph is the point.
    """
    payload = _with_asset(blank, size=(400, 100))  # far wider than any slot box
    _place(blank, "image_full", "bleed", {"media": payload})
    out = tmp_path / "aspect.pptx"
    export(blank.deck, out, assets=blank.assets)

    slide = Presentation(str(out)).slides[0]
    picture = _picture(slide)
    box = _slot_box(blank, "media")
    assert picture.width / picture.height == pytest.approx(400 / 100, rel=1e-3), \
        "the picture was stretched or cropped to fill its box"
    assert box[0] <= picture.left and box[1] <= picture.top
    assert picture.left + picture.width <= box[0] + box[2] + 1
    assert picture.top + picture.height <= box[1] + box[3] + 1


def test_bleed_and_inset_reach_the_file_as_different_pictures(blank: Session,
                                                              tmp_path: Path) -> None:
    """A margin, which is what the two words mean, and it has to be visible in the file.

    The expander has drawn this distinction for as long as `inset` has existed; nothing
    downstream carried it, so the pair was indistinguishable to the recipient. Strictly
    inside, on every edge — "different" would also be satisfied by a picture that moved.
    """
    payload = _with_asset(blank)
    _place(blank, "image_full", "bleed", {"media": payload})
    _place(blank, "image_full", "inset", {"media": payload})
    out = tmp_path / "bleed-inset.pptx"
    export(blank.deck, out, assets=blank.assets)

    one, two = list(Presentation(str(out)).slides)[:2]
    bleed, inset = _picture(one), _picture(two)
    assert inset.width < bleed.width and inset.height < bleed.height
    assert bleed.left < inset.left and bleed.top < inset.top


def test_media_scale_shrinks_the_picture_without_moving_it(blank: Session,
                                                           tmp_path: Path) -> None:
    """`media_scale` is a clamped override that had no reader anywhere in the harness.

    Concentric, because "fills 60% of its box" describes a smaller picture in the same place,
    not one pushed into a corner.
    """
    payload = _with_asset(blank)
    _place(blank, "image_full", "bleed", {"media": payload})
    block = blank.deck.slides[0].blocks[0]
    assert router.dispatch(blank, "set_override", {
        "slide_id": blank.deck.slides[0].id, "block_id": block.id,
        "key": "media_scale", "value": 0.6})["ok"]

    out = tmp_path / "scaled.pptx"
    export(blank.deck, out, assets=blank.assets)
    small = _picture(Presentation(str(out)).slides[0])

    blank.deck.slides[0].blocks[0].overrides = {}
    export(blank.deck, out, assets=blank.assets)
    full = _picture(Presentation(str(out)).slides[0])

    assert small.width < full.width and small.height < full.height
    assert (small.left + small.width / 2) == pytest.approx(full.left + full.width / 2, abs=2)
    assert (small.top + small.height / 2) == pytest.approx(full.top + full.height / 2, abs=2)


def test_a_media_slot_whose_asset_is_behind_nothing_is_reported(blank: Session,
                                                                tmp_path: Path) -> None:
    """The same claim as the chart bug, for the slot that replaced it as the silent one.

    A picture the deck names and the file does not have is content the user will not find,
    and the export report is where that has to be said. Nothing is invented in its place: a
    placeholder rectangle would be a shape the managed slide never had.
    """
    _place(blank, "image_full", "bleed",
           {"media": {"asset_id": "nothing-is-here.png", "alt": "A photograph"}})
    report = export(blank.deck, tmp_path / "missing.pptx", strict=False,
                    assets=blank.assets)

    dropped = [v for v in report.violations if v.rule == "slot_not_written"]
    assert dropped, "a media slot with no asset exported silently"
    assert "media" in dropped[0].detail


@pytest.mark.parametrize("payload,error", [
    ({"alt": "A photograph"}, "media_needs_asset"),
    ({"asset_id": ASSET}, "alt_required"),
    ({"asset_id": ASSET, "alt": "   "}, "alt_required"),
])
def test_a_media_slot_must_name_a_picture_and_describe_it(blank: Session, payload: dict,
                                                          error: str) -> None:
    """Alt text is required on the managed path too, and refused at the gate.

    `add_image` has refused a picture with no alt text since it shipped; a `media` slot is
    the same picture reached through a component, and a rejected write is the cheapest
    possible failure — the alternative is `ok` now and a missing picture at export.
    """
    result = router.dispatch(blank, "add_slide", {"layout": "stack", "blocks": [
        {"region": "body", "component": "image_full", "variant": "bleed",
         "slots": {"media": payload}}]})
    assert result["ok"] is False
    assert result["error"] == error


def test_the_preview_draws_the_picture_the_file_gets(blank: Session,
                                                     tmp_path: Path) -> None:
    """Preview equals export, which is the invariant everything else here is measured on.

    A media slot the preview left blank is how `bleed` and `inset` managed to look identical
    in *both*: the two disagreed with each other about nothing, because neither drew
    anything. The rectangle is checked, not just the presence of an `<img>` — the preview
    holds the slot box shrunk by `media_scale` and lets `object-fit: contain` letterbox
    inside it, which is the writer's `Box.fit` expressed in CSS.
    """
    payload = _with_asset(blank)
    _place(blank, "image_full", "bleed", {"media": payload})
    slide = blank.deck.slides[0]

    markup = blank.render_html(slide.id)
    assert "<img" in markup and payload["alt"] in markup
    assert "object-fit: contain" in markup

    box = _slot_box(blank, "media")
    export(blank.deck, tmp_path / "agree.pptx", assets=blank.assets)
    picture = _picture(Presentation(str(tmp_path / "agree.pptx")).slides[0])
    # The preview's container is the slot box itself at `media_scale: 1.0`; the file's
    # picture is letterboxed inside it. Same box, same rule.
    assert box[0] <= picture.left and picture.left + picture.width <= box[0] + box[2] + 1


def _slot_box(session: Session, slot: str) -> tuple[int, int, int, int]:
    """The EMU rectangle the expander gave one slot of the first slide's first block."""
    from ppt_harness.render import expand

    slide = session.deck.slides[0]
    laid_out = next(s for s in expand.expand_slide(session.theme, slide) if s.slot == slot)
    cx, cy = session.slide_size_emu()
    return laid_out.box.emu(*session.theme.grid.canvas, cx, cy)


# -------------------------------------------------------------------------- eject


def test_an_ejected_carded_slide_keeps_its_cards(blank: Session, tmp_path: Path) -> None:
    """Eject is one-way, so anything it drops is destroyed rather than mislaid.

    It froze TEXT shapes and only text shapes, so a slide built from a decorated variant left
    managed mode without its panels and there was no way back to them. Compared against the
    managed export of the same slide, because "lossless" is a claim about the file the
    recipient opens, not about the shape list.
    """
    _place(blank, "stat_row", "carded", {"items": STATS})
    slide_id = blank.deck.slides[0].id

    managed_out = tmp_path / "managed.pptx"
    export(blank.deck, managed_out, assets=blank.assets)
    before = [(s.left, s.top, s.width, s.height)
              for s in Presentation(str(managed_out)).slides[0].shapes
              if s.name.endswith(".panel")]
    assert before, "the managed slide had no cards to lose"

    assert router.dispatch(blank, "eject_slide", {"slide_id": slide_id})["ok"]
    assert blank.deck.slides[0].mode is Mode.FREEFORM
    frozen = [s for s in blank.deck.slides[0].shapes if s.geometry is not None]
    assert len(frozen) == len(before), "the cards were dropped on the way out"

    ejected_out = tmp_path / "ejected.pptx"
    export(blank.deck, ejected_out, strict=False, assets=blank.assets)
    slide = Presentation(str(ejected_out)).slides[0]
    names = [s.name for s in slide.shapes]
    panels = [s for s in slide.shapes if s.name.endswith("_panel")]
    assert [(s.left, s.top, s.width, s.height) for s in panels] == before, \
        "the frozen cards are not where the managed slide drew them"

    paint = decoration.paint_for(blank.theme, "card", "list")
    for panel in panels:
        assert str(panel.fill.fore_color.rgb) == paint.fill.lstrip("#").upper()
        assert names.index(panel.name) < names.index(panel.name.removesuffix("_panel")), \
            "a card frozen after its words would cover them"


def test_an_ejected_image_slide_keeps_its_picture(blank: Session, tmp_path: Path) -> None:
    """Same one-way loss, on the slot the slide is actually about.

    The frozen frame carries the picture's own proportions, so the ejected slide is the
    managed one rather than a re-derivation of it — `eject_slide` asks `Box.fit` exactly the
    question the writer does.
    """
    payload = _with_asset(blank)
    _place(blank, "image_full", "inset", {"media": payload})
    slide_id = blank.deck.slides[0].id

    managed_out = tmp_path / "managed-image.pptx"
    export(blank.deck, managed_out, assets=blank.assets)
    before = _picture(Presentation(str(managed_out)).slides[0])
    frame = (before.left, before.top, before.width, before.height)

    assert router.dispatch(blank, "eject_slide", {"slide_id": slide_id})["ok"]
    ejected_out = tmp_path / "ejected-image.pptx"
    export(blank.deck, ejected_out, strict=False, assets=blank.assets)

    after = _picture(Presentation(str(ejected_out)).slides[0])
    assert (after.left, after.top, after.width, after.height) == frame
    assert after._element._nvXxPr.cNvPr.get("descr") == payload["alt"]


def test_a_carded_variant_paints_its_cards_from_the_themes_roles(blank: Session,
                                                                 tmp_path: Path) -> None:
    """A card is a theme decision the writer carries out, not a colour it chooses.

    Checked against the theme's own answer and against the palette, because a decoration
    holding a hex value would be `set_font` wearing a different name. The ordering is part
    of the claim: the panel is drawn before the words, or it covers them.
    """
    _place(blank, "stat_row", "carded", {"items": STATS})
    out = tmp_path / "carded.pptx"
    export(blank.deck, out)

    slide = Presentation(str(out)).slides[0]
    names = [s.name for s in slide.shapes]
    panels = [s for s in slide.shapes if s.name.endswith(".panel")]
    assert len(panels) == len(STATS), "a card behind each figure, or none at all"

    paint = decoration.paint_for(blank.theme, "card", "list")
    palette = {v for v in blank.theme.palette.values() if isinstance(v, str)}
    assert {paint.fill, paint.line} <= palette, "a card was painted from outside the palette"

    for panel in panels:
        assert str(panel.fill.fore_color.rgb) == paint.fill.lstrip("#").upper()
        assert str(panel.line.color.rgb) == paint.line.lstrip("#").upper()
        figure = panel.name.removesuffix(".panel")
        assert names.index(panel.name) < names.index(figure), "the card is over the words"
        text = next(s for s in slide.shapes if s.name == figure)
        assert panel.left <= text.left and panel.top <= text.top
        assert panel.left + panel.width >= text.left + text.width


def test_zebra_bands_alternate_rows_and_plain_states_that_it_does_not(blank: Session,
                                                                     tmp_path: Path) -> None:
    """`bandRow` is inherited from whatever table style the package carries, so which of
    these two the recipient saw was their template's decision rather than the variant's.

    The stripe is written onto the cells as well, for the same reason: a style the harness
    did not author is not one it can rely on to paint anything.
    """
    _place(blank, "data_table", "plain", {"tabular": GRID})
    _place(blank, "data_table", "zebra", {"tabular": GRID})
    out = tmp_path / "zebra.pptx"
    export(blank.deck, out)

    plain, zebra = (next(s.table for s in slide.shapes if getattr(s, "has_table", False))
                    for slide in list(Presentation(str(out)).slides)[:2])
    assert plain.horz_banding is False and zebra.horz_banding is True
    assert all(plain.cell(r, 0).fill.type is None for r in range(len(plain.rows))), \
        "the plain variant painted a row"

    fills = [zebra.cell(r, 0).fill for r in range(len(zebra.rows))]
    banded = [r for r, fill in enumerate(fills) if fill.type is not None]
    assert banded, "zebra banded nothing"
    assert 0 not in banded, "the header was banded, so it reads as a second header"
    assert all(b - a == 2 for a, b in itertools.pairwise(banded)), \
        "the stripes are not alternate rows"
    band = decoration.paint_for(blank.theme, "banded", "tabular")
    assert str(fills[banded[0]].fore_color.rgb) == band.fill.lstrip("#").upper()


def test_a_slot_exports_in_the_face_its_role_asks_for(blank: Session, tmp_path: Path) -> None:
    """The budget measures the display stack, so the file has to render in it.

    Runs carried no face at all, and a bare text box inherits PowerPoint's *minor* font — so
    every `display` role exported in the body face while the preview drew it in the display
    one. On a serif-titled theme the harness was measuring Georgia and shipping Inter, which
    makes "measured against the font that will actually render it" false in the one place it
    is load-bearing.
    """
    from ppt_harness.state import templates

    session = Session.blank("Serif", theme=templates.load("editorial"))
    assert router.dispatch(session, "add_slide", {
        "layout": "stack",
        "blocks": [
            {"region": "header", "component": "slide_title", "slots": {"title": "A title"}},
            {"region": "body", "component": "bullets", "slots": {"items": ["One", "Two"]}},
        ],
    })["ok"]

    out = tmp_path / "faces.pptx"
    export(session.deck, out)
    faces = {shape.name.split(":")[-1]:
             shape.text_frame.paragraphs[0].runs[0].font.name
             for shape in Presentation(str(out)).slides[0].shapes if shape.has_text_frame}

    families = session.deck.theme.type.families
    assert faces["title"] == families["display"].split(",")[0].strip("'\"")
    assert faces["items"] == families["body"].split(",")[0].strip("'\"")
    assert faces["title"] != faces["items"], "the theme distinguishes them; the file must too"


def test_an_ejected_slide_can_still_be_exported(blank: Session, tmp_path: Path) -> None:
    """`eject_slide` used to make a generated deck permanently unexportable.

    Ejection is shipped, documented, and advertised as the way back from managed. On a deck
    the harness generated there is no original package to patch, and the writer raised —
    `ExportError`, past the router, as an exception rather than a refusal. The user's work
    became unrecoverable through the only door that leads out of the harness.

    The geometry ejection produces is absolute, so there was never anything missing; the
    writer simply had no branch for it.
    """
    assert router.dispatch(blank, "add_slide", {
        "layout": "stack",
        "blocks": [
            {"region": "header", "component": "slide_title", "slots": {"title": "A title"}},
            {"region": "body", "component": "bullets", "slots": {"items": ["One", "Two"]}},
        ],
    })["ok"]
    slide_id = blank.deck.slides[0].id
    assert router.dispatch(blank, "eject_slide", {"slide_id": slide_id})["ok"]

    out = tmp_path / "ejected.pptx"
    report = export(blank.deck, out)

    assert report.ok, [str(v) for v in report.violations]
    texts = {s.text_frame.text for s in Presentation(str(out)).slides[0].shapes
             if s.has_text_frame}
    assert "A title" in texts, "the ejected title never reached the file"
    assert any("One" in t for t in texts)


def test_a_duplicated_ejected_slide_still_exports(blank: Session, tmp_path: Path) -> None:
    """`duplicate_slide` + `eject_slide` are both shipped; together they broke export.

    The clone path exists for *imported* slides, where the OOXML holds what the harness
    cannot rebuild. On a generated deck there is no package to clone from, and the writer
    raised — the same hole as the ejected-slide branch, one level in. Found by a benchmark
    run, which it killed at task N of 32.
    """
    assert router.dispatch(blank, "add_slide", {
        "layout": "stack",
        "blocks": [{"region": "header", "component": "slide_title",
                    "slots": {"title": "A title"}}],
    })["ok"]
    slide_id = blank.deck.slides[0].id
    assert router.dispatch(blank, "eject_slide", {"slide_id": slide_id})["ok"]
    assert router.dispatch(blank, "duplicate_slide", {"slide_id": slide_id})["ok"]

    out = tmp_path / "dup.pptx"
    report = export(blank.deck, out)
    assert report.ok, [str(v) for v in report.violations]
    assert len(Presentation(str(out)).slides) == 2


def test_an_export_the_writer_refuses_is_a_result_not_an_exception(imported: Session,
                                                                   tmp_path: Path) -> None:
    """`ExportError` used to escape the router.

    Every other refusal comes back as `{ok: False, error: ...}`; this one propagated as an
    exception, so a caller had to know that one tool throws differently. Under MCP it would
    surface as a transport error rather than a tool result, and it took down a benchmark run.
    """
    imported.deck.slides.append(Slide(id="orphan", index=99, mode=Mode.FREEFORM, shapes=[]))
    result = router.dispatch(imported, "export", {"path": str(tmp_path / "x.pptx")})

    assert result["ok"] is False
    assert result["error"] == "export_failed"
    assert "no original" in result["message"]
