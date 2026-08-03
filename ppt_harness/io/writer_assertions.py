"""Writer assertions.

Four of the seven ways a deck diverges between preview and PowerPoint are closed here, and
they cost nothing. Applied in the writer *and* checked after export, on every export
including imported decks.

The three that must hold on every text box the harness writes:

1. **Autofit off.** `<a:noAutofit/>`, never `<a:normAutofit fontScale=…>`. Overflow must be
   loud; silent shrinking is the exact degradation this architecture exists to prevent.
2. **Line spacing in absolute points.** `spcPts`, never `spcPct` — percent resolves against
   font ascent/descent rather than the em, so it will not match CSS `line-height`, and the
   mismatch varies per font.
3. **Insets explicit.** Omitting them does not mean zero: PowerPoint defaults to 0.1"
   left/right, silently eating 14.4pt of width — enough to push a fitted line into a wrap.

`assert_fidelity` deliberately reads the emitted XML rather than the object model. The
object model is what we believe we wrote; the XML is what the recipient opens.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from ..state.document import EMU_PER_POINT

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


class FidelityError(AssertionError):
    pass


@dataclass(frozen=True)
class Violation:
    slide: int
    shape: str
    rule: str
    detail: str

    def __str__(self) -> str:
        return f"slide {self.slide} · {self.shape}: {self.rule} — {self.detail}"


# ------------------------------------------------------------------------- writing


def harden_text_frame(text_frame, line_pt: float | None = None) -> None:
    """Apply all three assertions to a text frame the harness owns.

    Only call this on frames the harness writes. Running it across an imported deck would
    silently restyle the original author's slides, which is not ours to do.
    """
    # 1. autofit off. Set through the enum rather than by appending XML: `<a:bodyPr>` has a
    #    required child order, and an element inserted in the wrong position is a schema
    #    violation PowerPoint may repair by discarding the whole slide.
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

    # 3. insets explicit — present, not merely defaulted to zero
    text_frame.margin_left = Emu(0)
    text_frame.margin_right = Emu(0)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)
    text_frame.word_wrap = True

    # 2. line spacing in points. `para.line_spacing = 1.4` emits spcPct and is the trap;
    #    only a Length emits spcPts.
    if line_pt is not None:
        for para in text_frame.paragraphs:
            para.line_spacing = Pt(line_pt)


def emu_to_pt(value: int) -> float:
    return value / EMU_PER_POINT


# ------------------------------------------------------------------------ checking


def check_frame(body_pr: etree._Element, paragraphs) -> list[tuple[str, str]]:
    """The three assertions against one text frame's XML. Returns (rule, detail) failures."""
    failures: list[tuple[str, str]] = []

    if body_pr.find(qn("a:normAutofit")) is not None:
        scale = body_pr.find(qn("a:normAutofit")).get("fontScale", "unset")
        failures.append(("autofit_left_on", f"normAutofit present (fontScale={scale})"))
    if body_pr.find(qn("a:noAutofit")) is None:
        failures.append(("autofit_not_asserted", "no <a:noAutofit/>; autofit is inherited"))

    missing = [k for k in ("lIns", "tIns", "rIns", "bIns") if body_pr.get(k) is None]
    if missing:
        failures.append(("implicit_insets", f"missing {', '.join(missing)}; PowerPoint "
                                            "defaults eat 14.4pt of width"))

    for i, para in enumerate(paragraphs):
        p_pr = para._pPr
        if p_pr is None:
            continue
        spc = p_pr.find(qn("a:lnSpc"))
        if spc is not None and spc.find(qn("a:spcPct")) is not None:
            val = spc.find(qn("a:spcPct")).get("val")
            failures.append(("percent_line_spacing",
                             f"paragraph {i} uses spcPct val={val}"))
    return failures


def assert_fidelity(
    path: Path | str, only_shapes: set[tuple[int, int]] | None = None
) -> list[Violation]:
    """Post-export check. Returns violations rather than raising, so a caller can decide.

    `only_shapes` limits the check to shapes the harness wrote, keyed by
    **(slide index, shape id)** — shape ids are unique within a slide, not within a package,
    so a flat set of ids silently pulls in untouched shapes that happen to share a number.

    Scoping to written shapes is the right default on an imported deck: the original
    author's untouched boxes are not ours to hold to this contract, and failing on them
    would make the check unusable on exactly the files it matters most for.
    """
    violations: list[Violation] = []
    prs = Presentation(str(path))
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if only_shapes is not None and (index - 1, int(shape.shape_id)) not in only_shapes:
                continue
            frame = shape.text_frame
            if not frame.text.strip():
                continue
            for rule, detail in check_frame(frame._bodyPr, frame.paragraphs):
                violations.append(Violation(index, shape.name or str(shape.shape_id), rule, detail))
    return violations


def raise_for_violations(violations: list[Violation]) -> None:
    if violations:
        listing = "\n  ".join(str(v) for v in violations)
        raise FidelityError(f"{len(violations)} fidelity violation(s):\n  {listing}")
