"""Import — DESIGN §7 steps 1–3.

Parse the package, extract the theme, and land **every** slide as `freeform`. Nothing is
adopted here; adoption reflows a slide and is always a user-visible proposal (v2).

The importer's discipline is that it must never lose the ability to write the file back.
Shapes the harness models get a role and editable text; everything else — SmartArt, groups,
charts, OLE, embedded video — is recorded as `opaque` with its geometry and left alone. An
opaque shape is still addressable and still exports, it just cannot be edited.
"""

from __future__ import annotations

import contextlib
import hashlib
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape

from ..components import icons as icon_set
from ..state.document import (
    ChartSpec,
    Deck,
    Frame,
    Geometry,
    Mode,
    Shape,
    Slide,
    TypeSpec,
)
from ..state.richtext import Run, to_markup
from ..state.store import DeckStore
from . import colors, media
from .theme_extract import extract_theme

#: Shape types the harness understands well enough to *edit*. Everything else round-trips
#: as opaque rather than being silently flattened or dropped.
EDITABLE = {MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PLACEHOLDER}
ASSET = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.LINKED_PICTURE}

#: Shapes with no text but a drawable outline. Connectors and freeforms carry preset or
#: custom geometry, so hatching them as "unmodelled" hides arrows and rules the slide is
#: partly made of.
DRAWABLE = {MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM}

#: Placeholder type -> theme type-scale role, so `restyle` and budgets have something to
#: consult on a slide that has no component.
ROLE_FROM_PLACEHOLDER = {
    PP_PLACEHOLDER.TITLE: "slide_title",
    PP_PLACEHOLDER.CENTER_TITLE: "deck_title",
    PP_PLACEHOLDER.SUBTITLE: "body",
    PP_PLACEHOLDER.BODY: "body",
    PP_PLACEHOLDER.OBJECT: "body",
}


class ImportError_(RuntimeError):
    pass


def _frame(shape: BaseShape) -> Frame:
    """Geometry in EMU. A placeholder that inherits its box from the layout reports `None`;
    python-pptx resolves that through the layout chain, and 0 is the honest fallback."""
    return Frame(
        x=int(shape.left or 0),
        y=int(shape.top or 0),
        cx=int(shape.width or 0),
        cy=int(shape.height or 0),
    )


def _kind(shape: BaseShape) -> str:
    try:
        st = shape.shape_type
    except (ValueError, AttributeError):
        return "unknown"
    return str(st).split(" ")[0].lower() if st is not None else "unknown"


def _role(shape: BaseShape) -> str | None:
    if not shape.is_placeholder:
        return None
    try:
        return ROLE_FROM_PLACEHOLDER.get(shape.placeholder_format.type)
    except (ValueError, AttributeError):
        return None


#: OOXML sets type in points; the canvas is px at 96dpi.
PT_TO_PX = 96 / 72

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

#: PowerPoint's default body size for a text box that states none, in points.
DEFAULT_TEXTBOX_PT = 18.0

#: `a:path` children an icon can be made of — exactly the four `export_mutate._custom_geometry`
#: emits, and no more. An `a:arcTo` or `a:quadBezTo` in the path means the shape was drawn by
#: something that is not this harness, and refusing to read it is how the match stays honest:
#: the vendored set is normalised to `M`/`L`/`C`/`Z` at vendor time precisely so no other
#: command can ever appear in one of ours.
PATH_COMMAND = {
    f"{{{A_NS}}}moveTo": "M",
    f"{{{A_NS}}}lnTo": "L",
    f"{{{A_NS}}}cubicBezTo": "C",
    f"{{{A_NS}}}close": "Z",
}

def _first_size_pt(txbody) -> tuple[float | None, bool, bool]:
    """The first explicit size in a text body, plus whether that run is bold and italic.

    Searched in the order PowerPoint resolves it: run properties, then the paragraph's
    default run properties, then the body's list style.
    """
    for path in (f".//{{{A_NS}}}r/{{{A_NS}}}rPr",
                 f".//{{{A_NS}}}pPr/{{{A_NS}}}defRPr",
                 f".//{{{A_NS}}}lstStyle//{{{A_NS}}}defRPr"):
        for node in txbody.findall(path):
            if node.get("sz"):
                return int(node.get("sz")) / 100, node.get("b") == "1", node.get("i") == "1"
    return None, False, False


def _inherited_size_pt(shape: BaseShape) -> tuple[float | None, bool, bool]:
    """Walk placeholder -> layout -> master.

    A placeholder with no size of its own is not sizeless; it is the layout's size, and
    then the master's. Falling back to a theme role instead would describe a slide the file
    does not contain.
    """
    if not shape.is_placeholder:
        return None, False, False
    try:
        idx = shape.placeholder_format.idx
        layout = shape.part.slide_layout
    except (ValueError, AttributeError):
        return None, False, False

    try:
        kind = shape.placeholder_format.type
    except (ValueError, AttributeError):
        kind = None

    sources = [layout]
    with contextlib.suppress(ValueError, AttributeError):
        sources.append(layout.slide_master)

    # `idx` first, then type. Date, footer, and slide-number placeholders are typed but
    # rarely share an idx with the master's, and rendering a 12pt footer at body size both
    # looks wrong and invents overflow that is not there.
    candidates = []
    for source in sources:
        candidates += [h for h in source.placeholders
                       if h.placeholder_format.idx == idx]
    if kind is not None:
        for source in sources:
            candidates += [h for h in source.placeholders
                           if h.placeholder_format.type == kind]

    for holder in candidates:
        if holder.has_text_frame:
            size, bold, italic = _first_size_pt(holder.text_frame._txBody)
            if size:
                return size, bold, italic
    return None, False, False


def _type_spec(shape: BaseShape) -> TypeSpec | None:
    """The shape's own type, in canvas px.

    An imported slide is being previewed, not restyled, so the file's sizes win over the
    theme's type scale. `None` means nothing in the chain stated a size and the caller
    should fall back to the theme role.
    """
    if not shape.has_text_frame:
        return None

    size_pt, bold, italic = _first_size_pt(shape.text_frame._txBody)
    if size_pt is None:
        size_pt, bold, italic = _inherited_size_pt(shape)
    if size_pt is None and not shape.is_placeholder:
        size_pt = DEFAULT_TEXTBOX_PT
    if size_pt is None:
        return None

    family = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            family = run.font.name
            break
        if family:
            break

    size_px = size_pt * PT_TO_PX
    line_px = size_px * 1.2
    spacing = shape.text_frame.paragraphs[0].line_spacing if shape.text_frame.paragraphs else None
    if spacing is not None:
        # A Length is absolute points; a bare float is a multiple of the font size.
        line_px = (spacing.pt * PT_TO_PX) if hasattr(spacing, "pt") else size_px * spacing
    return TypeSpec(family=family or "body", size=round(size_px, 2),
                    weight=700 if bold else 400, line=round(line_px, 2), italic=italic)


def _read_runs(frame) -> list[Run]:
    """Emphasis already in the file.

    Read so the model can *see* what is bold before deciding what to change, and so an edit
    to one word does not silently flatten the rest of the paragraph.
    """
    out: list[Run] = []
    for index, para in enumerate(frame.paragraphs):
        if index:
            out.append(Run(text="\n"))
        for run in para.runs:
            node = run._r.find(f"{{{A_NS}}}rPr")
            strike = (node.get("strike") or "") if node is not None else ""
            baseline = int(node.get("baseline") or 0) if node is not None else 0
            out.append(Run(
                text=run.text,
                bold=bool(run.font.bold),
                italic=bool(run.font.italic),
                underline=bool(run.font.underline),
                strike=strike.endswith("Strike"),
                script="super" if baseline > 0 else "sub" if baseline < 0 else "",
            ))
    # Nothing emphasised means nothing to remember; the plain string already says it.
    return [] if all(r.plain for r in out) else out


def _autofit_scale(shape: BaseShape) -> float | None:
    """The `fontScale` PowerPoint applied to make this shape's text fit, if any."""
    if not shape.has_text_frame:
        return None
    node = shape.text_frame._bodyPr.find(f"{{{A_NS}}}normAutofit")
    if node is None:
        return None
    raw = node.get("fontScale")
    return int(raw) / 100000 if raw else 1.0


#: Pictures larger than this are previewed as placeholders. A preview is not a reason to
#: hold 20 MB of image in memory, and a deck of photographs would otherwise be unbounded.
#: Defined in `io/media.py` and re-exported here, because the same number now governs the
#: other end of the pipe: `add_asset` refuses to ingest what import would have skipped, and
#: two copies of a limit is how the two ends come to disagree.
MAX_ASSET_BYTES = media.MAX_ASSET_BYTES


def _asset(shape: BaseShape) -> tuple[str, str, bytes] | None:
    """(key, content type, bytes) for a shape whose picture the preview can show.

    Video shapes expose their poster frame the same way, which is the right thing to draw:
    it is what PowerPoint itself shows before playback.
    """
    try:
        image = shape.image
    except (AttributeError, ValueError):
        return None
    try:
        blob = image.blob
    except (AttributeError, ValueError):
        return None
    if not blob or len(blob) > MAX_ASSET_BYTES:
        return None
    return hashlib.sha1(blob).hexdigest()[:16], image.content_type, blob


def _poster(shape: BaseShape) -> tuple[str, str, bytes] | None:
    """A video's poster frame.

    A movie is stored as a picture with a media relationship hung off it, so the frame
    PowerPoint shows before playback is right there in the package. Drawing it beats a grey
    rectangle labelled "media", and it is what the audience sees on the slide anyway.
    """
    blip = shape._element.find(f".//{{{A_NS}}}blip")
    if blip is None:
        return None
    rid = blip.get(f"{{{R_NS}}}embed")
    if not rid:
        return None
    try:
        part = shape.part.related_part(rid)
        blob = part.blob
        content_type = part.content_type
    except (KeyError, AttributeError, ValueError):
        return None
    if not blob or len(blob) > MAX_ASSET_BYTES or not content_type.startswith("image/"):
        return None
    return hashlib.sha1(blob).hexdigest()[:16], content_type, blob


def _icon_name(sp_pr) -> str:
    """Which vendored icon this shape's custom geometry *is*, or `""` — DESIGN §7 step 1.

    Import lands every slide as `freeform`, so an exported `icon_row` comes back as a handful
    of shapes carrying `a:custGeom`. Left at that, the harness forgets a mark it drew itself:
    `_geometry` finds no `a:prstGeom`, falls back to `preset="rect"`, and the preview draws a
    stroked rectangle where the file holds a chevron — the harness showing a rendering the
    file does not contain, which is the defect `components/icons` exists to have fixed.
    `Geometry.icon` is already the model's word for this and already has both consumers
    (`render/html` draws it from the path table, `export_mutate._add_icon` re-emits it); the
    importer was the only edge missing.

    Read off the shape's own `spPr` and nothing inherited: an icon states its geometry, and a
    custom path resolved through a layout would be a claim about a *different* element.

    Every check below is a way to say no, deliberately — DESIGN §7 is explicit that import
    infers nothing silently, and a wrong name here is worse than none. A miss leaves the shape
    exactly the anonymous freeform it is today; a false positive rewrites what the preview
    draws and what the writer would re-emit if the slide were ever frozen again.
    """
    geometry = sp_pr.find(f"{{{A_NS}}}custGeom")
    if geometry is None:
        return ""
    paths = geometry.findall(f"{{{A_NS}}}pathLst/{{{A_NS}}}path")
    if len(paths) != 1:
        return ""
    path = paths[0]
    # A square space with a stated side, because that is what the writer declares and what
    # `icons.identify` scales the table into. A path that states neither, or states a
    # rectangle, is not one of ours — and guessing the space would be the tolerance this
    # deliberately does not have.
    try:
        units = int(path.get("w", ""))
    except ValueError:
        return ""
    if units <= 0 or path.get("w") != path.get("h"):
        return ""

    commands: list[icon_set.Placed] = []
    for node in path:
        command = PATH_COMMAND.get(node.tag)
        if command is None:
            return ""
        values: list[int] = []
        for point in node.findall(f"{{{A_NS}}}pt"):
            # `int`, not `float`: ST_Coordinate also admits a universal measure like `"1in"`,
            # and a path written in those was not written here. Refuse rather than parse.
            try:
                values.extend((int(point.get("x", "")), int(point.get("y", ""))))
            except ValueError:
                return ""
        if len(values) != icon_set.ARITY[command]:
            return ""
        commands.append((command, tuple(values)))
    return icon_set.identify(icon_set.geometry_key(commands), units)


def _geometry(shape: BaseShape, scheme: dict[str, str]) -> Geometry | None:
    """How the file says to draw this shape.

    Fill and outline inherit. A slide's footer placeholder usually states no fill of its
    own — the coloured band behind it lives on the layout — so reading only the slide's
    `spPr` renders a deck stripped of its own furniture.
    """
    element = shape._element
    sp_pr = element.find(f"{{{P_NS}}}spPr")
    if sp_pr is None:
        return None

    prst = sp_pr.find(f"{{{A_NS}}}prstGeom")
    preset = (prst.get("prst") if prst is not None else None) or "rect"
    fill = colors.fill_of(sp_pr, scheme)
    line = colors.line_of(sp_pr, scheme)

    if fill is None or line is None:
        for holder in _inheritance_chain(shape):
            inherited = holder._element.find(f"{{{P_NS}}}spPr")
            if inherited is None:
                continue
            fill = fill or colors.fill_of(inherited, scheme)
            line = line or colors.line_of(inherited, scheme)
            if prst is None:
                node = inherited.find(f"{{{A_NS}}}prstGeom")
                if node is not None and node.get("prst"):
                    preset = node.get("prst", "rect")
            if fill and line:
                break

    xfrm = sp_pr.find(f"{{{A_NS}}}xfrm")

    geometry = Geometry(
        preset=preset,
        fill=fill[0] if fill else None,
        fill_alpha=fill[1] if fill else 1.0,
        line=line[0] if line else None,
        line_alpha=line[1] if line else 1.0,
        line_width_pt=line[2] if line else 0.75,
        flip_h=bool(xfrm is not None and xfrm.get("flipH") == "1"),
        flip_v=bool(xfrm is not None and xfrm.get("flipV") == "1"),
        rotation=(int(xfrm.get("rot", 0)) / 60000) if xfrm is not None else 0.0,
        icon=_icon_name(sp_pr),
    )
    return geometry if (geometry.visible or geometry.rotation) else None


def _chart(shape: BaseShape) -> ChartSpec | None:
    """Categories and series, enough to draw the chart in the preview.

    Deliberately a *copy* for rendering only. The authoritative data is the embedded
    worksheet in the package, which export never rebuilds — DESIGN §1.5 is explicit that a
    native chart must not degrade into a picture.
    """
    try:
        if not shape.has_chart:
            return None
        chart = shape.chart
    except (AttributeError, ValueError):
        return None

    try:
        kind = str(chart.chart_type).split()[0].lower()
        categories = [str(c) for c in chart.plots[0].categories] if chart.plots else []
        series = []
        for plotted in chart.series:
            values = [None if v is None else float(v) for v in plotted.values]
            series.append({"name": plotted.name or "", "values": values})
    except (AttributeError, ValueError, IndexError):
        return None
    return ChartSpec(kind=kind, categories=categories, series=series)


ALIGN = {"ctr": "center", "r": "right", "just": "justify", "l": "left"}
ANCHOR = {"t": "top", "ctr": "center", "b": "bottom"}


def _alignment(shape: BaseShape) -> str:
    """Paragraph alignment, resolved the way PowerPoint resolves it.

    The slide's own `pPr` wins; failing that the placeholder's list style on the layout,
    then the master. A title that the file centres and the preview left-aligns is the most
    visible way a preview can be wrong.
    """
    if not shape.has_text_frame:
        return "left"
    for para in shape.text_frame.paragraphs:
        p_pr = para._pPr
        if p_pr is not None and p_pr.get("algn"):
            return ALIGN.get(p_pr.get("algn", ""), "left")
    for holder in _inheritance_chain(shape):
        if not holder.has_text_frame:
            continue
        node = holder.text_frame._txBody.find(f".//{{{A_NS}}}lvl1pPr")
        if node is not None and node.get("algn"):
            return ALIGN.get(node.get("algn", ""), "left")
        for para in holder.text_frame.paragraphs:
            if para._pPr is not None and para._pPr.get("algn"):
                return ALIGN.get(para._pPr.get("algn", ""), "left")
    return "left"


def _anchor(shape: BaseShape) -> str:
    """Vertical anchoring. PowerPoint bottom-anchors most title placeholders."""
    if not shape.has_text_frame:
        return "top"
    for holder in [shape, *_inheritance_chain(shape)]:
        if not holder.has_text_frame:
            continue
        value = holder.text_frame._bodyPr.get("anchor")
        if value:
            return ANCHOR.get(value, "top")
    return "top"


def _text_color(shape: BaseShape, scheme: dict[str, str]) -> str | None:
    """The first explicit run colour, or `None` to fall back to the theme's ink."""
    if not shape.has_text_frame:
        return None
    for holder in [shape, *_inheritance_chain(shape)]:
        if not holder.has_text_frame:
            continue
        body = holder.text_frame._txBody
        # Run properties first, then the list style's default run properties — a layout
        # states its footer colour in `lstStyle`, never on a run.
        for path in (f".//{{{A_NS}}}rPr", f".//{{{A_NS}}}defRPr"):
            for node in body.findall(path):
                resolved = colors.resolve(node.find(f"{{{A_NS}}}solidFill"), scheme)
                if resolved:
                    return resolved[0]
    return None


def _inheritance_chain(shape: BaseShape) -> list[BaseShape]:
    """The layout and master placeholders this shape inherits from, nearest first."""
    if not shape.is_placeholder:
        return []
    try:
        idx = shape.placeholder_format.idx
        kind = shape.placeholder_format.type
        layout = shape.part.slide_layout
    except (ValueError, AttributeError):
        return []

    sources = [layout]
    with contextlib.suppress(ValueError, AttributeError):
        sources.append(layout.slide_master)

    exact, by_kind = [], []
    for source in sources:
        for holder in source.placeholders:
            if holder.placeholder_format.idx == idx:
                exact.append(holder)
            elif holder.placeholder_format.type == kind:
                by_kind.append(holder)
    return exact + by_kind


def _is_opaque(shape: BaseShape) -> bool:
    """True for anything the harness cannot faithfully rewrite.

    Groups are opaque because editing a child means recomputing the group's own extents;
    graphic frames cover SmartArt, charts, and OLE. Both are preserved verbatim on export.
    """
    try:
        st = shape.shape_type
    except (ValueError, AttributeError):
        return True
    # Charts stay editable data, so they are drawn rather than hatched — but they are still
    # never *rewritten*; export leaves the native chart part exactly as it found it.
    if st is MSO_SHAPE_TYPE.CHART:
        return False
    if st in (MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
              MSO_SHAPE_TYPE.LINKED_OLE_OBJECT, MSO_SHAPE_TYPE.DIAGRAM,
              MSO_SHAPE_TYPE.TABLE):
        return True
    if st in ASSET or st in EDITABLE or st in DRAWABLE:
        return False
    return not shape.has_text_frame


def _build_shape(shape: BaseShape, prefix: str, scheme: dict[str, str],
                 assets: dict[str, tuple[str, bytes]]) -> Shape:
    """One OOXML shape to the harness model. Shared by slide, layout and master."""
    opaque = _is_opaque(shape)
    text = None
    runs: list[Run] = []
    if not opaque and shape.has_text_frame:
        text = shape.text_frame.text
        runs = _read_runs(shape.text_frame)

    asset_key = None
    if not opaque:
        found = _asset(shape) or _poster(shape)
        if found is not None:
            asset_key, content_type, blob = found
            assets[asset_key] = (content_type, blob)

    return Shape(
        id=f"{prefix}_sh{shape.shape_id}",
        ooxml_id=int(shape.shape_id),
        type=_kind(shape),
        frame=_frame(shape),
        role=_role(shape),
        text=text,
        runs=runs,
        origin_text=text,
        origin_markup=to_markup(runs) if runs else None,
        type_spec=None if opaque else _type_spec(shape),
        autofit_scale=None if opaque else _autofit_scale(shape),
        geometry=None if opaque else _geometry(shape, scheme),
        chart=None if opaque else _chart(shape),
        align=_alignment(shape) if not opaque else "left",
        anchor=_anchor(shape) if not opaque else "top",
        color=_text_color(shape, scheme) if not opaque else None,
        asset=asset_key or (shape.name if _kind(shape) in ("picture", "media") else None),
        opaque=opaque,
    )


def _inherited_shapes(slide, index: int, scheme: dict[str, str],
                      assets: dict[str, tuple[str, bytes]]) -> list[Shape]:
    """Logos, footer bars, and decorative rules the slide inherits.

    Only *non-placeholder* shapes: a layout placeholder is either filled by the slide — in
    which case it is already among the slide's own shapes — or empty, in which case
    PowerPoint draws nothing. Master first, then layout, so layout art sits on top.

    `showMasterSp="0"` on the slide suppresses the master's contribution, and honouring it
    is the difference between a faithful preview and one with a logo the audience will not
    see.
    """
    out: list[Shape] = []
    layout = slide.slide_layout
    show_master = slide._element.get("showMasterSp") != "0"

    sources = []
    if show_master:
        with contextlib.suppress(ValueError, AttributeError):
            sources.append(("m", layout.slide_master))
    sources.append(("l", layout))

    for tag, source in sources:
        for shape in source.shapes:
            if shape.is_placeholder:
                continue
            built = _build_shape(shape, f"s{index + 1}{tag}", scheme, assets)
            if built.text is None and built.asset is None and built.geometry is None:
                continue
            out.append(built)
    return out


def import_pptx(path: Path | str) -> DeckStore:
    path = Path(path)
    if not path.exists():
        raise ImportError_(f"no such file: {path}")

    theme = extract_theme(path)
    prs = Presentation(str(path))

    with zipfile.ZipFile(path) as z:
        scheme = colors.scheme_map(etree.fromstring(z.read("ppt/theme/theme1.xml")))

    assets: dict[str, tuple[str, bytes]] = {}
    slides: list[Slide] = []
    for index, src in enumerate(prs.slides):
        shapes = [_build_shape(shape, f"s{index + 1}", scheme, assets)
                  for shape in src.shapes]
        slides.append(
            Slide(
                id=f"s{index + 1}",
                index=index,
                mode=Mode.FREEFORM,
                hidden=src._element.get("show") == "0",
                notes=src.notes_slide.notes_text_frame.text if src.has_notes_slide else "",
                origin={"part": src.part.partname.rsplit("/", 1)[-1], "index": index},
                shapes=shapes,
                inherited=_inherited_shapes(src, index, scheme, assets),
            )
        )

    deck = Deck(
        id=path.stem,
        title=(prs.core_properties.title or path.stem),
        theme=theme,
        slides=slides,
        source_path=str(path.resolve()),
    )
    store = DeckStore(deck)
    store.assets = assets
    return store
