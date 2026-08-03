"""Font embedding — PLAN B4.

The measured gap on the path the harness controls. Generated slides score 0.007 mean
difference against PowerPoint, and *all* of that residual is font substitution: the harness
lays out with the theme's faces, the recipient's machine does not have them, and PowerPoint
picks something else. Line breaks move, and a slide that fit stops fitting.

Embedding closes it — for everyone whose copy honours the embedded fonts. Two caveats that
belong in the open rather than in a footnote:

- **Licensing decides, not us.** `OS/2 fsType` says whether a face may travel. A restricted
  font is skipped and *reported*, because the honest outcome is a wider fidelity margin,
  not a quietly non-compliant file.
- **PowerPoint for Mac has historically ignored embedded fonts.** Embedding is necessary,
  not sufficient, which is why the measured margin stays.

OOXML embeds a font as its own part, referenced from `presentation.xml`:

    <p:embeddedFontLst>
      <p:embeddedFont>
        <p:font typeface="Inter"/>
        <p:regular r:id="rId9"/>
      </p:embeddedFont>
    </p:embeddedFontLst>
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from ..render import fonts

#: `fsType` bits that forbid travel. 2 = restricted licence; 0x0200 = bitmap embedding only.
#: Everything else — installable, editable, print-and-preview — may be embedded.
NO_EMBED = 0x0002
BITMAP_ONLY = 0x0200

CONTENT_TYPE = "application/x-fontdata"
RELTYPE = ("http://schemas.openxmlformats.org/officeDocument/2006/relationships/font")

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


@dataclass
class EmbedReport:
    embedded: list[str] = field(default_factory=list)
    bytes_added: int = 0
    skipped: dict[str, str] = field(default_factory=dict)
    """family -> why it could not be embedded. Reported, never silent."""

    @property
    def complete(self) -> bool:
        return not self.skipped

    def as_dict(self) -> dict[str, object]:
        out: dict[str, object] = {"embedded": sorted(self.embedded)}
        if self.bytes_added:
            out["kb_added"] = round(self.bytes_added / 1024)
        if self.skipped:
            out["not_embedded"] = self.skipped
            out["note"] = ("these faces will be substituted on a machine that lacks them, "
                           "which moves line breaks; the fidelity margin covers the rest")
        return out


def may_embed(family: str) -> tuple[bool, str]:
    """Whether this face's licence lets it travel, and why not when it does not."""
    path = fonts.find(family)
    if path is None:
        return False, "not installed on this machine"
    try:
        table = fonts.load(path, family)["OS/2"]
    except Exception:
        return False, "no OS/2 table; licensing cannot be determined"

    bits = int(getattr(table, "fsType", 0))
    if bits & NO_EMBED:
        return False, "licence forbids embedding"
    if bits & BITMAP_ONLY:
        return False, "licence permits bitmap embedding only"
    return True, ""


def families_in_use(theme) -> list[str]:
    """Every face the theme can put on a slide, across every script.

    Read from the resolved stacks rather than from the type scale: a CJK deck renders in a
    fallback face the scale never names, and embedding only the Latin one would leave the
    substitution it was meant to prevent.
    """
    seen: list[str] = []
    for stack in theme.type.families.values():
        for family in fonts.parse_stack(stack):
            if family not in seen and fonts.find(family) is not None:
                seen.append(family)
    return seen


def deck_text(deck) -> str:
    """Every character the deck can put on a slide.

    Read from the model rather than from the rendered file so that a subset covers text the
    user is about to type as well: the glyph set is padded with the printable Latin range,
    because a font missing a comma the moment someone adds one is worse than a slightly
    larger file.
    """
    from ..state import slots as slot_render

    chunks: list[str] = [
        "".join(chr(c) for c in range(0x20, 0x7F)),  # room to keep typing
        "\u2018\u2019\u201c\u201d\u2013\u2014\u2026\u2022",  # the punctuation decks use
    ]
    for slide in deck.slides:
        chunks.append(slide.notes or "")
        for shape in [*slide.shapes, *slide.inherited]:
            chunks.append(shape.text or "")
        for block in slide.blocks:
            for value in block.slots.values():
                chunks.append(slot_render.slot_text(value))
    return "".join(chunks)


def _subset(source: Path, family: str, text: str) -> bytes:
    """The face, cut down to the glyphs in use.

    Unsubsetted, two faces turned a 48 KB deck into 1.9 MB — a feature nobody would leave
    on. Subsetting is what makes embedding affordable enough to be the default.
    """
    import io

    from fontTools import subset

    font = fonts.load(source, family)
    options = subset.Options()
    options.layout_features = ["*"]      # keep kerning and ligatures: they change widths
    options.notdef_outline = True
    options.recalc_bounds = True
    options.drop_tables += ["FFTM"]

    subsetter = subset.Subsetter(options=options)
    subsetter.populate(text=text)
    subsetter.subset(font)

    buffer = io.BytesIO()
    font.save(buffer)
    return buffer.getvalue()


def embed(path: Path | str, theme, deck=None) -> EmbedReport:
    """Embed the theme's faces into an exported package.

    Called after the file is written, so it operates on the finished package the same way
    the fidelity audit does — nothing upstream has to know embedding happened.
    """
    from pptx import Presentation

    report = EmbedReport()
    families = families_in_use(theme)
    if not families:
        return report

    text = deck_text(deck) if deck is not None else ""

    prs = Presentation(str(path))
    part = prs.part
    listing = _font_list(part._element)

    for family in families:
        allowed, why = may_embed(family)
        if not allowed:
            report.skipped[family] = why
            continue
        source = fonts.find(family)
        if source is None:  # pragma: no cover - families_in_use already filtered
            report.skipped[family] = "not installed on this machine"
            continue
        try:
            payload = _subset(source, family, text) if text else source.read_bytes()
            rel_id = _attach(part, payload)
        except Exception as exc:  # a font that will not attach is a skip, not a crash
            report.skipped[family] = f"could not be attached: {exc}"
            continue
        _declare(listing, family, rel_id)
        report.embedded.append(family)
        report.bytes_added += len(payload)

    prs.save(str(path))
    return report


def _font_list(presentation_element):
    """`<p:embeddedFontLst>`, created if absent and placed where the schema wants it.

    OOXML is order-sensitive: the element has to sit after `sldSz`/`notesSz`, and a
    presentation with it in the wrong place is one PowerPoint offers to repair.
    """
    from lxml import etree

    existing = presentation_element.find(f"{{{P_NS}}}embeddedFontLst")
    if existing is not None:
        return existing

    listing = etree.SubElement(presentation_element, f"{{{P_NS}}}embeddedFontLst")
    anchor = presentation_element.find(f"{{{P_NS}}}notesSz")
    if anchor is None:
        anchor = presentation_element.find(f"{{{P_NS}}}sldSz")
    if anchor is not None:
        anchor.addnext(listing)
    return listing


def _attach(part, payload: bytes) -> str:
    """Add the font as a package part and return the relationship id."""
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    package = part.package
    index = 1
    while any(str(p.partname) == f"/ppt/fonts/font{index}.fntdata" for p in package.iter_parts()):
        index += 1

    partname = PackURI(f"/ppt/fonts/font{index}.fntdata")
    font_part = Part(partname, CONTENT_TYPE, package, payload)
    return part.relate_to(font_part, RELTYPE)


def _declare(listing, family: str, rel_id: str) -> None:
    """Name the face in `presentation.xml` so PowerPoint knows to use the embedded bytes."""
    from lxml import etree

    for entry in listing.findall(f"{{{P_NS}}}embeddedFont"):
        font = entry.find(f"{{{P_NS}}}font")
        if font is not None and font.get("typeface") == family:
            return

    entry = etree.SubElement(listing, f"{{{P_NS}}}embeddedFont")
    etree.SubElement(entry, f"{{{P_NS}}}font").set("typeface", family)
    # `regular` only: the harness never asks for a bold or italic *face*, it asks for a
    # bold run, and PowerPoint synthesises that from the regular when it must.
    etree.SubElement(entry, f"{{{P_NS}}}regular").set(f"{{{R_NS}}}id", rel_id)
