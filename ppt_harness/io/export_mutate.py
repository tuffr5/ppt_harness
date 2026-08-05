"""Mutating export — DESIGN §6.2.

Regenerating an imported file destroys everything unmodeled: SmartArt, animations,
transitions, embedded media, comments, custom XML, master variants, and the sensitivity
label an organisation may require. So the exporter **opens the original package and patches
only what changed**. State is an overlay on the real file, not a replacement.

The round-trip spike (DESIGN §6.2, `tests/test_roundtrip.py`) is what licenses this: on a
174 MB deck, open-then-save preserved every part, byte-identically for binaries and
canonically for XML, with the blast radius of a single text edit confined to that slide.

A deck created from scratch has no original package, so it starts from a blank presentation
sized to the theme and is written the same way afterwards. One code path, two entry points.
"""

from __future__ import annotations

import shutil
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation
from pptx.util import Emu, Pt

from ..components import decoration as decoration_mod
from ..components import icons as icon_set
from ..components import overrides as overrides_mod
from ..render import expand
from ..state import richtext, slots
from ..state.document import (
    EMU_PER_POINT,
    ChartSpec,
    Deck,
    Frame,
    Mode,
    Shape,
    Slide,
    TableSpec,
)
from . import embed_fonts
from . import media as media_mod
from . import writer_assertions as fidelity

#: The deck's pictures, keyed the way `media` slots and imported shapes name them:
#: `{asset_id: (content_type, bytes)}`. Threaded in from the session rather than read off the
#: deck because assets are deliberately not part of the document model — see `DeckStore`.
Assets = dict[str, tuple[str, bytes]]

ExportMode = Literal["editable", "pixel_locked"]


class ExportError(RuntimeError):
    pass


@dataclass
class ExportReport:
    path: Path
    mode: ExportMode
    slides_written: int = 0
    shapes_patched: int = 0
    shapes_added: int = 0
    violations: list[fidelity.Violation] = field(default_factory=list)
    fonts: dict[str, Any] = field(default_factory=dict)
    """What was embedded and what could not be, from `io/embed_fonts`."""

    @property
    def ok(self) -> bool:
        return not self.violations

    def summary(self) -> str:
        head = (f"{self.path.name} · {self.mode} · {self.slides_written} slides · "
                f"{self.shapes_patched} patched · {self.shapes_added} added")
        if self.violations:
            return head + f" · {len(self.violations)} FIDELITY VIOLATIONS"
        return head + " · fidelity clean"


# ---------------------------------------------------------------------- patching


def _write_visibility(slide, model: Slide) -> None:
    """`<p:sld show="0">` for a hidden slide.

    Absent means shown, so the attribute is removed rather than set to "1" — writing a
    default is churn the round-trip guarantee then has to account for.
    """
    if model.hidden:
        slide._element.set("show", "0")
    else:
        slide._element.attrib.pop("show", None)


def _write_notes(slide, model: Slide) -> None:
    """Speaker notes, which live in their own part.

    Written only when they differ, because touching `notes_slide` *creates* the part — an
    empty notes slide on every exported deck is churn the round-trip guarantee then has to
    explain away.
    """
    current = (slide.notes_slide.notes_text_frame.text
               if slide.has_notes_slide else "")
    if (model.notes or "") == current:
        return
    slide.notes_slide.notes_text_frame.text = model.notes or ""


def _patch_freeform(slide, index: int, model: Slide, written: set[tuple[int, int]],
                    assets: Assets | None = None) -> int:
    """Write back only shapes marked `dirty`.

    Everything else — including every opaque shape — is left exactly as the original author
    wrote it. This is the whole reason imported decks survive a round trip.
    """
    by_id = {int(s.shape_id): s for s in slide.shapes}
    patched = 0
    _write_notes(slide, model)
    _write_visibility(slide, model)

    # Deletions first: a shape removed from the model is removed from the file, whether or
    # not the harness understood it. The element is detached rather than emptied, so no
    # orphan is left behind for PowerPoint to complain about.
    for ooxml_id in model.removed:
        target = by_id.pop(ooxml_id, None)
        if target is not None:
            element = target._element
            element.getparent().remove(element)
            patched += 1

    # Then shapes the harness added, which are not in the original package at all.
    for shape_model in model.shapes:
        if shape_model.ooxml_id or shape_model.opaque:
            continue
        added = _add_shape(slide, shape_model, model, index, assets)
        if added is not None:
            shape_model.ooxml_id = added
            written.add((index, added))
            patched += 1

    for shape_model in model.shapes:
        if not shape_model.dirty or shape_model.opaque:
            continue
        target = by_id.get(shape_model.ooxml_id)
        if target is None:
            continue
        if not (target.has_text_frame or shape_model.table or shape_model.chart):
            continue
        if shape_model.table is not None and target.has_table:
            _patch_table(target.table, shape_model.table)
            written.add((index, int(target.shape_id)))
            patched += 1
            continue
        if shape_model.chart is not None and getattr(target, "has_chart", False):
            _patch_chart(target.chart, shape_model.chart)
            written.add((index, int(target.shape_id)))
            patched += 1
            continue

        frame = target.text_frame
        _write_runs(frame, shape_model.runs, shape_model.text or "")
        _write_paragraph_props(frame, shape_model)
        fidelity.harden_text_frame(frame)
        written.add((index, int(target.shape_id)))
        patched += 1
    return patched


def _write_ejected(prs, model: Slide, deck: Deck, written: set[tuple[int, int]],
                   assets: Assets | None = None) -> int:
    """Write a freeform slide that has no original — the output of `eject_slide`.

    The geometry is already absolute: ejection is the expander's answer written down, so this
    places each shape where the model says and does not re-derive anything.

    Three kinds of shape, in document order, because that *is* the z-order: a decoration
    panel, a picture, or text. Panels and pictures used to be skipped here, which made
    ejecting a `carded` or `image_full` slide a silent demolition — the cards and the
    photograph were simply not in the file, and eject is one-way, so there was nothing to go
    back to. Everything the expander drew has to survive the freeze or the door out of
    managed mode costs more than the catalog it escapes.
    """
    from pptx.util import Emu as _Emu

    theme = deck.theme
    blank = _blank_layout(prs)
    slide = prs.slides.add_slide(blank)
    slide_index = len(prs.slides._sldIdLst) - 1
    added = 0
    emu_per_px = int(prs.slide_width) / max(1, theme.grid.canvas[0])

    for shape in model.shapes:
        if shape.opaque:
            continue

        if shape.geometry is not None and shape.geometry.icon:
            # A frozen icon, written by the same function the managed path uses. The colour
            # and the weight come off the shape rather than being resolved again: eject is
            # one-way, and a mark that had to look its accent up would change colour the
            # moment the block it came from stopped existing.
            f = shape.frame
            _add_icon(slide, shape.geometry.icon, (f.x, f.y, f.cx, f.cy),
                      shape.geometry.line or theme.palette.get("ink", "#111111"),
                      round(shape.geometry.line_width_pt * EMU_PER_POINT),
                      f"{model.id}:{shape.id}")
            added += 1
            continue

        if shape.geometry is not None and shape.geometry.visible:
            f = shape.frame
            geometry = shape.geometry
            # The ramp is read off the frozen shape rather than resolved again from the
            # decoration it came from: eject is one-way, and a panel that had to look its
            # gradient up would lose its depth the moment that decoration changed.
            paint = decoration_mod.Paint(
                fill=geometry.fill or "",
                line=geometry.line or "",
                radius=float(theme.shape.get("radius") or 0.0),
                stops=tuple(decoration_mod.Stopped(at=s.at, colour=s.colour, alpha=s.alpha)
                            for s in geometry.stops),
                kind=geometry.gradient,
                angle=geometry.gradient_angle,
                preset=geometry.preset,
            )
            _paint_panel(slide, (f.x, f.y, f.cx, f.cy), paint,
                         f"{model.id}:{shape.id}", emu_per_px)
            added += 1
            continue

        if shape.asset or shape.source:
            # The frame already has the picture's proportions — `eject_slide` asked
            # `Box.fit` the same question the writer does — so this stretches into it and
            # the result is the rectangle the managed slide had.
            asset = media_mod.resolve(shape.asset or shape.source or "", assets)
            if asset is not None and _add_picture(slide, shape, image=asset.image()):
                added += 1
            continue

        if not (shape.text or "").strip():
            continue
        frame = shape.frame
        box = slide.shapes.add_textbox(_Emu(frame.x), _Emu(frame.y),
                                       _Emu(frame.cx), _Emu(frame.cy))
        box.name = f"{model.id}:{shape.id}"
        text_frame = box.text_frame
        _write_runs(text_frame, shape.runs, shape.text or "")
        _write_paragraph_props(text_frame, shape)
        # The shape's own type, not the theme's role: ejection captured what the slide
        # actually rendered at, and re-resolving it here would move the text.
        line_pt = shape.type_spec.line * 0.75 if shape.type_spec else None
        fidelity.harden_text_frame(text_frame, line_pt=line_pt)
        written.add((slide_index, int(box.shape_id)))
        added += 1

    _clear_empty_placeholders(slide)
    return added


def _write_managed(prs, model: Slide, deck: Deck, written: set[tuple[int, int]],
                   dropped: list[fidelity.Violation] | None = None,
                   assets: Assets | None = None) -> int:
    """Render a managed slide into real placeholders and text boxes.

    Geometry comes from the expander and nowhere else, and it is written as absolute boxes
    so PowerPoint never re-derives layout. That is what makes preview match export.
    """
    theme = deck.theme
    blank = _blank_layout(prs)
    slide = prs.slides.add_slide(blank)
    slide_index = len(prs.slides._sldIdLst) - 1

    canvas_w, canvas_h = theme.grid.canvas
    cx, cy = int(prs.slide_width), int(prs.slide_height)

    added = 0
    laid_out_slots = expand.expand_slide(theme, model)
    # Before anything else, because it is behind everything else: one shape, spanning the
    # slide, and the cheapest thing in the catalog that stops a slide looking like a blank
    # page with words on it.
    added += _paint_slide_layers(slide, model, theme, laid_out_slots, (cx, cy))
    for laid_out in laid_out_slots:
        block = model.block(laid_out.block_id)
        if block is None:
            continue
        value = block.slots.get(laid_out.slot)
        if not value:
            continue

        x, y, w, h = laid_out.box.emu(canvas_w, canvas_h, cx, cy)
        name = f"{model.id}:{laid_out.block_id}:{laid_out.slot}"

        # A `tabular`, `chart` or `media` slot is written as the thing it is — the same
        # structures `add_table`, `add_chart` and `add_image` produce — so a deck the harness
        # generates is one it can read back and keep editing. Writing them as text put a
        # Python dict on a slide.
        if not _is_text_slot(laid_out):
            shape_id = _write_object(slide, laid_out, value, name, (x, y, w, h), theme,
                                     assets, (cx, cy))
            if shape_id is not None:
                added += 1
            elif dropped is not None:
                # A filled slot that produced no object is content the user asked for and
                # will not find in the file. It has to be *said*: the chart-key bug survived
                # because `add_slide` returned ok, the slide measured clean, and the export
                # reported success while writing a title and nothing else.
                dropped.append(fidelity.Violation(
                    slide=slide_index, shape=name, rule="slot_not_written",
                    detail=f"{laid_out.component}.{laid_out.slot} is a "
                           f"{laid_out.shape} slot the writer could not build from "
                           f"{sorted(value) if isinstance(value, dict) else type(value).__name__}",
                ))
            continue

        # A list laid out across columns is written as one text box per cell, because that
        # is what a row of statistics *is* — four boxes side by side, not one box holding
        # four lines. Written as a single frame it exported as a vertical list, which then
        # came back through the preview (the preview is the export, rendered) and made every
        # `stat_row`, `card_grid`, `icon_row` and horizontal `timeline` in the catalog look
        # identical to `bullets`.
        written_cells = expand.written_cells(laid_out, value)
        # One panel per written box, and drawn before it: `spTree` is painted in document
        # order, so a card added after its text would cover the words it is behind.
        panels = laid_out.panels(len(written_cells))
        layers = decoration_mod.layers_for(theme, laid_out.decoration, laid_out.shape)
        for index, ((part, cell), panel) in enumerate(zip(written_cells, panels,
                                                          strict=True)):
            runs = richtext.parse(part)
            text = richtext.to_plain(runs)
            if not text:
                continue
            cell_name = name if laid_out.columns == 1 else f"{name}#{index}"
            # Back to front, in the order the decoration lists them: the ground it stands
            # on, the shade it drops, then the panel — and all of them before the words.
            for layer in layers:
                if layer.place == "slide":
                    continue  # drawn once for the slide, above
                box = expand.decoration_box(theme, panel, layer.place)
                _paint_panel(slide, box.emu(canvas_w, canvas_h, cx, cy), layer,
                             f"{cell_name}.{layer.place}", cx / canvas_w)
            cx_, cy_, cw, ch = cell.emu(canvas_w, canvas_h, cx, cy)
            box = slide.shapes.add_textbox(Emu(cx_), Emu(cy_), Emu(cw), Emu(ch))
            box.name = cell_name

            frame = box.text_frame
            _write_runs(frame, runs, text)
            _style(frame, laid_out, theme, stat=_is_stat_slot(value))
            fidelity.harden_text_frame(frame, line_pt=laid_out.spec.line * 0.75)
            written.add((slide_index, int(box.shape_id)))
            added += 1

        # After the cells, not before: the mark's square is disjoint from its own label by
        # construction — `LaidOutSlot._carve` cut one out of the other — but a decorated icon
        # variant would paint a panel across the whole cell, and a shape added before that
        # panel is a shape behind it. Drawn last, an icon is on its card rather than under it.
        #
        # Colour through the block's accent, which is a *role* on the theme's ordered ramp
        # and never a value in the catalog. That is the whole argument for geometry over a
        # picture: the same path is on-brand on every theme the deck is retargeted to,
        # where a raster is one colour for ever.
        accent = overrides_mod.accent_for(theme, laid_out.overrides)
        for index, (icon_name, icon_box) in enumerate(
                expand.written_icons(laid_out, value)):
            stroke = expand.icon_stroke_px(icon_box) * cx / canvas_w
            if _add_icon(slide, icon_name, icon_box.emu(canvas_w, canvas_h, cx, cy),
                         accent, round(stroke), f"{name}#{index}.icon"):
                added += 1

    _clear_empty_placeholders(slide)
    return added


def _paint_slide_layers(slide, model: Slide, theme, laid_out_slots, size) -> int:
    """The wash behind a whole slide, drawn once however many slots asked for it.

    Once, and first. A slide-scoped layer is not a property of the slot that named it — two
    slots wearing the same decoration would otherwise paint two washes, and two 13% ellipses
    are one 24% one, which is a different slide from the one the preview drew.
    """
    cx, cy = size
    canvas_w, canvas_h = theme.grid.canvas
    worn = [(item.decoration, item.shape) for item in laid_out_slots if item.decoration]
    added = 0
    for layer in decoration_mod.slide_layers(theme, worn):
        box = expand.decoration_box(theme, expand.content_box(theme), layer.place)
        _paint_panel(slide, box.emu(canvas_w, canvas_h, cx, cy), layer,
                     f"{model.id}.{layer.place}", cx / canvas_w)
        added += 1
    return added


def _patch_table(table, spec) -> None:
    """Rewrite the cells of an existing table, leaving its styling alone."""
    grid = ([spec.headers] if spec.headers else []) + spec.rows
    for r, line in enumerate(grid):
        if r >= len(table.rows):
            break
        for c in range(min(spec.columns, len(table.columns))):
            table.cell(r, c).text = line[c] if c < len(line) else ""


def _patch_chart(chart, spec) -> None:
    """Replace a chart's numbers **and** its worksheet.

    Rewriting the plot alone leaves the embedded workbook stale, so the recipient who opens
    the data sees numbers that no longer match the picture.
    """
    from pptx.chart.data import CategoryChartData

    data = CategoryChartData()
    data.categories = spec.categories
    for entry in spec.series:
        data.add_series(entry.get("name") or "", tuple(entry.get("values") or ()))
    chart.replace_data(data)


ALIGN_TO_OOXML = {"left": "l", "center": "ctr", "right": "r", "justify": "just"}

#: What PowerPoint itself writes for a default bulleted list.
BULLET_FONT = "Arial"
BULLET_CHAR = "\u2022"


def _write_paragraph_props(frame, shape_model) -> None:
    """Alignment and list style, which live on `<a:pPr>` rather than on runs."""
    from pptx.oxml.ns import qn

    for para in frame.paragraphs:
        p_pr = para._p.get_or_add_pPr()

        if shape_model.align and shape_model.align != "left":
            p_pr.set("algn", ALIGN_TO_OOXML[shape_model.align])
        else:
            p_pr.attrib.pop("algn", None)

        if shape_model.indent:
            p_pr.set("lvl", str(shape_model.indent))
        else:
            p_pr.attrib.pop("lvl", None)

        # A paragraph carries at most one bullet declaration, so clear before setting.
        for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
            for node in p_pr.findall(qn(tag)):
                p_pr.remove(node)

        if shape_model.bullet == "bullet":
            font = p_pr.makeelement(qn("a:buFont"), {"typeface": BULLET_FONT})
            char = p_pr.makeelement(qn("a:buChar"), {"char": BULLET_CHAR})
            p_pr.append(font)
            p_pr.append(char)
        elif shape_model.bullet == "number":
            p_pr.append(p_pr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"}))
        else:
            # Explicit `buNone`: omitting it inherits the layout's bullet, which is how a
            # paragraph the user cleared comes back bulleted.
            p_pr.append(p_pr.makeelement(qn("a:buNone"), {}))


def _write_link(written, url: str) -> None:
    """`<a:hlinkClick>` needs a relationship, not just an attribute."""
    if not url:
        return
    written.hyperlink.address = url


#: Harness chart kinds to the XL constants. Kept small on purpose: every kind here is one
#: the importer can read back, so a chart the harness writes stays editable by the harness.
CHART_TYPES = {
    "bar": "BAR_CLUSTERED",
    "column": "COLUMN_CLUSTERED",
    "line": "LINE",
    "pie": "PIE",
}


def _add_shape(slide, shape_model, model, index: int,
               assets: Assets | None = None) -> int | None:
    """Write a shape the harness created. Dispatches on what it is.

    A picture comes out of the deck's assets, and falls back to its `source` path only when
    the key resolves to nothing. The order matters: `source` is provenance — where the file
    was read from on the machine that read it — and writing from it means a deck whose
    picture exists on one laptop and nowhere else. `add_image` ingests, so the key is the
    real answer and the path is a leftover; the fallback is kept for a shape recorded before
    it did, and for that shape only.
    """
    if shape_model.table is not None:
        return _add_table(slide, shape_model)
    if shape_model.chart is not None:
        return _add_chart(slide, shape_model)
    if shape_model.asset or shape_model.source:
        asset = media_mod.resolve(shape_model.asset or "", assets)
        if asset is not None:
            return _add_picture(slide, shape_model, image=asset.image())
        return _add_picture(slide, shape_model) if shape_model.source else None
    return _add_textbox(slide, shape_model, model, index)


def _add_picture(slide, shape_model, image: Any = None) -> int | None:
    """A picture, with its alt text.

    `descr` is where PowerPoint keeps alternative text, and it is the only part of a picture
    that a screen reader can use.

    `image` is anything python-pptx opens — a path or a stream — and defaults to the shape's
    `source`, which is what `add_image` recorded. A managed `media` slot passes a stream
    instead, because its picture lives in the deck's assets and has no path on this machine.

    The frame is written verbatim: the picture is stretched into exactly the rectangle it is
    given. Every caller therefore owns the aspect ratio, and both of the harness's own
    callers hand over a rectangle `Box.fit` already gave the image's own proportions.
    """
    from pptx.util import Emu

    f = shape_model.frame
    picture = slide.shapes.add_picture(image if image is not None else shape_model.source,
                                       Emu(f.x), Emu(f.y), Emu(f.cx), Emu(f.cy))
    if shape_model.alt:
        picture._element._nvXxPr.cNvPr.set("descr", shape_model.alt)
    return int(picture.shape_id)


def _add_table(slide, shape_model, band: str | None = None) -> int | None:
    """A real `<a:tbl>`, so the recipient can edit cells rather than a picture of them.

    `band` is the colour alternate rows are filled with, or None for a table that states it
    has no banding. Stated either way: `bandRow` is inherited from whatever table style the
    package carries, so leaving it unsaid made `plain` and `zebra` the same table wearing the
    recipient's template. The fill is written onto the cells rather than left to the style
    for the same reason — a style the harness did not author is not a decision it can rely on.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    spec = shape_model.table
    grid = ([spec.headers] if spec.headers else []) + spec.rows
    f = shape_model.frame
    graphic = slide.shapes.add_table(len(grid), spec.columns,
                                     Emu(f.x), Emu(f.y), Emu(f.cx), Emu(f.cy))
    table = graphic.table
    table.first_row = bool(spec.headers)
    table.horz_banding = band is not None

    head = 1 if spec.headers else 0
    for r, line in enumerate(grid):
        for c in range(spec.columns):
            cell = table.cell(r, c)
            cell.text = line[c] if c < len(line) else ""
            # Every other *body* row: banding that counted the header would put the first
            # band immediately under it, and read as a second header rather than a stripe.
            if band and r >= head and (r - head) % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(band.lstrip("#"))
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
    return int(graphic.shape_id)


def _add_chart(slide, shape_model) -> int | None:
    """A native chart with its embedded worksheet.

    The worksheet is the point. A chart without it is a picture, and DESIGN §1.5 is explicit
    that the harness never degrades one into the other — there is no way back.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu

    spec = shape_model.chart
    data = CategoryChartData()
    data.categories = spec.categories
    for entry in spec.series:
        data.add_series(entry.get("name") or "", tuple(entry.get("values") or ()))

    kind = getattr(XL_CHART_TYPE, CHART_TYPES.get(spec.kind, "COLUMN_CLUSTERED"))
    f = shape_model.frame
    frame = slide.shapes.add_chart(kind, Emu(f.x), Emu(f.y), Emu(f.cx), Emu(f.cy), data)
    return int(frame.shape_id)


def _add_textbox(slide, shape_model, model, index: int) -> int | None:
    """Write a shape the harness created onto an imported slide.

    Its frame is already in EMU because the expander placed it from a named region — the
    tool never saw a coordinate, and neither did the model.
    """
    from pptx.util import Emu

    frame_model = shape_model.frame
    box = slide.shapes.add_textbox(Emu(frame_model.x), Emu(frame_model.y),
                                   Emu(frame_model.cx), Emu(frame_model.cy))
    box.name = f"{model.id}:{shape_model.id}"
    _write_runs(box.text_frame, shape_model.runs, shape_model.text or "")
    if shape_model.type_spec is not None:
        _apply_spec(box.text_frame, shape_model.type_spec)
    _write_paragraph_props(box.text_frame, shape_model)
    fidelity.harden_text_frame(
        box.text_frame,
        line_pt=(shape_model.type_spec.line * 0.75) if shape_model.type_spec else None,
    )
    return int(box.shape_id)


def _apply_spec(frame, spec) -> None:
    """Size and weight from the theme's type scale, converted to points."""
    from pptx.util import Pt

    for para in frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(spec.size * 0.75)
            if spec.weight >= 600:
                run.font.bold = True
            if spec.italic:
                run.font.italic = True


def _write_runs(frame, runs, plain: str) -> None:
    """Write text as OOXML runs, so emphasis survives as emphasis.

    `frame.text = ...` collapses a paragraph to a single run, which is why the harness
    could only ever write unemphasised text. Multi-line content still becomes one paragraph
    per line; runs live inside those.
    """

    if not runs:
        frame.text = plain
        return

    frame.clear()
    paragraph = frame.paragraphs[0]
    for run in runs:
        for index, chunk in enumerate(run.text.split("\n")):
            if index:
                paragraph = frame.add_paragraph()
            if not chunk:
                continue
            written = paragraph.add_run()
            written.text = chunk
            font = written.font
            font.bold = run.bold or None
            font.italic = run.italic or None
            font.underline = True if run.underline else None
            if run.strike or run.script:
                _raw_run_properties(written, run)
            _write_link(written, run.link)


def _raw_run_properties(written, run) -> None:
    """Strikethrough and super/subscript, which python-pptx does not expose.

    Both are plain attributes on `<a:rPr>`; `baseline` is thousandths of the font size and
    these are the values PowerPoint itself writes.
    """
    from ..state.richtext import SUBSCRIPT, SUPERSCRIPT

    properties = written._r.get_or_add_rPr()
    if run.strike:
        properties.set("strike", "sngStrike")
    if run.script == "super":
        properties.set("baseline", str(SUPERSCRIPT))
    elif run.script == "sub":
        properties.set("baseline", str(SUBSCRIPT))


def _is_text_slot(laid_out) -> bool:
    """Only text boxes carry the fidelity contract.

    A table's cells and a chart's plot area are PowerPoint's to lay out — that is the one
    place it is allowed to, because the alternative costs the recipient an editable object.
    Auditing them against a text-box rule would fail every deck that has one.
    """
    return laid_out.shape not in ("tabular", "chart", "media")


def _write_object(slide, laid_out, value, name: str, box, theme,
                  assets: Assets | None = None,
                  size: tuple[int, int] | None = None) -> int | None:
    """Write a non-text slot as a real object: a table, a chart, or a picture."""
    x, y, w, h = box
    frame = Frame(x=x, y=y, cx=w, cy=h)

    if laid_out.shape == "media":
        return _write_media(slide, laid_out, value, name, theme, assets, size)

    if laid_out.shape == "tabular" and isinstance(value, dict):
        spec = TableSpec(headers=[str(c) for c in (value.get("headers") or [])],
                         rows=[[str(c) for c in row] for row in (value.get("rows") or [])])
        if not (spec.headers or spec.rows):
            return None
        paint = decoration_mod.paint_for(theme, laid_out.decoration, laid_out.shape)
        return _add_table(slide, Shape(id=name, ooxml_id=0, type="table",
                                       frame=frame, table=spec),
                          band=paint.fill or None)

    if laid_out.shape == "chart" and isinstance(value, dict):
        # `categories` is what `ChartSpec`, the `add_chart` schema and the model all call
        # them. Reading `labels` here meant a managed chart slot authored exactly as
        # documented produced an empty `ChartSpec`, failed the guard below, and returned
        # None — `add_slide` answered `ok`, the write was measured, and the exported slide
        # had no chart on it. A silent drop, which is the worst kind: `labels` is accepted
        # too, because a deck already written that way should not lose its chart now.
        categories = value.get("categories") or value.get("labels") or []
        spec = ChartSpec(kind=str(value.get("kind") or "bar"),
                         categories=[str(c) for c in categories],
                         series=list(value.get("series") or []))
        if not (spec.categories and spec.series):
            return None
        return _add_chart(slide, Shape(id=name, ooxml_id=0, type="chart",
                                       frame=frame, chart=spec))
    return None


def _write_media(slide, laid_out, value, name: str, theme,
                 assets: Assets | None, size: tuple[int, int] | None) -> int | None:
    """A managed `media` slot as a real picture.

    This is the slot the writer used to skip, and skipping it was not a small gap: it made
    `image_full`'s `full` and `inset` the same empty frame, left `image_split` unable to put
    a picture beside its prose, and so retired both components while the catalog went on
    offering them. Imagery is the largest visual lever the harness has and it was writing
    none of it.

    The rectangle is the expander's, twice over. `media_scale` shrinks the slot box about its
    own centre and `Box.fit` letterboxes the picture inside what is left — see `io/media.py`
    for why proportions are held rather than cropped — so the picture is inside its box by
    construction, and the preview draws the same rule with `object-fit: contain` on the same
    scaled box. Neither of them is allowed its own arithmetic.
    """
    found = media_mod.payload(value)
    if found is None or size is None:
        return None
    asset_id, alt = found
    asset = media_mod.resolve(asset_id, assets)
    if asset is None:
        return None

    placed = (laid_out.box
              .scaled(overrides_mod.media_factor(laid_out.overrides))
              .fit(asset.aspect))
    x, y, w, h = placed.emu(*theme.grid.canvas, *size)
    return _add_picture(slide, Shape(id=name, ooxml_id=0, type="picture",
                                     frame=Frame(x=x, y=y, cx=w, cy=h), alt=alt),
                        image=asset.image())


def _is_stat_slot(value) -> bool:
    """Does this slot hold figure-and-label pairs? True only when they all are.

    All, not any: a slot mixing stats with plain strings would style the second paragraph of
    every cell as a caption, including cells whose second line is an ordinary wrapped
    sentence.
    """
    return (isinstance(value, list) and bool(value)
            and all(slots.is_stat(item) for item in value))


#: The coordinate space an icon's path is written in. `a:path` declares its own `w`/`h` and
#: the shape scales that space into its extent, so this is a *precision*, not a size: 21600
#: units across a 24-unit view box is 900 steps per authored unit, which is finer than any
#: renderer's device pixel at slide scale. It is also the number PowerPoint's own preset
#: geometry uses, so nothing here is exercising an unusual corner of the format.
ICON_PATH_UNITS = 21600


def _custom_geometry(sp_pr, name: str):
    """One icon's path as an `<a:custGeom>` element, ready to replace a preset one.

    This is the whole reason an icon is a *shape* rather than a picture. python-pptx cannot
    read SVG at all — `add_asset` refuses it with `svg_unsupported`, see `io/media.is_svg` —
    so the choice was never "vector or raster", it was "native geometry or a rasteriser in
    the build". Native wins on the thing that matters here: a `custGeom` stroke resolves from
    the theme at write time, so one path serves every palette, where a PNG is one colour
    forever and a PNG per accent is a cache that goes stale the first time a theme changes.

    Only `moveTo`, `lnTo`, `cubicBezTo` and `close` are emitted. The vendored paths are
    normalised to exactly those at vendor time precisely so this function has no arc
    conversion in it — SVG parametrises an elliptic arc by its endpoint and two flags and
    DrawingML by centre and sweep, and doing that conversion here would put trigonometry in
    the exporter that the HTML preview would then have to reproduce and could get wrong.

    `fill="none"` on the path, because these are open strokes. It is also what keeps the
    harness out of the one place DrawingML and SVG are not plainly the same shape: how a
    renderer resolves the winding of overlapping subpaths. An outline icon has no counters,
    so there is nothing to resolve.
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    geometry = sp_pr.makeelement(qn("a:custGeom"), {})
    for tag in ("a:avLst", "a:gdLst", "a:ahLst", "a:cxnLst"):
        etree.SubElement(geometry, qn(tag))
    etree.SubElement(geometry, qn("a:rect"), {"l": "0", "t": "0", "r": "r", "b": "b"})
    path_list = etree.SubElement(geometry, qn("a:pathLst"))
    path = etree.SubElement(path_list, qn("a:path"), {
        "w": str(ICON_PATH_UNITS), "h": str(ICON_PATH_UNITS), "fill": "none", "stroke": "1"})

    # `icon_set.placed`, not a scale factor applied here, because the importer matches an
    # incoming path against those same integers (`icons.identify`). Two roundings written
    # two ways is how a mark the harness drew stops being one it can read back — see the
    # note on `placed` for the half-unit case where the two spellings actually disagree.
    for command, values in icon_set.placed(name, ICON_PATH_UNITS):
        if command == "Z":
            etree.SubElement(path, qn("a:close"))
            continue
        node = etree.SubElement(path, qn({"M": "a:moveTo", "L": "a:lnTo",
                                          "C": "a:cubicBezTo"}[command]))
        for index in range(0, len(values), 2):
            # No Y flip. SVG and DrawingML both put the origin at the top left with y
            # increasing downwards, so a path copied across arrives the right way up — worth
            # stating, because the two disagree often enough elsewhere that it reads like an
            # oversight rather than a fact that was checked.
            etree.SubElement(node, qn("a:pt"), {
                "x": str(values[index]),
                "y": str(values[index + 1]),
            })
    return geometry


def _add_icon(slide, name: str, box: tuple[int, int, int, int], colour: str,
              stroke_emu: int, shape_name: str) -> bool:
    """Draw one icon as a real shape the recipient can select, recolour and resize.

    Built by taking an ordinary autoshape and swapping its `<a:prstGeom>` for the custom one,
    rather than by assembling a `<p:sp>` here: that borrows python-pptx's shape-id allocation
    and its empty text body, and the swap is the only part of this that the library has no
    API for. Same technique as `_gradient_fill` above and for the same reason — author the
    element the schema wants, in the place the schema wants it, and let the library keep
    owning everything around it.

    Every property is stated, exactly as `_paint_panel` states them: an autoshape left alone
    inherits the master's accent fill, an outline and a preset shadow, and an icon that
    arrived as a filled rounded rectangle with a drop shadow is not the mark the theme asked
    for.
    """
    from lxml import etree
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    if not icon_set.path(name):
        return False
    x, y, w, h = box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    shape.name = shape_name
    sp_pr = shape._element.spPr
    preset = sp_pr.find(qn("a:prstGeom"))
    if preset is None:  # pragma: no cover - add_shape always writes one
        return False
    sp_pr.replace(preset, _custom_geometry(sp_pr, name))

    shape.fill.background()
    shape.line.color.rgb = RGBColor.from_string(colour.lstrip("#").upper())
    shape.line.width = Emu(max(1, stroke_emu))
    # Round caps and joins, because that is how the set is drawn: Tabler's outline style sets
    # `stroke-linecap="round"` on every icon, and square ends turn a rounded chevron into
    # something that reads as a different, cruder drawing at small sizes.
    line = sp_pr.find(qn("a:ln"))
    if line is not None:
        line.set("cap", "rnd")
        etree.SubElement(line, qn("a:round"))
    shape.shadow.inherit = False
    return True


#: Preset geometry a decoration layer may ask for. Two entries, and that is the point: these
#: are shapes every renderer draws from the same parametric definition, so what the recipient
#: opens is what the preview showed. A soft edge or a blur would be neither.
LAYER_PRESETS = {"roundRect": "ROUNDED_RECTANGLE", "ellipse": "OVAL"}


def _gradient_fill(shape, paint: decoration_mod.Paint) -> None:
    """Write a decoration layer's ramp as a real `<a:gradFill>`.

    Authored as XML rather than through `shape.fill.gradient()` because that API gives two
    stops and no opacity, and both of those are the technique rather than a detail: three
    stops is what stops a ramp reading as a flat fill, and the alpha ramp is the whole of a
    contact shadow. Every element here is stock DrawingML — `gsLst`, `lin`, `path` — so it
    round-trips through PowerPoint, Keynote and LibreOffice without being rasterised.

    `background()` first, because python-pptx puts `<a:noFill>` at the position in `spPr`
    that the schema requires of a fill; replacing that element in place is what keeps the
    part valid without this function knowing the element order.
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    shape.fill.background()
    sp_pr = shape._element.spPr
    placeholder = sp_pr.find(qn("a:noFill"))
    if placeholder is None:  # pragma: no cover - background() always writes one
        return

    # Positions mean the same thing in all three renderings of this ramp: `pos=0` is the
    # focus of a path gradient, `0%` is the centre of a CSS `radial-gradient`, and `offset=0`
    # is the centre of an SVG `radialGradient`. Verified against a real renderer rather than
    # assumed — the convention is easy to get backwards, and backwards is a shadow that is a
    # dark ring with a clear middle.
    grad = sp_pr.makeelement(qn("a:gradFill"), {"rotWithShape": "0"})
    stop_list = etree.SubElement(grad, qn("a:gsLst"))
    for stop in paint.stops:
        node = etree.SubElement(stop_list, qn("a:gs"))
        node.set("pos", str(max(0, min(100000, round(stop.at * 100000)))))
        colour = etree.SubElement(node, qn("a:srgbClr"))
        colour.set("val", stop.colour.lstrip("#").upper())
        if stop.alpha < 1.0:
            etree.SubElement(colour, qn("a:alpha")).set(
                "val", str(max(0, round(stop.alpha * 100000))))

    if paint.kind == "radial":
        # `path=circle` with the focus rectangle collapsed to the centre, so the ramp runs
        # from the middle of the shape out to its rim — which is what makes the contact
        # shadow dark under the panel and gone by the time it reaches the edge.
        path = etree.SubElement(grad, qn("a:path"))
        path.set("path", "circle")
        focus = etree.SubElement(path, qn("a:fillToRect"))
        for edge in ("l", "t", "r", "b"):
            focus.set(edge, "50000")
    else:
        # 60000ths of a degree, clockwise from due east. `scaled=0` keeps the light coming
        # from one direction rather than one that tilts with each panel's aspect ratio.
        line = etree.SubElement(grad, qn("a:lin"))
        line.set("ang", str(round(paint.angle * 60000) % 21600000))
        line.set("scaled", "0")

    sp_pr.replace(placeholder, grad)


def _paint_panel(slide, box: tuple[int, int, int, int], paint: decoration_mod.Paint,
                 name: str, emu_per_px: float) -> None:
    """One shape a decorated variant draws behind a written box.

    A real autoshape, so the recipient can restyle it, and every property is stated. An
    autoshape left alone inherits the master's accent fill, its outline and a preset
    shadow — a card coloured by whichever template the deck is opened against is not the
    card the theme asked for, and it is exactly the divergence the writer exists to close.

    `box` is EMU and `paint.radius` is canvas px, so the scale is passed rather than the
    original px rectangle: `eject_slide` freezes panels as shapes whose geometry is already
    absolute, and it has no px box left to hand back.
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    x, y, w, h = box
    preset = getattr(MSO_SHAPE, LAYER_PRESETS.get(paint.preset, "ROUNDED_RECTANGLE"))
    shape = slide.shapes.add_shape(preset, Emu(x), Emu(y), Emu(w), Emu(h))
    shape.name = name
    # roundRect's adjustment is a fraction of the shorter side rather than a length, so the
    # theme's radius is divided by the box it is drawn on instead of converted to EMU. An
    # ellipse has no adjustment to set — its outline is already its own corner.
    if shape.adjustments:
        shorter_px = min(w, h) / max(emu_per_px, 1e-9)
        shape.adjustments[0] = min(0.5, paint.radius / max(1.0, shorter_px))
    if paint.stops:
        _gradient_fill(shape, paint)
    elif paint.fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(paint.fill.lstrip("#"))
    else:
        shape.fill.background()
    if paint.line:
        shape.line.color.rgb = RGBColor.from_string(paint.line.lstrip("#"))
        shape.line.width = Pt(decoration_mod.LINE_PX * 0.75)  # canvas px at 96dpi -> pt
    else:
        shape.line.fill.background()
    # The theme's shape language says `shadow: none`; inheriting the master's preset one
    # would be the template deciding what a themed card looks like. The depth a decoration
    # wants is drawn as its own shape — see `decoration.CONTACT` — precisely because a
    # preset shadow is a renderer effect and this is geometry.
    shape.shadow.inherit = False


def _style(frame, laid_out: expand.LaidOutSlot, theme, *, stat: bool = False) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    from ..render import fonts as font_stacks

    spec = laid_out.spec
    # The face the *role* asks for, named on the run. Without it a plain text box inherits
    # PowerPoint's minor font, so every `display` role — deck titles, slide titles, stats —
    # exported in the body face while the preview drew them in the display one. The budget
    # made it worse than cosmetic: it measures advance widths against the display stack, so
    # a serif-titled theme was measured in Georgia and rendered in Inter, and "measure the
    # font that will actually render it" quietly stopped being true.
    stack = font_stacks.parse_stack(theme.type.families.get(spec.family, ""))
    face = stack[0] if stack else None
    # Through the theme, never a literal: `emphasis` names an ink role and the palette
    # answers. An override carrying a colour would be `set_font` wearing another name.
    ink = overrides_mod.ink_for(theme, laid_out.overrides)
    align = PP_ALIGN.CENTER if laid_out.align == "center" else PP_ALIGN.LEFT
    # The figure of a stat carries the block's accent, from the theme's ordered ramp. Same
    # rule as the preview's, because the preview is this file rendered.
    figure_ink = overrides_mod.accent_for(theme, laid_out.overrides) if stat else ink
    for index, para in enumerate(frame.paragraphs):
        para.alignment = align
        # A stat cell is one figure and one label. The label rides the same ratio the
        # preview's CSS uses, from the same constant — set at the figure's size it would
        # export a 44px caption, and the file and the preview would disagree about a slide
        # neither of them was wrong about on its own.
        is_label = bool(stat and index)
        size = spec.size * (slots.STAT_LABEL_EM if is_label else 1.0)
        paint = ink if is_label else figure_ink
        for run in para.runs:
            run.font.size = Pt(size * 0.75)  # canvas px at 96dpi -> pt
            if face:
                run.font.name = face
            # The type scale sets the weight of the slot; a run that asked to be bold keeps
            # it. Assigning the spec's weight flat un-bolded every emphasised word on the
            # way out, which is the styling deciding what the author meant. A stat's label
            # is the exception: it is a caption under a heavy figure, and inheriting the
            # figure's weight is what made the file disagree with the preview about it.
            run.font.bold = (spec.weight >= 600 and not is_label) or run.font.bold is True
            run.font.color.rgb = RGBColor.from_string(paint.lstrip("#"))


def _blank_layout(prs):
    """The emptiest layout available. Managed slides carry their own geometry, so an
    opinionated layout would only contribute placeholders to delete."""
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ExportError("presentation has no slide layouts")
    blank = [lay for lay in layouts if not list(lay.placeholders)]
    if blank:
        return blank[0]
    return min(layouts, key=lambda lay: len(list(lay.placeholders)))


def _clear_empty_placeholders(slide) -> None:
    """Inherited placeholders left empty render as 'Click to add title' prompts in
    PowerPoint even though they look blank in a preview."""
    for shape in list(slide.placeholders):
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            shape._element.getparent().remove(shape._element)


# ------------------------------------------------------------------------ export


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _write_theme_part(prs, theme) -> None:
    """State the theme's palette and faces in `ppt/theme/theme1.xml`.

    Only for decks the harness generates. Every shape already carries resolved literal
    colors, so this changes nothing on screen — what it fixes is the *round trip*. Without
    it, reopening an exported deck re-extracts the theme from python-pptx's default
    template, so `new --from brand.pptx` produced brand-colored slides that then grew
    default-colored ones the moment anybody edited them.

    An imported deck keeps its original theme part untouched: mutating export patches
    shapes, and rewriting a theme somebody else authored is not a patch.

    The map back to OOXML is lossy in one direction only. `dk1`/`lt1`/`accent1..6` are read
    straight back by `extract_theme`; `rule`, `ink_muted`, `positive` and `negative` have no
    slot in `clrScheme` and are re-derived from the ones that do — deterministically, and by
    the same code that derived them the first time.
    """
    from ..render import fonts as font_stacks

    part = next((p for p in prs.part.package.iter_parts()
                 if str(p.partname).startswith("/ppt/theme/theme")), None)
    if part is None:  # pragma: no cover - python-pptx's template always has one
        return
    root = part._element if hasattr(part, "_element") else None
    if root is None:
        from lxml import etree
        root = etree.fromstring(part.blob)
        part._blob = None

    scheme = root.find(f".//{{{A_NS}}}clrScheme")
    if scheme is not None:
        accents = [theme.accent(i) for i in range(6)]
        for slot, value in (("dk1", theme.palette.get("ink")),
                            ("lt1", theme.palette.get("bg")),
                            ("dk2", theme.palette.get("ink")),
                            ("lt2", theme.palette.get("surface")),
                            *[(f"accent{i + 1}", accents[i]) for i in range(6)]):
            _set_scheme_color(scheme, slot, value)

    scheme = root.find(f".//{{{A_NS}}}fontScheme")
    if scheme is not None:
        for which, role in (("majorFont", "display"), ("minorFont", "body")):
            stack = font_stacks.parse_stack(theme.type.families.get(role, ""))
            if stack:
                node = scheme.find(f"{{{A_NS}}}{which}/{{{A_NS}}}latin")
                if node is not None:
                    node.set("typeface", stack[0])

    part._blob = None
    if not hasattr(part, "_element"):  # pragma: no cover - defensive
        from lxml import etree
        part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8",
                                    standalone=True)


def _set_scheme_color(scheme, slot: str, value) -> None:
    """Write one `clrScheme` slot as an `srgbClr`.

    `dk1`/`lt1` are usually `sysClr` (windowText/window) in a stock template, which carries
    a `lastClr` hint PowerPoint may prefer. Replacing the child outright rather than setting
    an attribute is what makes the value stick.
    """
    from lxml import etree

    if not isinstance(value, str) or not value.startswith("#"):
        return
    node = scheme.find(f"{{{A_NS}}}{slot}")
    if node is None:
        return
    for child in list(node):
        node.remove(child)
    etree.SubElement(node, f"{{{A_NS}}}srgbClr").set("val", value.lstrip("#").upper())


def export(
    deck: Deck,
    path: Path | str,
    mode: ExportMode = "editable",
    strict: bool = True,
    embed_fonts_: bool = True,
    assets: Assets | None = None,
) -> ExportReport:
    """Write `deck` to `path`.

    With a `source_path`, the original package is copied and patched. Without one, a new
    presentation is built at the theme's canvas size.

    `assets` is the deck's pictures, `{asset_id: (content_type, bytes)}`, and it is a
    parameter rather than something read off the deck because assets are deliberately not
    part of the document model — they live on the store (`DeckStore.assets`) so that a
    174 MB deck's images are not carried through every snapshot. Omitting them is legal and
    costs nothing to a deck with no `media` slot; a deck that has one gets a
    `slot_not_written` violation naming the asset, which is the loud version of the silence
    this exists to prevent.
    """
    path = Path(path)
    report = ExportReport(path=path, mode=mode)

    if deck.source_path:
        source = Path(deck.source_path)
        if not source.exists():
            raise ExportError(f"original package is gone: {source}")
        if source.resolve() == path.resolve():
            raise ExportError("refusing to export over the original package; choose another path")
        shutil.copy2(source, path)
        prs = Presentation(str(path))
    else:
        prs = Presentation()
        cx, cy = expand.slide_size_emu(deck.theme)
        prs.slide_width, prs.slide_height = Emu(cx), Emu(cy)

        _write_theme_part(prs, deck.theme)

    written: set[tuple[int, int]] = set()
    existing = list(prs.slides)
    #: The package slide backing each model slide, in model order.
    placed: list[Any] = []
    #: Slots that were filled and produced nothing. Reported as violations rather than
    #: swallowed — see `_write_managed`.
    dropped: list[fidelity.Violation] = []

    by_id = {s.id: s for s in deck.slides}

    for model in deck.slides:
        if model.mode is Mode.FREEFORM:
            if model.duplicate_of:
                source = by_id.get(model.duplicate_of)
                origin = (source.origin or {}).get("index") if source else None
                if (origin is None or origin >= len(existing)) and not deck.source_path:
                    # The copy of an ejected slide, on a deck that never had a package.
                    # Cloning is for *imported* slides, where the OOXML holds things the
                    # harness cannot rebuild; here the shapes are all there is, so write
                    # them. Same hole as the ejected-slide branch below, one level in:
                    # `duplicate_slide` and `eject_slide` are both shipped, and using them
                    # together on a generated deck made it unexportable.
                    before = len(prs.slides._sldIdLst)
                    report.shapes_added += _write_ejected(prs, model, deck, written, assets)
                    if len(prs.slides._sldIdLst) > before:
                        placed.append(list(prs.slides)[-1])
                    report.slides_written += 1
                    continue
                if origin is None or origin >= len(existing):
                    raise ExportError(
                        f"{model.id} is a copy of {model.duplicate_of}, which is not in "
                        "the original package"
                    )
                clone = _clone_slide(prs, existing[origin])
                index = list(prs.slides).index(clone)
                report.shapes_patched += _patch_freeform(clone, index, model, written, assets)
                report.slides_written += 1
                placed.append(clone)
                continue

            index = (model.origin or {}).get("index")
            if (index is None or index >= len(existing)) and deck.source_path:
                # The deck came from a file, so this slide *should* have had a part in it and
                # does not. That is corruption, and rebuilding it from the model would
                # silently drop whatever the harness never understood — the opposite of what
                # patching the original package is for. Fail loudly.
                raise ExportError(
                    f"freeform slide {model.id} has no original to patch, but this deck came "
                    f"from {Path(deck.source_path).name}; its slide part is missing"
                )
            if index is None or index >= len(existing):
                # A freeform slide with no original is an *ejected* one: the harness
                # generated it, `eject_slide` froze the expander's boxes into real shapes,
                # and there is no package part to patch because there never was one. It has
                # everything needed to be written from scratch — frames in EMU, roles, type
                # specs — so write it.
                #
                # This used to raise, which meant `eject_slide` — shipped, documented, and
                # advertised as the way back from managed — put a generated deck into a state
                # where it could never be exported again. The user's work was unrecoverable
                # through the one door that leads out of the harness.
                before = len(prs.slides._sldIdLst)
                report.shapes_added += _write_ejected(prs, model, deck, written, assets)
                if len(prs.slides._sldIdLst) > before:
                    placed.append(list(prs.slides)[-1])
                report.slides_written += 1
                continue
            report.shapes_patched += _patch_freeform(existing[index], index, model, written, assets)
            placed.append(existing[index])
        else:
            before = len(prs.slides._sldIdLst)
            report.shapes_added += _write_managed(prs, model, deck, written, dropped, assets)
            if len(prs.slides._sldIdLst) > before:
                placed.append(list(prs.slides)[-1])
        report.slides_written += 1

    _reorder_to_match(prs, placed)
    prs.save(str(path))

    # `written`, never `written or None`. An empty set means "we wrote nothing, so check
    # nothing"; collapsing it to None would audit the original author's untouched shapes
    # against a contract that was never theirs and fail every clean round trip.
    if embed_fonts_ and mode != "pixel_locked":
        # After the package is finished, like the audit — nothing upstream needs to know.
        report.fonts = embed_fonts.embed(path, deck.theme, deck).as_dict()

    report.violations = dropped + fidelity.assert_fidelity(path, only_shapes=written)
    if strict and report.violations:
        fidelity.raise_for_violations(report.violations)
    return report


def _clone_slide(prs, source):
    """Copy a slide part, relationships and all.

    The shapes are copied as XML rather than rebuilt from the model, because the harness
    does not understand every shape it can duplicate — a copy that quietly dropped the
    SmartArt would not be a copy. Relationships are re-pointed at the same targets, so the
    image parts and embedded objects are shared rather than duplicated.
    """
    import copy as copymod

    layout = source.slide_layout
    clone = prs.slides.add_slide(layout)

    # Replace the new slide's shape tree wholesale with the source's.
    tree = clone.shapes._spTree
    for child in list(tree):
        tag = child.tag.rsplit("}", 1)[-1]
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            tree.remove(child)
    for child in source.shapes._spTree:
        tag = child.tag.rsplit("}", 1)[-1]
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue
        tree.append(copymod.deepcopy(child))

    # Relationships are re-created and the cloned XML is re-pointed at the new ids. The
    # ids cannot be preserved — python-pptx allocates them — so remapping is the only way
    # `r:embed` still resolves. Targets are shared rather than copied, which is what
    # PowerPoint itself does when it duplicates a slide.
    remap: dict[str, str] = {}
    for rel in source.part.rels.values():
        if rel.reltype.endswith("/slideLayout"):
            continue  # the clone already has one, from add_slide
        target = rel.target_ref if rel.is_external else rel.target_part
        remap[rel.rId] = clone.part.relate_to(target, rel.reltype,
                                              is_external=rel.is_external)

    if remap:
        _remap_rels(tree, remap)
    return clone


def _remap_rels(element, remap: dict[str, str]) -> None:
    """Rewrite every `r:` attribute that names a relationship id."""
    prefix = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for node in element.iter():
        for name, value in list(node.attrib.items()):
            if name.startswith(prefix) and value in remap:
                node.set(name, remap[value])


def _reorder_to_match(prs, placed: list[Any]) -> None:
    """Put the package's slide order back in sync with the model.

    Driven by identity, not by position. The package holds imported slides in their original
    order and then everything appended — managed slides *and* duplicates — in creation
    order, so no rule about which class comes first survives contact with a copied slide.
    `placed` is the actual package slide backing each model slide, recorded as export walked
    them, which is the one thing that cannot drift.

    python-pptx has no reorder API; `sldIdLst` is the only place order lives.
    """
    id_list = prs.slides._sldIdLst
    entries = {int(e.get("id")): e for e in id_list}

    ordered = []
    for slide in placed:
        entry = entries.get(int(slide.slide_id))
        if entry is None:
            return  # leaving the order alone beats guessing at it
        ordered.append(entry)
    if len(ordered) != len(entries):
        return

    for entry in list(id_list):
        id_list.remove(entry)
    for entry in ordered:
        id_list.append(entry)
